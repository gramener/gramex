# -*- coding: utf-8 -*-
from __future__ import unicode_literals

import os
import tempfile
from six.moves.urllib_parse import urlparse
import requests
import numpy as np
import pandas as pd
import matplotlib.cm
import matplotlib.colors
from orderedattrdict import AttrDict
from tornado.template import Template
from tornado.escape import to_unicode
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.shapetree import SlideShapes
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from unittest import TestCase
from nose.tools import eq_, ok_, assert_raises
from pandas.util.testing import assert_frame_equal
from gramex import pptgen
from . import folder, sales_file


class TestPPTGen(TestCase):
    # Test Case module for pptgen

    @classmethod
    def setUp(cls):
        # Setup class method to initialize common variables.
        np.random.seed(0)
        cls.input = os.path.join(folder, 'input.pptx')
        cls.output = os.path.join(folder, 'output.pptx')
        cls.image = os.path.join(folder, 'small-image.jpg')
        cls.data = pd.read_excel(sales_file, encoding='utf-8')
        if os.path.exists(cls.output):
            os.remove(cls.output)

    def template(self, tmpl, data):
        # Function to generate tornado template.
        return to_unicode(Template(tmpl, autoescape=None).generate(**data))

    def get_shape(self, target, name):
        # Function to return required shape.
        return [shape for shape in target.slides[0].shapes if shape.name == name]

    def draw_chart(self, slidenumber, data, rule):
        # Function to change data in pptx native charts.
        return pptgen.pptgen(
            source=self.input,
            only=slidenumber,
            data={'data': data},
            edit_chart=rule)

    def check_if_number(self, data):
        # Function to check if data is a number
        decimal_points = 6
        try:
            return round(float(data), decimal_points)
        except TypeError:
            return data

    def check_opacity(self, opacity, prop):
        # Check opacity.
        opacity_constant = 100000
        if opacity:
            opacity_val = prop._xFill.srgbClr.xpath('./a:alpha')[0].values()[0]
            eq_(opacity, float(opacity_val) / opacity_constant)

    def check_chart(self, shape, data, series, xaxis, chart_name, chart_colors, opacity):
        # Function to validata pptx-native charts updated data.
        # Chart shape must have `chart_type` attribute
        ok_(hasattr(shape.chart, 'chart_type'))
        eq_(sorted([_series.name for _series in shape.chart.series]), series)
        for _series in shape.chart.series:
            datapints = [self.check_if_number(point) for point in data[_series.name].tolist()]
            seriesvals = [self.check_if_number(point) for point in list(_series.values)]
            eq_(datapints, seriesvals)
            # Checking color of chart
            for index, point in enumerate(_series.points):
                if chart_name == 'Scatter Chart':
                    marker, keyname = point.marker, _series.name
                elif chart_name == 'Area Chart':
                    marker, keyname = _series, _series.name
                elif chart_name in ['Pie Chart', 'Donut Chart']:
                    marker, keyname = point, data.loc[index][xaxis]
                else:
                    marker, keyname = point, _series.name
                color = marker.format.fill.fore_color
                stroke = marker.format.line.fill.fore_color
                self.check_opacity(opacity, color)
                eq_('#{}'.format(color.rgb), chart_colors[keyname])
                eq_('#{}'.format(stroke.rgb), chart_colors[keyname])

    def test_register(self):
        # Test case for register in pptgen
        with assert_raises(ValueError):
            pptgen.pptgen(source=self.input, only=1, register='dummy_command_register')
        pptgen.pptgen(
            source=self.input,
            only=1,
            register={
                'custom_function1': {
                    'function': 'lambda shape, spec, data: (shape, spec, data)'
                }
            })
        ok_('custom_function1' in pptgen.COMMANDS_LIST)
        pptgen.COMMANDS_LIST['custom_function1'](**{'shape': 1, 'spec': {}, 'data': {}})
        with assert_raises(TypeError):
            pptgen.COMMANDS_LIST['custom_function1'](**{'shape': 1, 'spec': {}, 'dummy': True})
        # Registering second command
        pptgen.pptgen(
            source=self.input,
            only=1,
            register={
                'custom_function2': 'lambda x: x(shape, spec, data)'
            })
        ok_('custom_function2' in pptgen.COMMANDS_LIST)
        pptgen.COMMANDS_LIST['custom_function2'](**{'shape': 1, 'spec': {}, 'data': {}})
        with assert_raises(TypeError):
            pptgen.COMMANDS_LIST['custom_function2'](**{'shape': 1, 'spec': {}, 'dummy': True})

    def test_data_format(self):
        # Testing data section. Data argument must be a `dict` like object
        with assert_raises(ValueError):
            pptgen.pptgen(source=self.input, only=1, data=[1, 2])
        for case in [{}, {'function': '{}'}, {'data': [1, 2, 3]}]:
            pptgen.pptgen(source=self.input, only=1, data=case)

    def test_source_without_target(self):
        # Test case to compare no change.
        target = pptgen.pptgen(source=self.input, only=1)
        eq_(len(target.slides), 1)
        eq_(target.slides[0].shapes.title.text, 'Input.pptx')

    def test_source_with_target(self):
        # Test case to compare target output title.
        pptgen.pptgen(source=self.input, target=self.output, only=1)
        target = Presentation(self.output)
        eq_(len(target.slides), 1)
        eq_(target.slides[0].shapes.title.text, 'Input.pptx')

    def test_change_title(self):
        # Title change test case with unicode value.
        text = '高σ高λس►'
        target = pptgen.pptgen(
            source=self.input, only=1,
            change={                  # Configurations are same as when loading from the YAML file
                'Title 1': {            # Take the shape named 'Title 1'
                    'text': text          # Replace its text with new text
                }
            })
        eq_(target.slides[0].shapes.title.text, text)

    def test_text_xml(self):
        # Test case for text xml object
        text = 'New Title<text color="#00ff00" bold="True" font-size="14"> Green Bold Text</text>'
        font_size, text_color, pix_to_inch = 12, '#ff0000', 10000
        target = pptgen.pptgen(
            source=self.input, only=1,
            change={
                'Title 1': {
                    'style': {
                        'color': text_color,
                        'font-size': font_size
                    },
                    'text': text
                }
            })
        eq_(target.slides[0].shapes.title.text, 'New Title Green Bold Text')
        text_shape = self.get_shape(target, 'Title 1')[0]
        total_runs = 2
        # Maximum one paragraph should be there, because entire text area is being changed
        eq_(len(text_shape.text_frame.paragraphs), 1)
        # Maximum two runs should be there, because we are adding two seperate style for texts
        eq_(len(text_shape.text_frame.paragraphs[0].runs), total_runs)

        first_run = text_shape.text_frame.paragraphs[0].runs[0]
        second_run = text_shape.text_frame.paragraphs[0].runs[1]

        eq_('{}'.format(first_run.font.color.rgb), 'FF0000')
        eq_('{}'.format(first_run.text), 'New Title')
        eq_(int(int('{}'.format(first_run.font.size)) / pix_to_inch), font_size - 1)

        second_run_font = 14
        eq_('{}'.format(second_run.font.color.rgb), '00FF00')
        eq_('{}'.format(second_run.text), ' Green Bold Text')
        eq_(int(int('{}'.format(second_run.font.size)) / pix_to_inch), second_run_font - 1)

    def test_replicate_slide(self):
        # Test slide replication.
        data = self.data.groupby('city')
        tmpl = 'Region: {{ city }} has Sales: {{ sales }} with Growth: {{ growth }}'
        target = pptgen.pptgen(
            source=self.input,
            only=3,
            data={'data': data},
            replicate_slide={
                'replicate': True,
                'data': 'data["data"]',
                'sales-text': {
                    'data': 'data[0]',
                    'text': tmpl
                }
            })
        eq_(data.ngroups, len(target.slides))
        contents_ex = [self.template(tmpl, grp.iloc[0]) for _, grp in data]
        contents_ppt = [shape.text for slide in target.slides
                        for shape in slide.shapes
                        if shape.name == 'sales-text']
        eq_(contents_ex, contents_ppt)

    def test_text_style(self):
        # Test case for testing text styles.
        target = pptgen.pptgen(
            source=self.input,
            only=4,
            replicate_slide={
                'Title 1': {
                    'text': 'New title',
                    'style': {
                        'color': '#ff0000'
                    }
                },
                'Text 1': {
                    'text': 'New text',
                    'style': {
                        'color': '#00ff00',
                        'bold': True
                    }
                }
            })
        slides = target.slides
        eq_(len(slides), 1)
        name_map = {'Title 1': {'color': 'FF0000', 'text': 'Title 1'},
                    'Text 1': {'color': '00FF00', 'text': 'Text 1'}}
        for shape in slides[0].shapes:
            if shape.name in ['Title 1', 'Text 1']:
                continue
            eq_(shape.text, name_map[shape.name]['text'])
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    eq_('{}'.format(run.font.color.rgb), name_map[shape.name]['color'])

    def test_group_and_image(self):
        # Test case for group objects.
        for img in [self.image, 'https://gramener.com/uistatic/img/store-supply-chain.png']:
            target = pptgen.pptgen(
                source=self.input,
                only=5,
                group_test={
                    'slide-title': 'Group Test',
                    'Group 1': {
                        'Caption': {
                            'text': 'New caption'
                        },
                        'Picture': {
                            'image': img
                        }
                    }
                })
            eq_(len(target.slides), 1)
            grp_shape = self.get_shape(target, 'Group 1')[0]
            for shape in SlideShapes(grp_shape.element, grp_shape):
                if shape.name == 'Caption':
                    eq_(shape.text, 'New caption')
                if shape.name == 'Picture':
                    if urlparse(img).netloc:
                        r = requests.get(img)
                        with tempfile.NamedTemporaryFile(delete=False) as handle:
                            handle.write(r.content)
                        with open(handle.name, 'rb') as f:
                            blob = f.read()
                        os.unlink(handle.name)
                    else:
                        with open(img, 'rb') as f:
                            blob = f.read()
                    eq_(shape.image.blob, blob)

    def test_stack(self):
        # Test case for stack elements.
        data = self.data.groupby('city', as_index=False)['sales', 'growth'].sum()
        data = data.to_dict(orient='records')
        tmpl = 'Region: {{ city }} has Sales: {{ sales }} with Growth: {{ growth }}'
        target = pptgen.pptgen(
            source=self.input,
            only=6,
            data={'data': data},
            stack_shapes={
                'slide-number': 1,
                'TextBox 1': {
                    'data': 'data["data"]',
                    'stack': 'vertical',
                    'margin': 0.10,
                    'text': tmpl
                }
            })
        eq_(len(target.slides), 1)
        contents_ex = [self.template(tmpl, item) for item in data]
        contents_ppt = [shape.text for shape in self.get_shape(target, 'TextBox 1')]
        eq_(contents_ex, contents_ppt)

    def test_replace(self):
        # Test case for replace command.
        target = pptgen.pptgen(
            source=self.input, only=7,
            change={
                'slide-number': [1],
                'TextBox 1': {
                    'replace': {
                        'Old': 'New',
                        'Title': 'Heading'
                    }
                }
            })
        eq_(len(target.slides), 1)
        txt_box = self.get_shape(target, 'TextBox 1')[0]
        ok_(txt_box.has_text_frame)
        eq_(txt_box.text, 'New Heading')

    def test_css(self):
        # Test case for `css` command.
        pix_to_inch = 10000
        spec = {'width': 200, 'height': 200, 'top': 100, 'left': 50}
        opacity_val, opacity_constant = 0.1, 100000
        target = pptgen.pptgen(
            source=self.input, only=8,
            change={
                'Rectangle 1': {
                    'css': {
                        'style': {
                            'opacity': opacity_val,
                            'color': '#ff0000',
                            'fill': '#ffff00',
                            'stroke': '#00ff00',
                            'width': spec['width'],
                            'height': spec['height'],
                            'left': spec['left'],
                            'top': spec['top']
                        }
                    }
                }
            })
        eq_(len(target.slides), 1)
        shape = self.get_shape(target, 'Rectangle 1')[0]
        opacity = shape.fill.fore_color._xFill.srgbClr.xpath('./a:alpha')[0].values()[0]
        for k in spec:
            eq_(getattr(shape, k), spec[k] * pix_to_inch)
        eq_(opacity_val, float(opacity) / opacity_constant)
        eq_('{}'.format(shape.fill.fore_color.rgb), 'FFFF00')
        eq_('{}'.format(shape.line.fill.fore_color.rgb), '00FF00')

    def test_table(self):
        # Function to test `table` command.
        with assert_raises(AttributeError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                change={
                    'Invalid Input': {
                        'table': {
                            'data': 'data'
                        }
                    }
                })

        target = pptgen.pptgen(
            source=self.input, only=9,
            data={'data': self.data},
            change={
                'Table 1': {
                    'table': {
                        'data': 'data["data"]',
                        'style': {
                            'fill': '#cccccc'
                        }
                    }
                }
            })
        eq_(len(target.slides), 1)
        shape = self.get_shape(target, 'Table 1')[0]
        eq_(len(shape.table.rows), len(self.data.index) + 1)
        eq_(len(shape.table.columns), len(self.data.columns))
        table_data, columns = {}, []
        for rno, row in enumerate(shape.table.rows):
            if rno == 0:
                columns = [c.text_frame.text for c in row.cells]
                continue
            for cno, cell in enumerate(row.cells):
                txt = cell.text_frame.text
                if columns[cno] not in table_data:
                    table_data[columns[cno]] = []
                table_data[columns[cno]].append(txt)
        table_data = pd.DataFrame(table_data)[columns].replace('nan', np.nan)
        for c in table_data:
            table_data[c] = table_data[c].astype(self.data[c].dtype)
        # Comparinig `dataframe` from table with original `dataframe`
        assert_frame_equal(table_data, self.data, check_names=True)

    def test_change_chart(self):
        # Test case for all native charts charts.
        with assert_raises(AttributeError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                change={
                    'Invalid Input': {
                        'chart': {
                            'data': 'data'
                        }
                    }
                })

        xaxis, opacity = 'city', 0.5
        slidenumbers = AttrDict(
            Bar_Chart=2, Column_Chart=10, Line_Chart=11, Area_Chart=12,
            Scatter_Chart=13, Bubble_Chart=14, Bubble_Chart_3D=15,
            Radar_Chart=16, Donut_Chart=17, Pie_Chart=18)
        for chart_name, slidenumber in slidenumbers.items():
            # Replacing `_` with a white space. Because chart names in input slides contains
            # spaces not `_`.
            chart_name = chart_name.replace('_', ' ')
            if chart_name in ['Pie Chart', 'Donut Chart']:
                series = ['sales']
                chart_colors = {
                    'Singapore': '#D34817',
                    'Hyderabad': '#9B2D1F',
                    'Bangalore': '#A28E6A',
                    'Coimbatore': '#956251',
                    'Newport Beach': '#918485',
                    'South Plainfield': '#855D5D'
                }
            else:
                series = ['growth', 'sales']
                chart_colors = {
                    'sales': '#D34817',
                    'growth': '#9B2D1F',
                }
            data = self.data.groupby(xaxis, as_index=False)[series].sum()
            rule = {
                chart_name: {
                    'chart': {
                        'data': 'data["data"]',
                        'x': xaxis,
                        'color': chart_colors,
                        'opacity': opacity
                    }
                }
            }
            target = self.draw_chart(slidenumber, data, rule)
            shape = self.get_shape(target, chart_name)[0]
            self.check_chart(shape, data, series, xaxis, chart_name, chart_colors, opacity)

    def test_treemap(self):
        # Test case for treemap.
        # If shape is not a rectangle
        with assert_raises(NotImplementedError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                drawtreemap={'Invalid Input': {'treemap': {'data': 'data'}}})

        target = pptgen.pptgen(
            source=self.input, only=19,
            data={'data': self.data},
            drawtreemap={
                'Treemap Rectangle': {
                    'treemap': {
                        'data': 'data["data"]',
                        'keys': ['city'],
                        'values': '{"sales": "sum", "growth": "sum"}',
                        'size': {'function': 'lambda v: v["sales"]'},
                        'sort': {
                            'function': 'lambda v: v.sort_values(by=["sales"], ascending=False)'
                        },
                        'color': {
                            'function': 'lambda v: _color.gradient(v["growth"]/100, _color.RdYlGn)'
                        },
                        'text': {
                            'function': 'lambda v: "{}".format(v["city"])'
                        }
                    }
                }
            })
        data = self.data.groupby('city', as_index=False).agg({'sales': 'sum', 'growth': 'sum'})
        # Sorting as per treemap logic in chart and taking a subset to compare with
        # chart data
        data = data.sort_values(by=['sales'], ascending=False)[['city']]
        data = data.reset_index(drop=True)
        # Only one slide should be there
        eq_(len(target.slides), 1)
        text, width = [], []
        for shape in target.slides[0].shapes:
            if shape.shape_type == MSO_SHAPE.RECTANGLE:
                # Gettiing text as per treemap draw logic
                text.append(shape.text)
                # Getting rectangles with
                width.append(shape.width)
        # Creating a new dataframe from treemap chart
        treemapdata = pd.DataFrame({'city': text, 'sales': width})
        # Treemap data should to exactly in same order as per data
        assert_frame_equal(data, treemapdata[['city']], check_names=True)

    def test_horizontal_vertical_bullet_chart(self):
        # Test case for horizontal and vertical bullet chart
        # If shape is not a rectangle
        with assert_raises(NotImplementedError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                drawbullet={'Invalid Input': {'bullet': {'data': 'data'}}})

        default_size = 10000
        change_data = [
            {
                'text': {'function': 'lambda v: "%.1f" % v'},
                'poor': 20,
                'slidenumber': 20,
                'shapename': 'Bullet Rectangle Horizontal',
                'orient': 'horizontal'
            },
            {
                'text': False,
                'poor': 0,
                'shapename': 'Bullet Rectangle Vertical',
                'orient': 'vertical',
                'slidenumber': 21
            },
            {
                'text': False,
                'poor': 0,
                'hi': 10,
                'lo': 10,
                'shapename': 'Bullet Rectangle Vertical',
                'orient': 'vertical',
                'slidenumber': 21
            }
        ]

        input_rect = Presentation(self.input)
        for update_data in change_data:
            # Getting required param from config to compare with output
            orient = update_data['orient']
            shpname = update_data['shapename']
            slidenumber = update_data['slidenumber']
            # Getting shape name
            shapes = input_rect.slides[slidenumber - 1].shapes
            _shp = [shape for shape in shapes if shape.name == shpname][0]
            # width, height = _shp.width, _shp.height
            height = _shp.height if orient == 'horizontal' else _shp.width
            width = _shp.width if orient == 'horizontal' else _shp.height
            lo = update_data.get('lo', 0)
            hi = update_data.get('hi', self.data['sales'].max())
            good = 100
            target = pptgen.pptgen(
                source=self.input, only=slidenumber,
                data={'data': self.data},
                draw_bullet={
                    shpname: {
                        'bullet': {
                            'data': 'data["data"]["sales"].loc[0]',
                            'max-width': 1,
                            'poor': update_data['poor'],
                            'good': good,
                            'lo': lo,
                            'hi': hi,
                            'target': 'data["data"]["sales"].max()',
                            'average': 'data["data"]["sales"].mean()',
                            'orient': orient,
                            'gradient': 'Oranges',
                            'text': update_data['text']
                        }
                    }
                })
            eq_(len(target.slides), 1)
            textboxes, rectangles = 0, 0
            rects_width, rects_height = [], []
            _average = self.data['sales'].mean()
            _data = self.data['sales'].loc[0]
            for rectdata in [good, _average, update_data['poor'], _data]:
                if not rectdata:
                    continue
                shp_width = ((rectdata - lo) / ((hi - lo) or np.nan)) * width
                rects_width.append(shp_width if orient == 'horizontal' else height)
                rects_height.append(height if orient == 'horizontal' else shp_width)

            # Removing and updatiing actual data point's height
            if orient == 'horizontal':
                rects_height.pop(-1)
                rects_height.append(int(height / 2.0))
            if orient == 'vertical':
                rects_width.pop(-1)
                rects_width.append(int(height / 2.0))
            # Adding target's width and height
            rects_height.append(default_size if orient == 'vertical' else height)
            rects_width.append(default_size if orient == 'horizontal' else height)

            rects_width_from_output, rects_height_from_output = [], []
            for shape in target.slides[0].shapes:
                if shape.name == 'Title 1':
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    textboxes += 1
                elif shape.shape_type == MSO_SHAPE.RECTANGLE:
                    rectangles += 1
                    rects_width_from_output.append(shape.width)
                    rects_height_from_output.append(shape.height)
            # Comapring number of text boxes to be shown
            total_text_boxes = 2 if update_data['text'] else 0
            eq_(textboxes, total_text_boxes)
            # Comapring number of rectangle shapes
            total_rects = 5 if update_data['poor'] else 4
            if hi == lo:
                total_rects, rects_width, rects_height = 0, [], []
            eq_(rectangles, total_rects)
            # Comparing rectangle's height
            eq_(rects_height_from_output, list(map(int, rects_height)))
            # Comparing rectangle's width
            eq_(rects_width_from_output, list(map(int, rects_width)))

    def test_heatgrid(self):
        # Test case for heatgrid
        # If shape is not a rectangle
        with assert_raises(NotImplementedError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                drawheatgrid={'Invalid Input': {'heatgrid': {'data': 'data'}}})

        pix_to_inch = 10000
        white = '#cccccc'
        default_margin, cell_width, cell_height, leftmargin, font = 5, 60, 50, 0.20, 14
        row, column, value, color = 'देश', 'city', 'sales', 'RdYlGn'
        data = self.data.groupby([row, column], as_index=False).agg({value: 'sum'})
        target = pptgen.pptgen(
            source=self.input, only=22,
            data={'data': data},
            drawheatgrid={
                'Heatgrid Rectangle': {
                    'heatgrid': {
                        'data': 'data["data"]',
                        'row': row,
                        'column': column,
                        'value': value,
                        'text': {
                            'function': '"{}".format(data["' + column + '"])'
                        },
                        'left-margin': leftmargin,
                        'cell-width': cell_width,
                        'cell-height': cell_height,
                        'na-text': 'NA',
                        'na-color': white,
                        'style': {
                            'gradient': {
                                'function': 'lambda data, handler: "{}"'.format(color)
                            },
                            'font-size': font,
                            'text-align': 'center',
                            'stroke': '#cccccc',
                            'padding': {
                                'left': 5
                            }
                        }
                    }
                }
            })
        eq_(len(target.slides), 1)
        gradient = matplotlib.cm.get_cmap(color)
        cross = pptgen.utils.scale(data.pivot(row, column, value))
        colors_ex = [white if np.isnan(x) else matplotlib.colors.to_hex(gradient(x))
                     for x in cross.flatten()]
        label_rows = np.sort(data[row].unique())
        label_cols = np.sort(data[column].unique())
        labels = label_cols
        for r in label_rows:
            labels = np.hstack((labels, label_cols, r))
        labels_ex = labels.tolist()
        rect_shape_width = (cell_width - 2 * default_margin) * pix_to_inch
        rect_shape_height = (cell_height - default_margin) * pix_to_inch
        labels_ppt, colors_ppt = [], []
        for shape in target.slides[0].shapes:
            if shape.name == 'Title 1':
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                ok_(shape.has_text_frame)
                p = shape.text_frame.paragraphs[0]
                eq_(int(p.runs[0].font.size / pix_to_inch), font - 1)
                ok_(p.alignment == PP_ALIGN.CENTER)
                labels_ppt.append(p.runs[0].text)
            elif shape.shape_type == MSO_SHAPE.RECTANGLE:
                eq_(shape.width, rect_shape_width)
                eq_(shape.height, rect_shape_height)
                colors_ppt.append('#{}'.format(shape.fill.fore_color.rgb).lower())
        eq_(labels_ex, labels_ppt)
        eq_(colors_ex, colors_ppt)

    def test_sankey(self):
        # Test case for sankey
        # If shape is not a rectangle
        with assert_raises(NotImplementedError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                drawsankey={'Invalid Input': {'sankey': {'data': 'data'}}})

        shpname = 'Sankey Rectangle'
        slidenumber = 23
        input_shapes = Presentation(self.input).slides[slidenumber - 1].shapes
        width = [shape for shape in input_shapes if shape.name == shpname][0].width
        groups = ['देश', 'city', 'product']
        order = {'function': 'lambda g: g["sales"].sum()'}
        text = {'function': 'lambda g: g.apply(lambda x: x.name)'}
        color = {'function': 'lambda g: _color.gradient(g["growth"].sum(), _color.RdYlGn)'}
        data = self.data.fillna(0)
        target = pptgen.pptgen(
            source=self.input, only=slidenumber,
            data={'data': data},
            drawsankey={
                shpname: {
                    'sankey': {
                        'data': 'data["data"]',
                        'sort': True,
                        'groups': groups,
                        'color': color,
                        'text': text,
                        'order': order
                    }
                }
            })
        eq_(len(target.slides), 1)
        total_cust_shapes, get_rect_order, get_rect_width = 0, [], []
        grp_order = 'lambda g: g["sales"].sum()'

        for index, grp in enumerate(groups):
            grpobj = data.groupby(grp)
            frame = pd.DataFrame({'size': grpobj[grp].count(), 'seq': eval(grp_order)(grpobj)})
            frame['width'] = frame['size'] / float(frame['size'].sum()) * width
            frame = frame.sort_values(by=['seq'])
            get_rect_width.extend([int(i) for i in frame['width'].tolist()])
            get_rect_order.extend(list(frame.index))
            if index < len(groups) - 1:
                grpby = [groups[index], groups[index + 1]]
                total_cust_shapes += len(data.groupby(grpby, as_index=False)['sales'].sum())

        cust_shape_count, rect_order, rect_width = 0, [], []
        for shape in target.slides[0].shapes:
            if shape.name == 'Title 1':
                continue
            # Custom shapes count
            if len(shape.element.xpath('.//a:custGeom')):
                cust_shape_count += 1
            # Rectangles count
            elif shape.shape_type == MSO_SHAPE.RECTANGLE:
                rect_order.append(shape.text)
                rect_width.append(shape.width)
        # Comparing rectangle's plot order
        eq_(get_rect_order, rect_order)
        # Comparing rectangle's width
        eq_(get_rect_width, rect_width)
        # Comparing custom shape's count
        eq_(total_cust_shapes, cust_shape_count)

    def test_calendarmap(self):
        # Test case for calendarmap
        # If shape is not a rectangle
        with assert_raises(NotImplementedError):
            pptgen.pptgen(
                source=self.input, only=25,
                data={'data': {}},
                drawcalendarmap={'Invalid Input': {'calendarmap': {'data': 'data'}}})

        width, slidenumber, periods, weekstart = 34, 24, 150, 2
        data = pd.DataFrame({'date': pd.date_range('1/1/2017', periods=periods, freq='D')})
        data['column'] = 100. * np.random.rand(len(data))
        data = data.sort_values(by=['date']).set_index('date')
        label = 40
        labels = [[0, 0], [0, label], [label, 0], [label, label]]
        for case in labels:
            left, top = case
            target = pptgen.pptgen(
                source=self.input, only=slidenumber,
                data={'data': data},
                drawcalendar={
                    'Calendar Rectangle': {
                        'calendarmap': {
                            'data': 'data["data"]["column"]',
                            'width': width,
                            'weekstart': weekstart,
                            'label_left': left,
                            'label_top': top,
                            'text-color': '#000000',
                            'startdate': 'data.index[0]'
                        }
                    }
                })
            startweekday = (data.index[0].weekday() - weekstart) % 7
            eq_(len(target.slides), 1)
            texts_ex = data.index.strftime('%d').tolist()
            weekly_mean, weekday_mean = [], []
            if top:
                weekly_mean = pd.Series([data['column'][max(0, x):x + 7].mean()
                                        for x in range(-startweekday, len(data['column']), 7)])
                frmt = pptgen.utils.decimals(weekly_mean.values)
                weekly_mean = weekly_mean.map(('{:,.%df}' % frmt).format).tolist()
            if left:
                weekday_mean = pd.Series([data['column'][(x - startweekday) % 7::7].mean()
                                         for x in range(7)])
                frmt = pptgen.utils.decimals(weekday_mean.values)
                weekday_mean = weekday_mean.map(('{:,.%df}' % frmt).format).tolist()
            gradient = matplotlib.cm.get_cmap('RdYlGn')
            colors_ex = [matplotlib.colors.to_hex(gradient(x))
                         for x in pptgen.utils.scale(data['column'])]
            texts_ppt, colors_ppt, toplabels_ppt, leftlabels_ppt = [], [], [], []
            for shape in target.slides[0].shapes:
                # Rect cells text and color
                if shape.name.startswith('Rectangle'):
                    texts_ppt.append(shape.text)
                    colors_ppt.append('#{}'.format(shape.fill.fore_color.rgb).lower())
                # Top bar labels
                elif shape.name == 'summary.top.label':
                    toplabels_ppt.append(shape.text)
                # Left bar labels
                elif shape.name == 'summary.left.label':
                    leftlabels_ppt.append(shape.text)
            # Comparing text's order(which is day)
            eq_(texts_ex, texts_ppt)
            # Caparing rect colors
            eq_(colors_ex, colors_ppt)
            # Comparing top bar chart's labels
            eq_(weekly_mean, toplabels_ppt)
            # Comparing left bar chart's labels
            eq_(weekday_mean, leftlabels_ppt)

    @classmethod
    def tearDown(cls):
        # Teardown classmethod.
        if os.path.exists(cls.output):
            os.remove(cls.output)
