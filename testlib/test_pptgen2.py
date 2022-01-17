import io
import gramex.data
import numpy as np
import os
import pandas as pd
import pptx
from gramex.config import objectpath
from gramex.pptgen2 import pptgen, load_data, commands, commandline
from nose.tools import eq_, ok_, assert_raises
from orderedattrdict import AttrDict
from pptx import Presentation
from pptx.dml.color import _NoneColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR as MVA
from pptx.oxml.ns import _nsmap, qn
from testfixtures import LogCapture, OutputCapture
from unittest import TestCase
from . import folder, sales_file, afe

units = ('inches', 'cm', 'mm', 'pt', 'emu', 'centipoints')
aeq_ = lambda a, b: ok_(abs(a - b) <= 1)        # noqa


class TestPPTGen(TestCase):
    # Test Case module for pptgen

    @classmethod
    def setUp(cls):
        # Setup class method to initialize common variables.
        np.random.seed(0)
        cls.input = os.path.join(folder, 'input2.pptx')
        cls.prs = Presentation(cls.input)
        cls.output = os.path.join(folder, 'output.pptx')
        cls.image = os.path.join(folder, 'small-image.jpg')
        cls.data = pd.read_excel(sales_file, engine='openpyxl')
        if os.path.exists(cls.output):
            os.remove(cls.output)

    @classmethod
    def remove_output(cls):
        if os.path.exists(cls.output):
            os.remove(cls.output)

    def get_shapes(self, shapes, name):
        '''Return shapes with given name among shapes. Find shapes in groups using ">"'''
        names = name.split('>')
        if len(names) == 1:
            return [shape for shape in shapes if shape.name == name]
        else:
            for name in names[:-1]:
                shapes = sum((tuple(s.shapes) for s in self.get_shapes(shapes, name)), ())
            return self.get_shapes(shapes, names[-1])

    def get_shape(self, shapes, name):
        try:
            return self.get_shapes(shapes, name)[0]
        except IndexError:
            raise ValueError('Cannot find shape %s among %r' % (name, [s.name for s in shapes]))

    def check_opacity(self, colorformat, opacity):
        str_val = colorformat._xFill.xpath('.//a:alpha')[0].values()[0]
        eq_(commands.ST_Percentage.convert_from_xml(str_val), opacity)

    def test_data(self):
        # Empty dict is returned as-is
        eq_(load_data({}), {})
        # [non-str, non-dict] datasets are loaded as-is
        vals = [(eq_, None), (eq_, 1), (eq_, [1, 2]), (afe, self.data)]
        for test_fn, val in vals:
            test_fn(load_data(val)['data'], val)
        for test_fn, val in vals:
            test_fn(load_data({'key': val})['key'], val)
        # Strings can be treated as functions with _default_key=function
        eq_(load_data('[1, 2]', _default_key='function'), {'data': [1, 2]})
        eq_(load_data({'key': '[1, 2]'}, _default_key='function'), {'key': [1, 2]})
        # Strings can be treated as URLs with _default_key=url
        sales_data = gramex.data.filter(sales_file)
        afe(load_data(sales_file, _default_key='url')['data'], sales_data)
        afe(load_data({'key': sales_file}, _default_key='url')['key'], sales_data)
        # Strings raise an exception without _default_key
        with assert_raises(Exception):
            load_data('text')
        with assert_raises(Exception):
            load_data({'key': 'text'})
        # Dicts with url: are processed via gramex.data.filter
        afe(load_data({'url': sales_file})['data'], sales_data)
        afe(load_data({'d': {'url': sales_file}})['d'], sales_data)
        transform = 'data.set_index(["देश", "city", "product"])'
        afe(load_data({'d': {'url': sales_file, 'transform': transform}})['d'],
            sales_data.set_index(['देश', 'city', 'product']))
        afe(load_data({'d': {'url': sales_file, 'args': {'product': 'Eggs'}}})['d'],
            gramex.data.filter(sales_file, args={'product': 'Eggs'}))
        # Dicts with function: are executed
        afe(load_data({'function': 'gramex.cache.open(%r)' % sales_file})['data'], sales_data)
        afe(load_data({'d': {'function': 'gramex.cache.open(%r)' % sales_file}})['d'], sales_data)
        eq_(load_data({'d': {'function': 'str(handler)'}}, handler='abc')['d'], 'abc')
        # Functions can chain data keys, and also accept kwargs
        eq_(load_data({'x': 'a', 'y': 'x'}, a=1, _default_key='function')['y'], 1)
        eq_(load_data({'x': 'a + 1', 'y': 'x + 1'}, a=1, _default_key='function')['y'], 3)
        # kwargs are overridden by data unless they're None
        eq_(load_data(2, data=1)['data'], 2)
        eq_(load_data({'x': 2}, x=1)['x'], 2)
        eq_(load_data(None, data=1)['data'], 1)
        eq_(load_data({'x': None}, x=1)['x'], 1)

    def test_expr(self):
        # Test expr mode - when strings are expressions
        t = lambda v: commands.expr(v, data={'_expr_mode': True, 'x': 1})    # noqa
        eq_(t('x + 0'), 1)                      # Value is a variable
        eq_(t('"x + 0"'), 'x + 0')              # String is a literal
        eq_(t('"{x + 0}"'), '{x + 0}')          # String is a literal
        eq_(t('f"x + 0"'), 'x + 0')             # f-string is a template
        eq_(t('f"{x + 0}"'), '1')               # f-string is a template using data
        for val in [None, True, 1, []]:         # Non-string returns as-is
            eq_(t(val), val)
        afe(t(self.data), self.data)
        eq_(t({'value': 'x + 0'}), 'x + 0')     # value: is a template
        # NOTE: literal cannot process x + 0. Just test with x
        eq_(t({'value': '{x}'}), '1')           # value: is a template using data
        for val in [None, 1, [], {}]:           # value: non-string returns as-is
            eq_(t({'value': val}), val)
        afe(t({'value': self.data}), self.data)

        # Test literal mode - when strings are values
        t = lambda v: commands.expr(v, data={'_expr_mode': False, 'x': 1})   # noqa
        eq_(t('x + 0'), 'x + 0')                    # Value is a literal
        eq_(t('"x + 0"'), '"x + 0"')                # String value is a literal
        # NOTE: literal cannot process expressions like x + 0. Just test with x
        eq_(t('{x}'), '1')                          # String template is formatted
        eq_(t('f"x + 0"'), 'f"x + 0"')              # f-string value is a literal
        for val in [None, True, 1, []]:             # Non-string returns as-is
            eq_(t(val), val)
        afe(t(self.data), self.data)
        eq_(t({'expr': 'x + 0'}), 1)                # expr: is a variable
        eq_(t({'expr': '"{x + 0}"'}), '{x + 0}')    # expr: quoted string becomes string literal
        eq_(t({'expr': 'f"{x + 0}"'}), '1')         # expr: f-string is a template using data
        for val in [None, 1, [], {}]:               # expr: non-string returns as-is
            eq_(t({'expr': val}), val)

    def test_length(self):
        length = commands.length
        eq_(length(3.2), pptx.util.Inches(3.2))
        eq_(length(np.int64(3)), pptx.util.Inches(3))
        for unit in ('', '"', 'in', 'inch'):
            eq_(length('3.2' + unit), pptx.util.Inches(3.2))
            eq_(length('3.2  ' + unit), pptx.util.Inches(3.2))
        for unit in ('cp', 'centipoint'):
            eq_(length('3.2' + unit), pptx.util.Centipoints(3.2))
            eq_(length('3.2  ' + unit), pptx.util.Centipoints(3.2))
        for unit in units:
            eq_(length('3.2' + unit), getattr(pptx.util, unit.title())(3.2))
            eq_(length('3.2  ' + unit), getattr(pptx.util, unit.title())(3.2))
            eq_(length('+3.2 ' + unit), getattr(pptx.util, unit.title())(3.2))
            eq_(length('-3.2 ' + unit), getattr(pptx.util, unit.title())(-3.2))
        with assert_raises(ValueError):
            length('3.4 nonunits')
        with assert_raises(ValueError):
            length(None)
        length_class = commands.length_class
        for unit in ('"', 'in', 'inch', 'inches', 'IN', 'Inch', 'INCHes', ''):
            eq_(length_class(unit), pptx.util.Inches)
        for unit in ('emu', 'Emu', 'EMU'):
            eq_(length_class(unit), pptx.util.Emu)
        for unit in ('cp', 'CentiPoint', 'CENTIPoints'):
            eq_(length_class(unit), pptx.util.Centipoints)
        with assert_raises(ValueError):
            eq_(length_class('nonunits'))

    def test_unit(self):
        rule = {'Title 1': {'width': 10}}
        for unit in units:
            prs = pptgen(source=self.input, only=1, rules=[rule], unit=unit)
            eq_(commands.length_unit.__name__, unit.title())
            eq_(prs.slides[0].shapes[0].width, commands.length_unit(10))

    def test_register(self, slides=3):
        # register= must be a dict
        with assert_raises(TypeError):
            pptgen(source=self.input, only=1, register='dummy')
        # register= compiles the functions into commands.cmdlist
        prs = pptgen(source=self.input, only=slides, register={
            'cmd1': '(shape, spec, data)',
            'cmd2': 'shape.get(spec, data)',
            'rename': 'setattr(shape, "name", spec)',
            'rotate': 'setattr(shape, "rotation", spec)',
        }, rules=[
            {'Rectangle 1': {'rotate': 45, 'rename': 'abc'}}
        ])
        ok_('cmd1' in commands.cmdlist)
        eq_(commands.cmdlist['cmd1'](shape=1, spec={}), (1, {}, None))
        ok_('cmd2' in commands.cmdlist)
        eq_(commands.cmdlist['cmd2'](shape={}, spec='x', data='y'), 'y')
        shape = self.get_shape(prs.slides[0].shapes, 'abc')
        eq_(shape.rotation, 45)

    def test_only(self, slides=[2, 4]):
        # Delete slides except those specified in ``only``
        with assert_raises(TypeError):
            pptgen(source=self.input, only={})
        with assert_raises(TypeError):
            pptgen(source=self.input, only='4')
        # Test single only= value
        only = slides[0]
        prs = pptgen(source=self.input, only=only)
        eq_(len(prs.slides), 1)
        eq_(prs.slides[0].shapes.title.text, self.prs.slides[only - 1].shapes.title.text)
        # Test multiple only= value
        only = slides
        prs = pptgen(source=self.input, only=only)
        eq_(len(prs.slides), len(only))
        for i, slide in enumerate(only):
            eq_(prs.slides[i].shapes.title.text, self.prs.slides[slide - 1].shapes.title.text)

    def test_target(self, slides=1):
        # pptgen returns target presentation
        prs = pptgen(source=self.input, only=slides)
        eq_(len(prs.slides), 1)
        eq_(prs.slides[0].shapes.title.text, self.prs.slides[0].shapes.title.text)
        # pptgen ALSO saves at target= if it is specified
        prs = pptgen(source=self.input, target=self.output, only=slides)
        eq_(len(prs.slides), 1)
        eq_(prs.slides[0].shapes.title.text, self.prs.slides[0].shapes.title.text)
        prs = Presentation(self.output)
        eq_(len(prs.slides), 1)
        eq_(prs.slides[0].shapes.title.text, self.prs.slides[0].shapes.title.text)

    def test_incorrect(self, slides=1):
        with LogCapture() as logs:
            pptgen(source=self.input, only=slides, rules=[
                {'No-Shape': {'left': 0}},
                {'Title 1': {'no-command': 0}}
            ])
        logs.check_present(
            ('gramex', 'WARNING', 'pptgen2: No shape matches pattern: No-Shape'),
            ('gramex', 'WARNING', 'pptgen2: Unknown command: no-command on shape: Title 1')
        )

    def test_slide_filter(self, slides=[1, 2, 3]):
        # Rules are specified as rule-name={shape: {rule}, ...}
        data = {'x': [2, 3]}
        rule1 = {'slide-number': 1, 'Title 1': {'width': 10}}
        rule2 = {'slide-number': {'expr': 'x'},
                 'Title 1': {'width': 20},
                 'Rectangle 1': {'width': 20}}
        prs = pptgen(source=self.input, only=slides, data=data, rules=[rule1, rule2])
        eq_(self.get_shape(prs.slides[0].shapes, 'Title 1').width, pptx.util.Inches(10))
        eq_(self.get_shape(prs.slides[1].shapes, 'Title 1').width, pptx.util.Inches(20))
        eq_(self.get_shape(prs.slides[2].shapes, 'Title 1').width, pptx.util.Inches(20))
        eq_(self.get_shape(prs.slides[2].shapes, 'Rectangle 1').width, pptx.util.Inches(20))
        rule1 = {'slide-title': '*pptx*', 'Title 1': {'width': 10}}
        rule2 = {'slide-title': ['*pos*', '*pptx*'], 'Rectangle 1': {'width': 20}}
        prs = pptgen(source=self.input, only=slides, rules=[rule1, rule2])
        eq_(self.get_shape(prs.slides[0].shapes, 'Title 1').width, pptx.util.Inches(10))
        eq_(self.get_shape(prs.slides[2].shapes, 'Rectangle 1').width, pptx.util.Inches(20))
        with LogCapture() as logs:
            pptgen(source=self.input, only=slides, rules=[{'slide-number': 5}])
        logs.check_present(
            ('gramex', 'WARNING', 'pptgen2: No slide with slide-number: 5, slide-title: None'),
        )

    def test_transition(self, slides=[1, 2, 3]):
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'slide-number': 1, 'transition': 'glitter'},
            {'slide-number': 2, 'transition': {'type': 'morph', 'duration': 1}},
            {'slide-number': 3, 'transition': {'type': 'wind left', 'advance': 2}},
        ])
        node = 'mc:AlternateContent/mc:Choice/p:transition'

        tr = prs.slides[0].element.find(node, _nsmap)
        eq_(tr.attrib['{%s}dur' % _nsmap['p14']], '300')
        el = tr.find('p14:glitter', _nsmap)
        eq_(el.attrib, {'pattern': 'hexagon', 'dir': 'l'})

        tr = prs.slides[1].element.find(node, _nsmap)
        eq_(tr.attrib['{%s}dur' % _nsmap['p14']], '1000')
        el = tr.find('p159:morph', _nsmap)
        eq_(el.attrib, {'option': 'byObject'})
        el = prs.slides[1].element.find(node, _nsmap)

        tr = prs.slides[2].element.find(node, _nsmap)
        eq_(tr.attrib['advTm'], '2000')
        el = tr.find('p15:prstTrans', _nsmap)
        eq_(el.attrib, {'prst': 'wind', 'invX': '1'})

        prs = pptgen(source=prs, target=self.output, rules=[{'transition': 'none'}])
        for slide in prs.slides:
            eq_(slide.element.find(node, _nsmap), None)

    def test_normalize_group(self, slides=3):
        def coords(grp):
            c = grp.element.find(qn('p:grpSpPr')).find(qn('a:xfrm'))
            return AttrDict(off=c.find(qn('a:off')),
                            ext=c.find(qn('a:ext')),
                            choff=c.find(qn('a:chOff')),
                            chext=c.find(qn('a:chExt')))
        grp = self.get_shape(self.prs.slides[slides - 1].shapes, 'Group 2')
        subgrps = self.get_shapes(grp.shapes, 'SubGroup')
        for g in [grp] + subgrps:
            c = coords(g)
            assert c.off.x != c.choff.x, 'x offset is initially different'
            assert c.off.y != c.choff.y, 'y offset is initially different'
            assert c.ext.cx != c.chext.cx, 'width is initially different'
            assert c.ext.cy != c.chext.cy, 'height is initially different'
        # Just opening via pptgen normalizes the groups
        prs = pptgen(source=self.input, target=self.output)
        grp = self.get_shape(prs.slides[slides - 1].shapes, 'Group 2')
        subgrps = self.get_shapes(grp.shapes, 'SubGroup')
        for g in [grp] + subgrps:
            c = coords(g)
            assert c.off.x == c.choff.x, 'x offset is same after normalization'
            assert c.off.y == c.choff.y, 'y offset is same after normalization'
            assert c.ext.cx == c.chext.cx, 'width is same after normalization'
            assert c.ext.cy == c.chext.cy, 'height is same after normalization'

    def test_shape_names(self, slides=3):
        prs = pptgen(source=self.input, only=slides, rules=[
            {'group 1': {'left': 99}},          # Case-sensitive match is ignored
            {'Group ?': {'left': 1}},           # Group 1, Group 2
            {'?extBox ?': {'left': 2}},         # TextBox 1
            {'*le 1': {'left': 3}},             # Title 1, Rectangle 1
            {'*form*': {'left': 4}},            # Freeform 1
            {'[BC]har[tu] 1': {'left': 5}},     # Chart 1
        ])
        eq_(self.get_shape(prs.slides[0].shapes, 'Group 1').left, pptx.util.Inches(1))
        eq_(self.get_shape(prs.slides[0].shapes, 'Group 2').left, pptx.util.Inches(1))
        eq_(self.get_shape(prs.slides[0].shapes, 'TextBox 1').left, pptx.util.Inches(2))
        eq_(self.get_shape(prs.slides[0].shapes, 'Title 1').left, pptx.util.Inches(3))
        eq_(self.get_shape(prs.slides[0].shapes, 'Rectangle 1').left, pptx.util.Inches(3))
        eq_(self.get_shape(prs.slides[0].shapes, 'Freeform 1').left, pptx.util.Inches(4))
        eq_(self.get_shape(prs.slides[0].shapes, 'Chart 1').left, pptx.util.Inches(5))

    def test_name_position(self, slides=3):
        pos = {'width': 4, 'height': 3, 'top': 2, 'left': 1, 'rotation': 30,
               'name': {'expr': 'shape.name + " X"'}}
        add = {'add-width': -0.1, 'add-height': 0.05, 'add-top': 0.3, 'add-left': -0.1,
               'add-rotation': 30, 'name': {'expr': 'shape.name + " X"'}}
        for name in ['Rectangle 1', 'TextBox 1', 'Picture 1', 'Chart 1', 'Group 1', 'Table 1',
                     'Diagram 1', 'Audio 1', 'Freeform 1', 'Word Art 1']:
            # 'Zoom 1', 'Equation 1' are not supported by python-pptx
            prs = pptgen(source=self.input, only=slides, rules=[{name: pos}])
            shp = self.get_shape(prs.slides[0].shapes, name + ' X')
            for attr, val in pos.items():
                if attr != 'name':
                    convert = float if attr == 'rotation' else commands.length_unit
                    eq_(getattr(shp, attr), convert(val))

            prs = pptgen(source=self.input, only=slides, rules=[{name: add}])
            shp = self.get_shape(prs.slides[0].shapes, name + ' X')
            src = self.get_shape(self.prs.slides[slides - 1].shapes, name)
            for cmd, val in add.items():
                if attr != 'name':
                    attr = cmd.split('-')[-1]
                    convert = float if cmd == 'add-rotation' else commands.length_unit
                    eq_(getattr(shp, attr), convert(val) + getattr(src, attr))

            for zoom in (0.6, 1.2):
                prs = pptgen(source=self.input, only=slides, rules=[{name: {'zoom': zoom}}])
                shp = self.get_shape(prs.slides[0].shapes, name)
                src = self.get_shape(self.prs.slides[slides - 1].shapes, name)
                aeq_(shp.left, int(src.left - (zoom - 1) * src.width / 2))
                aeq_(shp.top, int(src.top - (zoom - 1) * src.height / 2))
                aeq_(shp.width, int(src.width * zoom))
                aeq_(shp.height, int(src.height * zoom))

        # Adjust position within group and subgroups
        text_pos = {'left': 1, 'top': 1, 'width': 2, 'height': 0.5}
        img_pos = {'left': 0, 'top': 1, 'width': 0.5, 'height': 0.5}
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'Group 2': {
                'left': 1,
                'Subgroup': {'left': 1, 'Text': text_pos, 'Picture': img_pos}
            }}
        ])
        eq_(self.get_shape(prs.slides[0].shapes, 'Group 2>Subgroup').left, pptx.util.Inches(1))
        shape = self.get_shape(prs.slides[0].shapes, 'Group 2>Subgroup>Text')
        for attr, val in text_pos.items():
            eq_(getattr(shape, attr), pptx.util.Inches(val))
        shape = self.get_shape(prs.slides[0].shapes, 'Group 2>Subgroup>Picture')
        for attr, val in img_pos.items():
            eq_(getattr(shape, attr), pptx.util.Inches(val))

    def test_image(self, slides=3):
        shape = self.get_shape(self.prs.slides[slides - 1].shapes, 'Picture 1')
        width = shape.width
        for img, aspect in (('small-image.jpg', 1), ('small-image.png', 2)):
            path = os.path.join(folder, img)
            prs = pptgen(source=self.input, only=slides,
                         rules=[{'Picture 1': {'image': path}}])
            shape = self.get_shape(prs.slides[0].shapes, 'Picture 1')
            rid = shape._pic.blipFill.blip.rEmbed
            part = shape.part.related_parts[rid]
            with open(path, 'rb') as handle:
                eq_(part.blob, handle.read())
            eq_(shape.width, width)
            self.assertAlmostEqual(shape.width / shape.height, aspect, places=5)

    def test_image_width_height(self, slides=3):
        shape = self.get_shape(self.prs.slides[slides - 1].shapes, 'Picture 1')
        aspect = shape.width / shape.height
        for size in (3, '3 inches', '7.62 cm', '76.2 mm', '216 pt', '2743200 emu', '21600 cp'):
            # image-width preserves aspect ratio
            prs = pptgen(source=self.input, only=slides,
                         rules=[{'Picture 1': {'image-width': size}}])
            shape = self.get_shape(prs.slides[0].shapes, 'Picture 1')
            eq_(shape.width, pptx.util.Inches(3))
            self.assertAlmostEqual(shape.width / shape.height, aspect, places=5)
            # image-height preserves aspect ratio
            prs = pptgen(source=self.input, only=slides,
                         rules=[{'Picture 1': {'image-height': size}}])
            shape = self.get_shape(prs.slides[0].shapes, 'Picture 1')
            eq_(shape.height, pptx.util.Inches(3))
            self.assertAlmostEqual(shape.width / shape.height, aspect, places=5)

    def test_image_clone_copy(self, slides=7):
        repeat = (0, 1, 2)
        prs = pptgen(source=self.input, rules=[
            {'copy-slide': repeat, 'slide-number': slides, 'Group 1': {'Picture': {
                'clone-shape': repeat,
                'image': self.image,
            }}}])
        for index in repeat:
            slide = prs.slides[slides + index - 1]
            group = [shape for shape in slide.shapes if shape.name == 'Group 1'][0]
            pics = [shape for shape in group.shapes if shape.name == 'Picture']
            eq_(len(pics), len(repeat))
            for pic in pics[1:]:
                pic.width == pics[0].width
                pic.height == pics[0].height

    def test_fill_stroke(self, slides=3):
        colors = (
            ('red', {'rgb': (255, 0, 0)}),
            ('#f00', {'rgb': (255, 0, 0)}),
            ('#ff0000', {'rgb': (255, 0, 0)}),
            ('rgb(255, 0, 0)', {'rgb': (255, 0, 0)}),
            ((255, 0, 0), {'rgb': (255, 0, 0)}),
            ([255, 0, 0], {'rgb': (255, 0, 0)}),
            ((1.0, 0.5, 0), {'rgb': (255, 128, 0)}),
            ([1.0, 0.5, 0], {'rgb': (255, 128, 0)}),
            ('ACCENT_1', {'theme_color': MSO_THEME_COLOR.ACCENT_1, 'brightness': 0}),
            ('ACCENT_2+40', {'theme_color': MSO_THEME_COLOR.ACCENT_2, 'brightness': 0.4}),
            ('ACCENT_3-20', {'theme_color': MSO_THEME_COLOR.ACCENT_3, 'brightness': -0.2}),
        )
        for color, result in colors:
            for name in ['TextBox 1', 'Rectangle 1', 'Word Art 1', 'Freeform 1']:
                # Doesn't work for 'Group 1', 'Table 1', 'Audio 1', 'Chart 1', 'Diagram 1'
                prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
                    {name: {
                        'fill': color,
                        'stroke': color,
                        'fill-opacity': 0.5,
                        'stroke-opacity': 0.4,
                        'stroke-width': '1 pt',
                    }}
                ])
                shape = self.get_shape(prs.slides[0].shapes, name)
                for key in (('fill.fore_color.', 'line.fill.fore_color.')):
                    for attr, val in result.items():
                        if attr == 'brightness':
                            self.assertAlmostEqual(objectpath(shape, key + attr), val, places=5)
                        else:
                            eq_(objectpath(shape, key + attr), val)
                self.check_opacity(shape.fill.fore_color, 0.5)
                self.check_opacity(shape.line.fill.fore_color, 0.4)
                eq_(shape.line.width, pptx.util.Pt(1))
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'Rectangle 1': {'fill': 'none', 'stroke': 'none'}}
        ])
        shape = self.get_shape(prs.slides[0].shapes, 'Rectangle 1')
        eq_(shape.fill.type, MSO_FILL.BACKGROUND)
        eq_(shape.line.fill.type, MSO_FILL.BACKGROUND)

    def test_clone_shape(self, slides=3):
        data = {'a': -0.5, 'b': 1, 'c': 2.5}
        clone = {'clone-shape': data, 'top': 1, 'add-top': {'expr': 'clone.val'}}
        # Clone shapes
        for name in ['TextBox 1', 'Group 1', 'Table 1', 'Audio 1', 'Freeform 1', 'Word Art 1']:
            # TODO: 'Chart 1', 'Diagram 1' don't work yet
            prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{name: clone}])
            shapes = [shape for shape in prs.slides[0].shapes if shape.name == name]
            eq_(shapes[0].top, pptx.util.Inches(0.5))
            eq_(shapes[1].top, pptx.util.Inches(2.0))
            eq_(shapes[2].top, pptx.util.Inches(3.5))
        # Clone groups
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'Group 1': {
                'clone-shape': data,
                'data': {'myclone': 'clone'},
                'top': {'expr': '1 + myclone.val'},
                'Picture': {
                    'clone-shape': data,
                    'data': {'subclone': 'clone'},
                    'left': {'expr': '1 + subclone.val / 2'},
                    'image-width': 0.2,
                },
                'Caption': {
                    'clone-shape': data,
                    'data': {'subclone2': 'clone'},
                    'left': {'expr': '1 + subclone2.val / 2'},
                    'text': '{clone.pos}, {clone.key}, {clone.val}, {clone.shape.text}, ' +
                            '{clone.parent.key}, {clone.parent.val}',
                    'fill': 'red',
                }
            }}
        ])
        groups = self.get_shapes(prs.slides[0].shapes, 'Group 1')
        picture = self.get_shapes(prs.slides[0].shapes, 'Group 1>Picture')
        caption = self.get_shapes(prs.slides[0].shapes, 'Group 1>Caption')
        n = len(data)
        for i, (ik, iv) in enumerate(data.items()):
            eq_(groups[i].top, pptx.util.Inches(1 + iv))
            for j, (jk, jv) in enumerate(data.items()):
                eq_(picture[i * n + j].left, pptx.util.Inches(1 + jv / 2))
                eq_(picture[i * n + j].width, pptx.util.Inches(0.2))
                eq_(caption[i * n + j].left, pptx.util.Inches(1 + jv / 2))
                eq_(caption[i * n + j].fill.fore_color.rgb, (255, 0, 0))
                eq_(caption[i * n + j].text, f'{j}, {jk}, {jv}, Grouped image, {ik}, {iv}')

    def test_text(self, slides=4):
        # Non-strings are converted to str
        for val in (1, ['x'], ):
            prs = pptgen(source=self.input, only=slides, rules=[
                {'TextBox 1': {'text': str(val)}}])
            shape = self.get_shape(prs.slides[0].shapes, 'TextBox 1')
            eq_(shape.text, str(val))
        # Empty strings clear text
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'TextBox 1': {'text': ''}}])
        shape = self.get_shape(prs.slides[0].shapes, 'TextBox 1')
        eq_(shape.text, '')
        # Unicode characters work
        text = '高σ高λس►'
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'TextBox 1': {'text': text}}])
        shape = self.get_shape(prs.slides[0].shapes, 'TextBox 1')
        eq_(shape.text, text)
        # Para and run formatting works
        text = '''P0R0 <a>P0R1</a> P0R2 <a>P0R3</a> P0R4
            <p align="left" bold="y" color="#ff0000" font-name="Arial" font-size="8 pt" italic="y"
                level="0" line-spacing="3 pt" space-before="16 pt" space-after="20 pt"
                underline="y"
            >
                P1R0
                <a baseline="superscript" bold="n" color="#00ff00" font-name="Calibri"
                    font-size="18 pt" italic="n" strike="double" underline="n"> P1R1 </a>
                P1R2
            </p>
            <p>P2R0
                <a bold="y" baseline="-35%" strike="single"> P2R1 </a>
                <a baseline="subscript" strike="none"> P2R2 </a>
                P2R3</p>
            P3R0 <a color="#00f">P3R1</a> P3R2
        '''
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'TextBox 1': {'text': text}}])
        shape = self.get_shape(prs.slides[0].shapes, 'TextBox 1')
        eq_(shape.text.split(),
            'P0R0 P0R1 P0R2 P0R3 P0R4 P1R0 P1R1 P1R2 P2R0 P2R1 P2R2 P2R3 P3R0 P3R1 P3R2'.split())
        paras = shape.text_frame.paragraphs
        # Para 0 has the right attributes and text
        eq_(paras[0].text, 'P0R0 P0R1 P0R2 P0R3 P0R4 ')
        eq_(paras[0].runs[0].text, 'P0R0 ')
        # Para 0 attributes are preserved
        eq_(paras[0].level, 0)
        eq_(paras[0].alignment, PP_ALIGN.CENTER)
        eq_(paras[0].runs[0].font.size, pptx.util.Pt(28))
        eq_(paras[0].runs[0].font.name, 'Consolas')
        for attr in ('line_spacing', 'space_after', 'space_before'):
            eq_(getattr(paras[0], attr), None)
        # Para 1 has the right attributes and text
        eq_(paras[1].text.split(), 'P1R0 P1R1 P1R2'.split())
        eq_(paras[1].alignment, PP_ALIGN.LEFT)
        eq_(paras[1].font.bold, True)
        eq_(paras[1].font.color.rgb, (255, 0, 0))
        eq_(paras[1].font.name, 'Arial')
        eq_(paras[1].font.size, pptx.util.Pt(8))
        eq_(paras[1].font.italic, True)
        eq_(paras[1].level, 0)
        eq_(paras[1].line_spacing, pptx.util.Pt(3))
        eq_(paras[1].space_before, pptx.util.Pt(16))
        eq_(paras[1].space_after, pptx.util.Pt(20))
        eq_(paras[1].font.underline, True)
        eq_(paras[1].runs[0].text, ' P1R0 ')
        # Para 1 run 2 has the specified attributes
        eq_(paras[1].runs[1].text, ' P1R1 ')
        eq_(paras[1].runs[1].font.bold, False)
        eq_(paras[1].runs[1].font.color.rgb, (0, 255, 0))
        eq_(paras[1].runs[1].font.name, 'Calibri')
        eq_(paras[1].runs[1].font.size, pptx.util.Pt(18))
        eq_(paras[1].runs[1].font.italic, False)
        eq_(paras[1].runs[1].font.underline, False)
        eq_(paras[1].runs[1].font._rPr.get('baseline'), '30000')
        eq_(paras[1].runs[1].font._rPr.get('strike'), 'dblStrike')
        # Auto-created runs have no attributes
        for attr in ('bold', 'name', 'italic', 'underline'):
            eq_(getattr(paras[1].runs[0].font, attr), None)
            eq_(getattr(paras[1].runs[2].font, attr), None)
            eq_(getattr(paras[2].runs[0].font, attr), None)
            eq_(getattr(paras[2].runs[2].font, attr), None)
        # ... except font size, which is taken from the first run (set in the PPTX)
        eq_(paras[1].runs[0].font.size, pptx.util.Pt(28))
        eq_(paras[1].runs[2].font.size, pptx.util.Pt(28))
        eq_(paras[2].runs[0].font.size, pptx.util.Pt(28))
        eq_(paras[2].runs[2].font.size, pptx.util.Pt(28))
        # Para 2 runs have the right attrs
        eq_(paras[2].runs[1].text, ' P2R1 ')
        eq_(paras[2].runs[1].font.bold, True)
        eq_(paras[2].runs[1].font._rPr.get('baseline'), '-35000')
        eq_(paras[2].runs[1].font._rPr.get('strike'), 'sngStrike')
        eq_(paras[2].runs[2].text, ' P2R2 ')
        eq_(paras[2].runs[2].font._rPr.get('baseline'), '-25000')
        eq_(paras[2].runs[2].font._rPr.get('strike'), 'noStrike')
        # Para 3: runs are auto-wrapped into paras
        eq_(paras[3].runs[0].text, ' P3R0 ')
        eq_(paras[3].runs[1].text, 'P3R1')
        eq_(paras[3].runs[1].font.color.rgb, (0, 0, 255))
        eq_(paras[3].runs[2].text, ' P3R2 ')

    def test_text_style(self, slides=4):
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'TextBox 1': {
                'bold': 0,
                'italic': 1,
                'underline': 0,
                'color': 'blue',
                'font-name': 'Calibri',
                'font-size': '10 pt',
            }}
        ])
        shape = self.get_shape(prs.slides[0].shapes, 'TextBox 1')
        for para in shape.text_frame.paragraphs:
            eq_(para.font.bold, False)
            eq_(para.font.fill.fore_color.rgb, (0, 0, 255))
            eq_(para.font.italic, True)
            eq_(para.font.name, 'Calibri')
            eq_(para.font.size, pptx.util.Pt(10))
            eq_(para.font.underline, False)
            for run in para.runs:
                eq_(run.font.bold, None)
                # PPT needs colors on runs too, not only paras
                eq_(run.font.fill.fore_color.rgb, (0, 0, 255))
                eq_(run.font.italic, None)
                eq_(run.font.name, None)
                eq_(run.font.size, None)
                eq_(run.font.underline, None)

    def test_replace(self, slides=4):
        with assert_raises(ValueError):
            pptgen(source=self.input, only=slides, rules=[{'TextBox 1': {'replace': 'text'}}])
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'TextBox 1': {
                'replace': {
                    '[Oo]ld': 'Old1',
                    '(Title|italic)': '<a underline="y" bold="n">New</a>',
                    'title': 'ignored',
                    'der': 'd<a bold="y" underline="n">E<a color="green">R</a>',
                    'c.l.r': '<a font-size="18 pt" font-name="Arial">COLOR</a>',
                }
            }}
        ])
        defaults = {'font-name': 'Consolas', 'font-size': '28 pt'}
        expected_runs = [
            {'text': 'Old1', **defaults},
            {'text': ' ', **defaults},
            {'text': 'New', 'underline': True, **defaults},
            {'text': ' ', **defaults},
            {'text': 'un', 'underline': True, **defaults},
            {'text': 'd', 'underline': True, **defaults},
            {'text': 'E', 'bold': True, **defaults},
            {'text': 'R', 'color': (0, 128, 0), 'underline': True, **defaults},
            {'text': 'line', 'underline': True, **defaults},
            {'text': ' ', **defaults},
            {'text': 'New', 'italic': True, 'underline': True, **defaults},
            {'text': ' ', **defaults},
            {'text': 'COLOR', 'color': (255, 0, 0), 'font-size': '18 pt', 'font-name': 'Arial'},
        ]
        shape = self.get_shape(prs.slides[0].shapes, 'TextBox 1')
        for expected, actual in zip(expected_runs, shape.text_frame.paragraphs[0].runs):
            eq_(expected['text'], actual.text)
            eq_(expected.get('bold', False), bool(actual.font.bold))
            eq_(expected.get('italic', False), bool(actual.font.italic))
            eq_(expected.get('underline', False), bool(actual.font.underline))
            if 'color' in expected:
                eq_(expected['color'], actual.font.color.rgb)
            else:
                ok_(isinstance(actual.font.color._color, _NoneColor))
            eq_(expected.get('font-name', None), actual.font.name)
            if 'font-size' in expected:
                eq_(commands.length(expected['font-size']), actual.font.size)
            else:
                eq_(actual.font.size, None)

    def test_link_hover_tooltip(self, slides=[1, 2, 3, 4, 5, 6, 7], main_slide=5):
        prefixes = {'Link ': 'link', 'Hover ': 'hover', 'Has ': 'link', 'Tooltip ': 'tooltip'}
        vals = {
            'first': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=firstslide'},
            'last': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=lastslide'},
            'next': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=nextslide'},
            'previous': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=previousslide'},
            'prev': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=previousslide'},
            # 'end': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=endshow'},
            'back': {'target': '', 'action': 'ppaction://hlinkshowjump?jump=lastslideviewed'},
            'noaction': {'target': '', 'action': 'ppaction://noaction'},
            '1': {'target': self.prs.slides[1 - 1].shapes.title.text,
                  'action': 'ppaction://hlinksldjump'},
            '2': {'target': self.prs.slides[2 - 1].shapes.title.text,
                  'action': 'ppaction://hlinksldjump'},
            '4': {'target': self.prs.slides[4 - 1].shapes.title.text,
                  'action': 'ppaction://hlinksldjump'},
            'https://t.co/': {'target': 'https://t.co/', 'action': None},
            'file.pptx': {'target': 'file.pptx',
                          'action': 'ppaction://hlinkpres?slideindex=1&slidetitle='},
            'file.xlsx': {'target': 'file.xlsx', 'action': 'ppaction://hlinkfile'},
        }
        shape_rule = {prefix + val: {key: val}
                      for val in vals for prefix, key in prefixes.items()}
        text_rule = {prefix + 'Text': {
            'replace': {val + '$': f'<a {key}="{val}">{val}</a>' for val in vals}
        } for prefix, key in prefixes.items()}
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            shape_rule, text_rule])
        slide = prs.slides[main_slide - 1]
        shapes = slide.shapes
        for prefix, key in prefixes.items():
            for val, attr in vals.items():
                # Shape rules
                shape = self.get_shape(shapes, prefix + val)
                tag = 'a:hlinkClick' if key in {'link', 'tooltip'} else 'a:hlinkHover'
                self.check_link(slide, shape._element, tag, key, val, attr)
                # Text rules
                shape = self.get_shape(shapes, prefix + 'Text')
                tag = 'a:hlinkClick' if key in {'link', 'tooltip'} else 'a:hlinkMouseOver'
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text == val:
                            self.check_link(slide, run._r, tag, key, val, attr)

    def check_link(self, slide, el, tag, key, val, attr):
        link = el.find('.//' + tag, _nsmap)
        # Tooltip converts 'noaction' into 'next'. Handle that
        if key == 'tooltip' and val == 'noaction':
            action = 'ppaction://hlinkshowjump?jump=nextslide'
        else:
            action = attr['action']
        eq_(link.get('action'), action)
        rid = link.get(qn('r:id'))
        target = slide.part.rels[rid]._target if rid else ''
        if isinstance(target, pptx.parts.slide.SlidePart):
            target = target.slide.shapes.title.text
        eq_(target, attr['target'])

    def test_table(self, slides=9):
        data = self.data.head(10)       # The 10th row has NaNs. Ensure the row is included
        headers = ['<a color="red">देश</a>', 'city', '<p>prod</p><p>uct</p>', 'Sales']
        prs = pptgen(source=self.input, target=self.output, only=slides, mode='expr',
                     data={'data': data},
                     rules=[
                         {'Table 1': {'table': {'data': 'data', 'header-row': headers}}},
                         {'Table 2': {'table': {'data': 'data', 'width': 2}}},
                     ])
        for row_offset, shape_name in ((1, 'Table 1'), (0, 'Table 2')):
            table = self.get_shape(prs.slides[0].shapes, shape_name).table
            for i, (index, row) in enumerate(data.iterrows()):
                for j, (column, val) in enumerate(row.iteritems()):
                    cell = table.rows[i + row_offset].cells[j]
                    eq_(cell.text, '{}'.format(val))
        # Test table header
        header = self.get_shape(prs.slides[0].shapes, 'Table 1').table.rows[0].cells
        eq_(header[0].text, 'देश')
        eq_(header[0].text_frame.paragraphs[0].runs[0].font.color.rgb, (255, 0, 0))
        eq_(header[2].text_frame.paragraphs[0].text, 'prod')
        eq_(header[2].text_frame.paragraphs[1].text, 'uct')
        eq_(header[4].text, 'Table 1')      # Inherited from the template
        # Test column widths
        gridcols = self.get_shape(prs.slides[0].shapes, 'Table 2').table._tbl.tblGrid.gridCol_lst
        all(v.get('w') == pptx.util.Inches(2) for v in gridcols)

        # If there's no table data, text is copied from source
        prs = pptgen(source=self.input, target=self.output, only=slides, mode='expr', rules=[
            {'Table 2': {'table': {
                'header-row': True,
                'fill': '"red" if "Val" in cell.val else "yellow"',
            }}}
        ])
        table = self.get_shape(prs.slides[0].shapes, 'Table 2').table
        eq_(table.rows[1].cells[0].fill.fore_color.rgb, (255, 0, 0))
        eq_(table.rows[1].cells[1].fill.fore_color.rgb, (255, 255, 0))
        prs = pptgen(source=self.input, target=self.output, only=slides, mode='expr', rules=[
            {'Table 2': {'table': {
                'fill': '"red" if "Table" in cell.val else "yellow"',
            }}}
        ])
        table = self.get_shape(prs.slides[0].shapes, 'Table 2').table
        eq_(table.rows[0].cells[0].fill.fore_color.rgb, (255, 0, 0))
        eq_(table.rows[0].cells[1].fill.fore_color.rgb, (255, 255, 0))

        # Test all table commands comprehensively
        cmds = {'table': {
            'data': data,
            'header-row': False,
            'total-row': True,
            'first-column': True,
            'last-column': True,
            'width': {
                'देश': '1 in',
                'city': {'expr': '"2 in" if cell.column == "city" else "1 in"'},
                'product': {'expr': '"2 in" if cell.column == "city" else "1.5 in"'},
            },
            'align': {'expr': '"left" if cell.pos.row % 2 else "right"'},
            'bold': {'expr': 'cell.pos.row % 2'},
            'color': {'expr': '"red" if cell.pos.row % 3 else "green"'},
            'fill': {'expr': '"#eee" if cell.pos.row % 2 else "#ccc"'},
            'fill-opacity': 0.4,
            'font-name': 'Arial',
            'font-size': {'expr': '"10 pt" if cell.column == "देश" else "8 pt"'},
            'italic': {
                'देश': True,
                'city': {'expr': 'cell.pos.row % 2'},
            },
            'margin-left': {
                'देश': '0.05 in',
                'city': {'expr': '0 if cell.pos.column % 2 else "0.1 in"'},
            },
            'margin-right': '1 pt',
            'margin-top': {'expr': '0 if cell.pos.column % 2 else "0.1 in"'},
            'margin-bottom': 0,
            'underline': {'expr': 'cell.pos.column % 2'},
            'vertical-align': {
                'देश': 'middle',
                'city': {'expr': '"top" if cell.pos.row % 2 else "bottom"'},
            },
            # Add text: at the end to verify that it over-rides bold:, italic:, etc
            'text': '{cell.pos.row} {cell.pos.column} <a italic="y">{cell.index}</a> ' +
                    '{cell.column} {cell.val} {cell.row.size} {cell.data.size}',
        }}
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[
            {'Table 1': cmds}, {'Table 2': cmds}
        ])
        for shape_name in ('Table 1', 'Table 2'):
            src_table = self.get_shape(self.prs.slides[slides - 1].shapes, shape_name).table
            table = self.get_shape(prs.slides[0].shapes, shape_name).table
            # Table shape is extended or contracted
            eq_(len(table.rows), len(data))
            eq_(len(table.columns), len(data.columns))
            # Special rows / columns are set
            eq_(table.first_row, False)
            eq_(table.last_row, True)
            eq_(table.first_col, True)
            eq_(table.last_col, True)
            # Check column widths for changed columns
            gridcols = table._tbl.tblGrid.gridCol_lst
            eq_(int(gridcols[0].get('w')), pptx.util.Inches(1))
            eq_(int(gridcols[1].get('w')), pptx.util.Inches(2))
            eq_(int(gridcols[2].get('w')), pptx.util.Inches(1.5))

            # Check cell contents
            maxrow, maxcol = len(src_table.rows) - 1, len(src_table.columns) - 1
            for i, (index, row) in enumerate(data.iterrows()):
                # Row height is the same as in the source table (or its last row)
                eq_(table.rows[i].height, src_table.rows[min(i, maxrow)].height)
                for j, (column, val) in enumerate(row.iteritems()):
                    # Unspecified col width is the same as in the source table (or its last col)
                    if column in {'sales', 'growth'}:
                        eq_(table.columns[j].width, src_table.columns[min(j, maxcol)].width)
                    # Text matches, and all cell.* attributes are correct
                    cell = table.rows[i].cells[j]
                    paras = cell.text_frame.paragraphs
                    eq_(cell.text, f'{i} {j} {index} {column} {val} {row.size} {data.size}')
                    eq_(paras[0].font.bold, bool(i % 2))
                    # Check para font color, but not run font color. Run is overwritten because
                    # text: command is given AFTER color: command
                    eq_(paras[0].font.color.rgb, (255, 0, 0) if i % 3 else (0, 128, 0))
                    eq_(cell.fill.fore_color.rgb, (238, 238, 238) if i % 2 else (204, 204, 204))
                    self.check_opacity(cell.fill.fore_color, 0.4)
                    eq_(paras[0].font.size,
                        pptx.util.Pt(10) if column == 'देश' else pptx.util.Pt(8))
                    eq_(paras[0].font.name, 'Arial')
                    eq_(paras[0].font.italic,
                        True if column == 'देश' else
                        bool(i % 2) if column == 'city' else
                        None)
                    eq_(paras[0].runs[1].font.italic, True)
                    eq_(paras[0].font.underline, bool(j % 2))
                    eq_(paras[0].alignment, PP_ALIGN.LEFT if i % 2 else PP_ALIGN.RIGHT)
                    eq_(cell.vertical_anchor,
                        MVA.MIDDLE if column == 'देश' else
                        (MVA.TOP if i % 2 else MVA.BOTTOM) if column == 'city' else
                        None)
                    eq_(cell.margin_left,
                        pptx.util.Inches(0.05) if column == 'देश' else
                        pptx.util.Inches(0 if j % 2 else 0.1) if column == 'city' else
                        pptx.util.Inches(0.1))
                    eq_(cell.margin_right, pptx.util.Pt(1))
                    eq_(cell.margin_top, pptx.util.Inches(0 if j % 2 else 0.1))
                    eq_(cell.margin_bottom, 0)

        # table: can only apply to a table element, not text
        with assert_raises(ValueError):
            pptgen(source=self.input, only=slides, rules=[{'Title 1': {'table': {}}}])
        # table.data: must be a DataFrame
        with assert_raises(ValueError):
            pptgen(source=self.input, only=slides, rules=[{'Table 1': {'table': {'data': []}}}])
        # Invalid column names raise a warning
        with LogCapture() as logs:
            pptgen(source=self.input, only=slides, rules=[{'Table 1': {'table': {
                'data': self.data.head(3),
                'width': {'NA1': 1},
                'text': {'NA2': 0},
            }}}])
        logs.check_present(
            ('gramex', 'WARNING', 'pptgen2: No column: NA1 in table: Table 1'),
            ('gramex', 'WARNING', 'pptgen2: No column: NA2 in table: Table 1'),
        )

    # TODO: if we delete slide 6 and use slides=[6, 7], this causes an error
    def test_copy_slide(self, slides=[7, 8]):
        data = [1, 1.5, 2]
        prs = pptgen(source=self.input, target=self.output, only=slides, mode='expr', rules=[
            {
                'slide-numbers': [1, 2],
                'copy-slide': data,
                'data': {'mycopy': 'copy'},
                'Title 1': {'text': 'f"{copy.pos}: {copy.key} - {copy.val}: {len(copy.slides)}"'},
                'TL': {'top': 'copy.val', 'left': 'mycopy.val'},
                'TC': {'top': 'copy.val', 'left': 'mycopy.val * 2'},
                'TR': {'top': 'copy.val', 'left': 'mycopy.val * 3'},
                'CL': {'top': 'copy.val * 2', 'left': 'mycopy.val'},
                'CC': {'top': 'copy.val * 2', 'left': 'mycopy.val * 2'},
                'CR': {'top': 'copy.val * 2', 'left': 'mycopy.val * 3'},
                'BL': {'top': 'copy.val * 3', 'left': 'mycopy.val'},
                'BC': {'top': 'copy.val * 3', 'left': 'mycopy.val * 2'},
                'BR': {'top': 'copy.val * 3', 'left': 'mycopy.val * 3'},
            }
        ])
        # All shapes are copied into 3 slides?
        eq_(len(prs.slides), len(slides) * len(data))
        names = [[shape.name for shape in slide.shapes] for slide in prs.slides]
        eq_(names[0], names[2])
        eq_(names[0], names[4])
        eq_(names[1], names[3])
        eq_(names[1], names[5])
        # Titles are copied?
        eq_(prs.slides[0].shapes.title.text, '0: 0 - 1: 2')
        eq_(prs.slides[2].shapes.title.text, '1: 1 - 1.5: 2')
        eq_(prs.slides[4].shapes.title.text, '2: 2 - 2: 2')
        # Position commands are applied?
        for val, slide in zip(data, (1, 3, 5)):
            eq_(self.get_shape(prs.slides[slide].shapes, 'TL').left, pptx.util.Inches(val))
            eq_(self.get_shape(prs.slides[slide].shapes, 'TL').top, pptx.util.Inches(val))
            eq_(self.get_shape(prs.slides[slide].shapes, 'CC').left, pptx.util.Inches(2 * val))
            eq_(self.get_shape(prs.slides[slide].shapes, 'CC').top, pptx.util.Inches(2 * val))
            eq_(self.get_shape(prs.slides[slide].shapes, 'BR').left, pptx.util.Inches(3 * val))
            eq_(self.get_shape(prs.slides[slide].shapes, 'BR').top, pptx.util.Inches(3 * val))
        # Background is copied?
        for n in (1, 3, 5):
            eq_(prs.slides[n].background.fill.fore_color.theme_color, MSO_THEME_COLOR.ACCENT_3)
            self.assertAlmostEqual(prs.slides[n].background.fill.fore_color.brightness, 0.8)
        # Links are copied?
        for n in (0, 2, 4):
            shape = self.get_shape(prs.slides[n].shapes, 'Freeform 1')
            eq_(shape.click_action.hyperlink.address, 'https://t.co/')
            para = self.get_shape(prs.slides[n].shapes, 'TextBox 1').text_frame.paragraphs[0]
            eq_(para.runs[1]._r.find('.//' + qn('a:hlinkClick')).get('action'),
                'ppaction://hlinkshowjump?jump=firstslide')

    def chart_data(self, shape):
        return pd.read_excel(
            io.BytesIO(shape.chart.part.chart_workbook.xlsx_part.blob),
            index_col=0,
            engine='openpyxl',
        ).fillna('')

    def test_chart_data(self, slides=[10]):
        data = pd.DataFrame({
            'Alpha': [1, 2, 3],
            'Beta': [4, 5, 6],
            'Gamma': [7, 9, ''],
        }, index=['X', 'Y', 'Z'])
        charts_2d = ['Column Chart', 'Line Chart', 'Bar Chart']
        charts_1d = ['Pie Chart', 'Donut Chart']
        prs1 = pptgen(source=self.input, target=self.output, only=slides, rules=[
            *({chart: {'chart-data': data[['Alpha']]}} for chart in charts_1d),
            *({chart: {'chart-data': data}} for chart in charts_2d)
        ])
        prs2 = pptgen(source=self.input, target=self.output, only=slides, rules=[
            *({chart: {'chart': {'data': data[['Alpha']]}}} for chart in charts_1d),
            *({chart: {'chart': {'data': data}}} for chart in charts_2d)
        ])
        for prs in (prs1, prs2):
            shapes = prs.slides[0].shapes
            for chart in charts_1d:
                afe(self.chart_data(self.get_shape(shapes, chart)), data[['Alpha']])
            for chart in charts_2d:
                afe(self.chart_data(self.get_shape(shapes, chart)), data)

    def test_chart_attrs(self, slides=[10]):
        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Bar Chart': {
                'chart': {
                    'fill': pd.DataFrame({
                        'Series 1': 'red',
                        'Series 2': 'yellow',
                        'Series 3': ['#f88', '#f44', '#f00', '#800'],
                    }),
                    'text': pd.DataFrame({
                        'Series 1': ['<a bold="y">1.1</a>', '1.2', '1.3', '1.4'],
                        'Series 2': ['<a bold="y">2.1</a>', '2.2', '2.3', '2.4'],
                        'Series 3': ['<a bold="y">3.1</a>', '3.2', '3.3', '3.4'],
                    })
                }
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Bar Chart').chart.series
        eq_(series[0].points[0].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[2].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[3].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[1].points[0].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[1].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[2].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[3].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[2].points[0].format.fill.fore_color.rgb, (255, 136, 136))
        eq_(series[2].points[1].format.fill.fore_color.rgb, (255, 68, 68))
        eq_(series[2].points[2].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[2].points[3].format.fill.fore_color.rgb, (136, 0, 0))
        for series_index in range(3):
            for point_index in range(4):
                frame = series[series_index].points[point_index].data_label.text_frame
                eq_(frame.paragraphs[0].runs[0].text, f'{series_index + 1}.{point_index + 1}')
                eq_(frame.paragraphs[0].runs[0].font.bold, True if point_index == 0 else None)

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Column Chart': {
                'chart': {'fill': pd.DataFrame({
                    'Series 1': 'red',
                    'Series 2': 'yellow',
                    'Series 3': ['#f88', '#f44', '#f00', '#800'],
                })}
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Column Chart').chart.series
        eq_(series[0].points[0].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[2].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[3].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[1].points[0].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[1].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[2].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[3].format.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[2].points[0].format.fill.fore_color.rgb, (255, 136, 136))
        eq_(series[2].points[1].format.fill.fore_color.rgb, (255, 68, 68))
        eq_(series[2].points[2].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[2].points[3].format.fill.fore_color.rgb, (136, 0, 0))

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Radar Chart': {
                'chart': {'stroke': pd.DataFrame({
                    'Series 1': ['red', 'blue', 'green', '#777'],
                    'Series 2': ['yellow', 'black', 'red', '#088'],
                    'Series 3': ['#ccc', '#aaa', '#777', '#444'],
                })}
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Radar Chart').chart.series
        eq_(series[0].points[0].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[2].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[0].points[3].format.line.fill.fore_color.rgb, (119, 119, 119))
        eq_(series[1].points[0].format.line.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[1].format.line.fill.fore_color.rgb, (0, 0, 0))
        eq_(series[1].points[2].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[1].points[3].format.line.fill.fore_color.rgb, (0, 136, 136))
        eq_(series[2].points[0].format.line.fill.fore_color.rgb, (204, 204, 204))
        eq_(series[2].points[1].format.line.fill.fore_color.rgb, (170, 170, 170))
        eq_(series[2].points[2].format.line.fill.fore_color.rgb, (119, 119, 119))
        eq_(series[2].points[3].format.line.fill.fore_color.rgb, (68, 68, 68))

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'XY Chart': {
                'chart': {
                    'fill': pd.DataFrame({
                        'Y-Values': ['red', 'blue', 'green', '#777'],
                        'Z-Values': ['yellow', 'black', 'red', '#088'],
                    }),
                    'stroke': pd.DataFrame({
                        'Y-Values': ['red', 'blue', 'green', '#777'],
                        'Z-Values': ['yellow', 'black', 'red', '#088'],
                    })
                }
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'XY Chart').chart.series
        # TODO: Test fill
        eq_(series[0].points[0].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[2].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[1].points[0].format.line.fill.fore_color.rgb, (255, 255, 0))
        eq_(series[1].points[1].format.line.fill.fore_color.rgb, (0, 0, 0))
        eq_(series[1].points[2].format.line.fill.fore_color.rgb, (255, 0, 0))

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Pie Chart': {
                'chart': {
                    'fill': pd.DataFrame({
                        'Sales': ['red', 'blue', 'green', '#777'],
                    }),
                    'stroke': pd.DataFrame({
                        'Sales': ['red', 'blue', 'green', '#777'],
                    })
                }
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Pie Chart').chart.series
        eq_(series[0].points[0].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[2].format.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[0].points[3].format.fill.fore_color.rgb, (119, 119, 119))
        eq_(series[0].points[0].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[2].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[0].points[3].format.line.fill.fore_color.rgb, (119, 119, 119))

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Donut Chart': {
                'chart': {
                    'fill': pd.DataFrame({
                        'Sales': ['red', 'blue', 'green', '#777'],
                    }),
                    'stroke': pd.DataFrame({
                        'Sales': ['red', 'blue', 'green', '#777'],
                    })
                }
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Donut Chart').chart.series
        eq_(series[0].points[0].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[2].format.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[0].points[3].format.fill.fore_color.rgb, (119, 119, 119))
        eq_(series[0].points[0].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[2].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[0].points[3].format.line.fill.fore_color.rgb, (119, 119, 119))

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Line Chart': {
                'chart': {
                    'stroke': pd.DataFrame({
                        'Series 1': ['red'] * 4,
                        'Series 2': ['blue'] * 4,
                        'Series 3': ['green'] * 4,
                    })
                }
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Line Chart').chart.series
        eq_(series[0].points[0].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[2].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[3].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[1].points[0].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[1].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[2].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[3].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[2].points[0].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[2].points[1].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[2].points[2].format.line.fill.fore_color.rgb, (0, 128, 0))
        eq_(series[2].points[3].format.line.fill.fore_color.rgb, (0, 128, 0))

        prs = pptgen(source=self.input, target=self.output, only=slides, rules=[{
            'Area Chart': {
                'chart': {
                    'fill': pd.DataFrame({
                        'Series 1': ['red'] * 4,
                        'Series 2': ['blue'] * 4,
                    }),
                    'stroke': pd.DataFrame({
                        'Series 1': ['red'] * 4,
                        'Series 2': ['blue'] * 4,
                    })
                }
            },
        }])
        series = self.get_shape(prs.slides[0].shapes, 'Area Chart').chart.series
        eq_(series[0].points[0].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[2].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[3].format.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[1].points[0].format.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[1].format.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[2].format.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[3].format.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[0].points[0].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[1].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[2].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[0].points[3].format.line.fill.fore_color.rgb, (255, 0, 0))
        eq_(series[1].points[0].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[1].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[2].format.line.fill.fore_color.rgb, (0, 0, 255))
        eq_(series[1].points[3].format.line.fill.fore_color.rgb, (0, 0, 255))

    def test_commandline(self):
        # "slidesense" prints usage
        with OutputCapture() as logs:
            commandline([])
        ok_(logs.captured.startswith('usage: slidesense'))
        # "slidesense nonexistent.yaml" prints an error
        with LogCapture() as logs:
            commandline(['nonexistent.yaml'])
        logs.check_present(
            ('gramex', 'ERROR', 'No rules found in file: nonexistent.yaml')
        )
        # "slidesense gramex.yaml nonexistent-url" prints an error
        with LogCapture() as logs:
            path = os.path.join(folder, 'slidesense-gramex.yaml')
            commandline([path, 'nonexistent-url'])
        logs.check_present(
            ('gramex', 'ERROR', 'No PPTXHandler matched in file: ' + path)
        )

        target = os.path.join(folder, 'output.pptx')
        non_target = os.path.join(folder, 'nonexistent.pptx')

        for args in (
            ('slidesense-config.yaml', ),
            ('slidesense-gramex.yaml', ),
            ('slidesense-gramex.yaml', 'slidesense-test'),
        ):
            self.remove_output()
            commandline([os.path.join(folder, args[0]), *args[1:],
                         f'--target={target}', '--no-open'])
            ok_(os.path.exists(target))
            ok_(not os.path.exists(non_target))

    @classmethod
    def tearDown(cls):
        cls.remove_output()
