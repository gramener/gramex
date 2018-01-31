# -*- coding: utf-8 -*-
from __future__ import unicode_literals

import io
import pptx
from nose.tools import eq_
from . import TestGramex


class TestPPTXHandler(TestGramex):
    def get_pptx(self, url):
        content = self.get(url).content
        prs = pptx.Presentation(io.BytesIO(content))
        text = {
            shp.name: shp.text
            for shp in prs.slides[0].shapes
            if shp.has_text_frame
        }
        return prs, text

    def test_source(self):
        prs, text = self.get_pptx('/pptx/source')
        eq_(text, {
            'Text1': 'Text1',
            'Arabic': 'Arabic',
            'Chinese': 'Chinese',
            'Hindi': 'Hindi',
            'Russian': 'Russian',
        })

    def test_change_text(self):
        prs, text = self.get_pptx('/pptx/change-text')
        eq_(text, {
            'Text1': 'New Text1',
            'Arabic': 'سعيد',
            'Chinese': '高兴',
            'Hindi': 'खुश',
            'Russian': 'счастливый',
        })

    def test_change_text_args(self):
        prs, text = self.get_pptx('/pptx/change-text-args?text=Text')
        eq_(text, {
            'Text1': 'Text',
            'Arabic': 'سعيد',
            'Chinese': '高兴',
            'Hindi': 'खुश',
            'Russian': 'счастливый',
        })
