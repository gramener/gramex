import io
import os
import re
import json
import time
import logging
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from nose.tools import eq_, ok_
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from pdfminer.converter import PDFPageAggregator
from pdfminer.pdfinterp import PDFResourceManager
from pdfminer.pdfinterp import PDFPageInterpreter
from pdfminer.high_level import extract_text_to_fp
from shutilwhich import which
from tornado.web import create_signed_value
import gramex
from gramex.handlers import Capture
from . import TestGramex, server

_captures = {}
paths = {'phantomjs': which('phantomjs'), 'node': which('node')}


def get_capture(name, **kwargs):
    '''Return a cached Capture() object constructed with kwargs'''
    if name in _captures:
        return _captures[name]
    capture = _captures[name] = Capture(**kwargs)
    end_time, delay = time.time() + 10, 0.05
    logging.info('Waiting for capture.js...')
    while not capture.started:
        if time.time() > end_time:
            raise RuntimeError('capture.js took too long to start')
        time.sleep(delay)
    return capture


def get_text(content):
    infile, outfile = io.BytesIO(content), io.BytesIO()
    extract_text_to_fp(infile, outfile)
    return outfile.getvalue().decode('utf-8')


def normalize(s):
    '''
    Ignore Unicode characters. These are normalized by PhantomJS.
    Ignore whitespaces. They may translate to \t sometimes.
    Just stick to ASCII.
    '''
    return re.sub(r'[^\x33-\x7f]+', '', s)


def get_layout_elements(content):
    '''Take content of pdf and return list of text in order that it occurs'''
    rsrcmgr = PDFResourceManager()
    device = PDFPageAggregator(rsrcmgr, laparams=LAParams())
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    # Create a PDF device object.
    parser = PDFParser(io.BytesIO(content))
    # Create a PDF page aggregator object.
    page = next(PDFPage.create_pages(PDFDocument(parser)))
    interpreter.process_page(page)
    # receive the LTPage object for the page.
    return [child.get_text() for child in device.get_result() if hasattr(child, 'get_text')]


def test_dependencies():
    assert paths['phantomjs'], 'phantomjs is not installed'
    assert paths['node'], 'node is not installed'


class TestCaptureHandler(TestGramex):
    # Note: This only tests PhantomJS. See TestCaptureHandlerChrome for Chrome
    size = (1200, 768)
    max_colors = 65536
    url = '/dir/capture'
    src = '/capture'

    @classmethod
    def setupClass(cls):
        cls.capture = get_capture('default', port=9402)
        cls.folder = os.path.dirname(os.path.abspath(__file__))

    def test_init(self):
        ok_(self.capture.started)

    def fetch(self, url, code=200, **kwargs):
        r = self.get(url, **kwargs)
        eq_(r.status_code, code)
        return r

    def check_filename(self, result, name):
        disp = result.headers.get('Content-Disposition', '')
        eq_(disp, 'attachment; filename="%s"' % name)

    def check_pdf(self, content):
        # Ensure that each line in dir/capture.html is in the text
        # On Linux, spaces seem to appear as tabs. Normalize that
        text = normalize(get_text(content))
        source = os.path.join(self.folder, 'dir', 'capture.html')
        for line in io.open(source, 'r', encoding='utf-8'):
            frag = normalize(line.strip())
            # Only compare non-HTML non-empty lines
            if len(frag) > 0 and '<' not in line:
                self.assertIn(frag, text)

    def test_pdf(self):
        # Test capture library API
        content = self.capture.pdf(url=server.base_url + self.url)
        self.check_pdf(content)

        # Test service: relative and absolute URLs
        for url in (server.base_url + self.url, '..' + self.url, self.url):
            result = self.fetch(self.src, params={'url': url})
            self.check_filename(result, 'screenshot.pdf')
            self.check_pdf(result.content)

        # delay=. After 500ms, page changes text and color to blue
        # file=.  Changes filename
        result = self.fetch(self.src, params={'url': self.url, 'delay': 600, 'file': 'delay'})
        self.check_filename(result, 'delay.pdf')
        self.assertIn('Blueblock', normalize(get_text(result.content)))

        # --format and --orientation
        result = self.fetch(self.src, params={
            'url': self.url, 'format': 'A3', 'orientation': 'landscape'})
        parser = PDFParser(io.BytesIO(result.content))
        page = next(PDFPage.create_pages(PDFDocument(parser)))
        self.assertIn([round(x) for x in page.attrs['MediaBox']], (
            [0, 0, 1188, 842],      # noqa: Chrome uses 1188 x 842 for A3
            [0, 0, 1191, 842],      # noqa: PhantomJS uses 1191 x 842 for A3
        ))

        # cookie=. The Cookie is printed on the screen via JS
        result = self.fetch(self.src, params={'url': self.url + '?show-cookie', 'cookie': 'a=x'})
        self.assertIn('a=x', normalize(get_text(result.content)))
        # Cookie: header is the same as ?cookie=.
        # Old request cookies vanish. Only new ones remain
        result = self.fetch(self.src, params={'url': self.url + '?show-cookie'},
                            headers={'Cookie': 'b=z'})
        result_text = normalize(get_text(result.content))
        self.assertIn('js:cookie=b=z', result_text)
        self.assertIn('server:cookie=b=z', result_text)

    def check_img(self, content, color=None, min=100, size=None):
        img = Image.open(io.BytesIO(content)).convert(mode='RGBA')
        if size:
            eq_(img.size, size)
        # dir/index.html has a style="color:red" block in it. So check that there's enough red
        if color:
            # Ensure that the color is present at least in min pixels in the image
            colors = {clr: freq for freq, clr in img.getcolors(self.max_colors)}
            self.assertGreater(colors.get(color, 0), min)
        return img

    def test_png(self):
        content = self.capture.png(url=server.base_url + self.url)
        self.check_img(content, color=(255, 0, 0, 255), min=100, size=self.size)

        # Check file=, ext=, width=, height=
        result = self.fetch(self.src, params={
            'url': self.url, 'width': 600, 'height': 1200, 'file': 'capture', 'ext': 'png'})
        self.check_filename(result, 'capture.png')
        self.check_img(result.content, size=(600, 1200))

        # selector=. has a 100x100 green patch
        result = self.fetch(self.src, params={
            'url': self.url, 'selector': '.subset', 'ext': 'png'})
        self.check_img(result.content, color=(0, 128, 0, 255), min=9000, size=(100, 100))

    def test_jpg(self):
        # TODO: check colors
        content = self.capture.jpg(url=server.base_url + self.url)
        self.check_img(content)

        # Check file=, ext=, width=, height=
        result = self.fetch(self.src, params={
            'url': self.url, 'width': 800, 'height': 1000, 'file': 'capture', 'ext': 'jpg'})
        self.check_filename(result, 'capture.jpg')
        self.check_img(result.content, size=(800, 1000))


class TestCaptureHandlerChrome(TestCaptureHandler):
    src = '/capturechrome'

    @classmethod
    def setupClass(cls):
        cls.capture = get_capture('chrome', port=9412, engine='chrome', timeout=10)
        cls.folder = os.path.dirname(os.path.abspath(__file__))

    @staticmethod
    def check_text(shape, text=None, font_size=None):
        if text is not None:
            eq_(shape.text, text)
        if font_size is not None:
            eq_(shape.text_frame.paragraphs[0].runs[0].font.size, Pt(font_size))

    def test_pptx(self):
        content = self.capture.pptx(url=server.base_url + self.url)
        prs = Presentation(io.BytesIO(content))
        eq_(len(prs.slides), 1)
        self.check_img(prs.slides[0].shapes[0].image.blob)

        # Non-existent selectors raise a HTTP error
        self.fetch(self.src, code=500, params={
            'url': self.url, 'selector': 'nonexistent', 'ext': 'pptx'})

        # Check selector=.subset has a 100x100 green patch
        title = '高=σ'
        result = self.fetch(self.src, params={
            'url': self.url, 'title': title, 'selector': '.subset', 'ext': 'pptx'})
        prs = Presentation(io.BytesIO(result.content))
        eq_(len(prs.slides), 1)
        self.check_img(prs.slides[0].shapes[0].image.blob,
                       color=(0, 128, 0, 255), min=9000, size=(100, 100))
        self.check_text(prs.slides[0].shapes[1], text=title, font_size=18)

        # Check title_size
        result = self.fetch(self.src, params={
            'url': self.url, 'title': title, 'selector': '.subset', 'ext': 'pptx',
            'title_size': 24})
        prs = Presentation(io.BytesIO(result.content))
        self.check_text(prs.slides[0].shapes[1], text=title, font_size=24)

        # Check multi-slide generation with multi-title and multi-selector
        result = self.fetch(self.src, params={
            'url': self.url, 'ext': 'pptx', 'title': ['高', 'σ'], 'selector': ['.subset', 'p']})
        prs = Presentation(io.BytesIO(result.content))
        eq_(len(prs.slides), 2)
        self.check_img(prs.slides[0].shapes[0].image.blob,
                       color=(0, 128, 0, 255), min=9000, size=(100, 100))
        para_img = self.check_img(prs.slides[1].shapes[0].image.blob)
        para_size = (prs.slides[1].shapes[0].width, prs.slides[1].shapes[0].height)
        eq_(prs.slides[0].shapes[1].text, '高')
        eq_(prs.slides[1].shapes[1].text, 'σ')
        # Check multi-dpi generation and default dpi as 96
        result = self.fetch(self.src, params={
            'url': self.url, 'ext': 'pptx',
            'selector': ['.subset', 'p'],
            'dpi': ['96', '192'],               # 2nd image has half the original size
        })
        prs = Presentation(io.BytesIO(result.content))
        eq_(len(prs.slides), 2)
        self.check_img(prs.slides[0].shapes[0].image.blob,
                       color=(0, 128, 0, 255), min=9000, size=(100, 100))
        self.check_img(prs.slides[1].shapes[0].image.blob, size=para_img.size)
        # 2nd image with twice the dpi is half the size
        self.assertAlmostEqual(prs.slides[1].shapes[0].width * 2, para_size[0], delta=1)
        self.assertAlmostEqual(prs.slides[1].shapes[0].height * 2, para_size[1], delta=1)

    def test_delay_render(self):
        # delay=. After 1 second, renderComplete is set. Page changes text and color to green
        result = self.fetch(self.src, params={'url': self.url, 'delay': 'renderComplete',
                                              'file': 'delay-str'})
        self.check_filename(result, 'delay-str.pdf')
        self.assertIn('Greenblock', normalize(get_text(result.content)))

    def test_headers(self):
        # Headers are passed through, but cookie overrides it
        result = self.fetch(self.src, params={'url': '/httpbin', 'cookie': 'right=1'}, headers={
            'dnt': '1',
            'user-agent': 'gramex-test',
            'cookie': 'wrong=1'
        })
        text = get_text(result.content)
        ok_(re.search(r'"dnt":\s+"1"', text, re.IGNORECASE))
        ok_(re.search(r'"user\-agent":\s+"gramex-test"', text, re.IGNORECASE))
        # ?cookie= should override header. Result must be right=1, not wrong=1
        ok_('wrong' not in text)
        ok_('right' in text)

        # Check that the request can mimic a user
        user = {'id': 'login@example.org', 'role': 'manager'}
        secret = gramex.service.app.settings['cookie_secret']
        result = self.fetch(self.src, params={'url': '/auth/session'}, headers={
            'x-gramex-user': create_signed_value(secret, 'user', json.dumps(user)),
        })
        text = get_text(result.content)
        ok_(user['id'] in text)
        ok_(user['role'] in text)

    def err(self, code, **params):
        return self.fetch(self.src, code=code, params=params)

    def test_errors(self):
        # nonexistent URLs should capture the 404 page and return a screenshot
        self.err(code=200, url='/nonexistent')
        self.err(code=500, url='http://nonexistent')
        self.err(code=400, url=self.url, ext='nonexistent')
        self.err(code=500, url=self.url, selector='nonexistent', ext='png')

    def test_pdf_header_footer(self):
        result = self.fetch(self.src, params={
            'url': self.url,
            'header': '\u00A9|Header|$pageNumber',
            'footer': '|$title|'
        })
        # Check if elements are ordered correctly.
        layout = get_layout_elements(result.content)
        ok_('©' in layout[0])
        ok_('Header' in layout[1])
        ok_('1' in layout[2])
        ok_('This is the footer' in layout[-1])
        # Check templates
        result = self.fetch(self.src, params={
            'url': self.url,
            'headerTemplate': '<div><span class="url"></span><h1>Header</h1></div>',
            'footerTemplate': '''<div style="font-size:5px">
                <span class="pageNumber"></span>/<span class="totalPages"></span></div>'''
        })
        layout = get_layout_elements(result.content)
        ok_(self.url in layout[0])
        # Text elements are grouped differently in PDFs generated by Window and Linux
        ok_('Header' in layout[0] or 'Header' in layout[1])
        ok_(re.search(r'1\s*/\s*1', layout[-1]))

    def test_browser_disconnect(self, delay=0.2):
        # Disconnect the browser after the response is received
        content = self.capture.jpg(url=server.base_url + self.url, _test_disconnect=1)
        self.check_img(content)
        # Re-connecting to the instance takes a bit of time. This may vary based on system speed
        # If this test fails due to ReadTimeout exception, increase the delay.
        time.sleep(delay)
        # Capture the same image again. The server should re-connect automatically
        content = self.capture.jpg(url=server.base_url + self.url)
        self.check_img(content)
