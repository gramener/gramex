'''
pptgen2.pptgen() modifies a ``source`` PPTX using ``rules`` and ``data``.

The main loop is in :py:func:`pptgen`, which calls other functions as required.
'''

import copy
import gramex
import gramex.data
import os
import pandas as pd
import pptx
import sys
from fnmatch import fnmatchcase
from gramex.config import app_log
from gramex.transforms import build_transform
from orderedattrdict import AttrDict
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import _nsmap, qn
from pptx.shapes.shapetree import SlideShapes
from textwrap import dedent
from typing import Union, List, Dict
from . import commands


def pptgen(
    source: Union[str, pptx.presentation.Presentation],
    rules: List[dict] = [],
    data: dict = {},
    target: str = None,
    only: Union[int, List[int]] = None,
    register: Dict[str, str] = {},
    unit: str = 'Inches',
    mode: str = 'literal',
    handler=None,
    **config,
) -> pptx.presentation.Presentation:
    '''
    Process a configuration. This loads a Presentation from source, applies the
    (optional) configuration changes and (optionally) saves it into target. Returns the modified

    :arg PPTX source: string or pptx.Presentation object to transform
    :arg list rules: list of rules to apply to the ``source`` PPTX. Each rule
    :arg str target: optional path save file
    :arg int/list only: slide number(s) to process. 1 is the first slide. [1, 3] is slides 1 & 3
    :arg dict register: new commands to register via :py:func:`register_commands`.
    :arg str unit: default length unit (Inches, Cm, Centipoints, etc)
    :arg str mode: default expression mode. Values in Python are treated as 'literals'
        (e.g. 'red' is the STRING red). But PPTXHandler passes the mode as `expr`. Values are
        treated as expressions (e.g. 'red' is the VARIABLE red).
    :arg handler: if PPTXHandler passes a handler, make it available to the commands as a variable
    :return: target PPTX
    '''
    # TODO: source can be an expression. PPTXHandler may need to use multiple themes
    prs = source if isinstance(source, pptx.presentation.Presentation) else Presentation(source)
    # Load data with additional variables:
    #   prs: source presentation
    #   handler: if PPTXHandler passes a handler, allow commands to use it as a variable
    #   _expr_mode: Allows commands.expr() to evaluate specs as literals or expressions correctly
    data = load_data(data, handler=handler, prs=prs, _expr_mode='expr' in mode.lower())

    register_commands(register)
    commands.length_unit = commands.length_class(unit)

    slides = pick_only_slides(prs, only)
    # PPTX applies transforms to groups. Flatten these so that changing position works as expected
    for slide in slides:
        for shape in slide.shapes:
            commands.flatten_group_transforms(shape)
    # copy-slide: can copy any set of slides multiple times. To track of which source slide maps to
    # which target slide, we use `slide_map`. slide_map[target_slide_index] = source_slide_index
    slide_map = list(range(len(slides)))

    # Loop through each rule (copying them to protect from modification)
    for rule in copy.deepcopy(rules):
        slides_in_rule = tuple(slide_filter(slides, rule, data))
        # If no slides matched, warn the user
        if len(slides_in_rule) == 0:
            app_log.warn(
                f'pptgen2: No slide with slide-number: {rule.get("slide-number")}, '
                f'slide-title: {rule.get("slide-title")}'
            )
            continue
        # Copy slides after the last mapped position of the last slide in this rule
        max_index = max(index for index, slide in slides_in_rule)
        copy_pos = next(i for i, v in reversed(tuple(enumerate(slide_map))) if max_index == v)
        # Copy all slides into the `copies` list BEFORE applying any rules. Ensures that rules
        # applied to slide 1 don't propagate into 2, 3, etc.
        copies = []
        copy_seq = iterate_on(rule.get('copy-slide', [None]), data)
        for i, (copy_key, copy_val) in enumerate(copy_seq):
            copy_row = AttrDict(pos=i, key=copy_key, val=copy_val, slides=[])
            copies.append(copy_row)
            for index, slide in slides_in_rule:
                if i > 0:
                    copy_pos += 1
                    slide = copy_slide(prs, slide, copy_pos)
                    slide_map.insert(copy_pos, index)
                copy_row.slides.append(slide)
        # Apply rules on all copied slides
        for copy_row in copies:
            # Include rule-level `data:`. Add copy, slide as variables
            slide_data = load_data(
                rule.get('data', {}), _default_key='function', copy=copy_row, **data
            )
            for slide in copy_row.slides:
                slide_data['slide'] = slide  # Rule can use 'slide' as a variable
                transition(slide, rule.get('transition', None), data)
                apply_commands(rule, slide.shapes, slide_data)
    if target:
        prs.save(target)
    return prs


# List of commands that can be used in a rule. Shapes cannot have these names
rule_cmdlist = {
    'copy-slide',
    'data',
    'slide-number',
    'slide-title',
    'transition',
}
# List of special commands that can be used that are not in commands.cmdlist
special_cmdlist = {
    'clone-shape',
    'data',
}


def apply_commands(rule: Dict[str, dict], shapes, data: dict):
    '''
    Apply commands in rule to change shapes using data.

    :arg dict rule: a dict of shape names, and commands to apply on each.
        e.g. ``{"Oval 1": {"fill": "red"}, "Oval 2": {"text": "OK"}}``
    :arg Shapes shapes: a slide.shapes or group.shapes object on which the rule should be applied
    :arg dict data: data context for the commands in the rule
    '''
    # Apply every rule to every pattern -- as long as the rule key matches the shape name
    for pattern, spec in rule.items():
        if pattern in rule_cmdlist:
            continue
        matched_shapes = [shape for shape in shapes if fnmatchcase(shape.name, pattern)]
        for shape in matched_shapes:
            # Clone all slides into the `clones` list BEFORE applying any command. Ensures that
            # commands applied to the shape don't propagate into its clones
            clones = []
            clone_seq = iterate_on(spec.get('clone-shape', [None]), data)
            parent_clone = data.get('clone', None)
            for i, (clone_key, clone_val) in enumerate(clone_seq):
                if i > 0:
                    # This copies only a shape, group or image. Not table, chart, media, equation,
                    # or zoom. But we don't see a need for these yet.
                    el = copy.deepcopy(shape.element)
                    shape.element.addnext(el)
                    shape = pptx.shapes.autoshape.Shape(el, shape._parent)
                clones.append(
                    AttrDict(pos=i, key=clone_key, val=clone_val, shape=shape, parent=parent_clone)
                )
            # Run commands in the spec on all cloned shapes
            is_group = shape.element.tag.endswith('}grpSp')
            for clone in clones:
                # Include shape-level `data:`. Add shape, clone as variables
                shape_data = load_data(
                    spec.get('data', {}),
                    _default_key='function',
                    shape=shape,
                    clone=clone,
                    **{k: v for k, v in data.items() if k not in {'shape', 'clone'}},
                )
                for cmd in spec:
                    if cmd in commands.cmdlist:
                        commands.cmdlist[cmd](clone.shape, spec[cmd], shape_data)
                    # Warn on unknown commands. But don't warn on groups -- they have sub-shapes
                    elif cmd not in special_cmdlist and not is_group:
                        app_log.warn(f'pptgen2: Unknown command: {cmd} on shape: {pattern}')
                # If the shape is a group, apply spec to each sub-shape
                if is_group:
                    apply_commands(spec, SlideShapes(clone.shape.element, shapes), shape_data)
        # Warn if the pattern is neither a shape nor a command
        if (
            not matched_shapes
            and pattern not in special_cmdlist
            and pattern not in commands.cmdlist
        ):
            app_log.warn(f'pptgen2: No shape matches pattern: {pattern}')


def load_data(_conf, _default_key: str = None, **kwargs) -> dict:
    '''
    Loads datasets based on configuration and returns a dict of those datasets.

    :arg dataset _conf: The dataset configuration
    :arg str _default_key: Can be ``function``, ``url`` or ``None`` (default).
        If specified, it converts string data configurations into ``{_default_key: _conf}``.
    :return: A dict of datasets loaded based on the configuration.

    ``_conf`` is processed as follows:

    - String ``'data.xlsx'`` is loaded via :py:func:`gramex.cache.open` into ``{data: ...}`` if
        ``_default_key == 'url'``
    - String ``'data[0]'`` is evaluated via :py:func:`gramex.transforms.build_transform` into
        ``{data: ...}``` if ``_default_key == 'function'``
    - String ``anything``` raises an Exception if ``_default_key`` is None
    - Dict ``{url: ...}`` is loaded with :py:func:`gramex.data.filter` into ``{data: ...}``
    - Dict ``{function: ...}`` is evaluated via :py:func:`gramex.transforms.build_transform`
        into ``{data: ...}``
    - Dict ``{x: ..., y: ...}`` loads the respective datasets into ``x`` and ``y`` instead of
        ``data``. Each dataset is processed using the above rules.
    - Any other datatype passed is returned as is in ``{data: ...}``

    Any keyword arguments passed are also added to the resulting dataset, but overwritten only if
    ``_conf`` loaded a dataset that's not ``None``.
    '''

    def str2conf(data, key):
        '''Convert string configurations to {url: str} or {function:str} based on _default_key'''
        # If data is not a string, return data as-is
        if not isinstance(data, str):
            return data
        # If data is a string, return {_default_key: data} (or raise a TypeError)
        if _default_key is not None:
            return {_default_key: data}
        raise TypeError(f'{key}: must be a dict, not {data!r}')

    data = str2conf(_conf, 'data')
    if not isinstance(data, dict) or 'url' in data or 'function' in data:
        data = {'data': data}
    data = {key: str2conf(conf, key) for key, conf in data.items()}
    for key, conf in data.items():
        if isinstance(conf, dict):
            conf = copy.copy(conf)
            if 'url' in conf:
                if 'transform' in conf:
                    conf['transform'] = build_transform(
                        {'function': conf['transform']},
                        vars={'data': None, 'handler': None},
                        filename=f'PPTXHandler:data.{key}',
                        iter=False,
                    )
                data[key] = gramex.data.filter(**conf)
            if 'function' in conf:
                # Let functions use previously defined data variables, including current one
                _kwargs = {**kwargs, **data}
                _vars = {key: None for key in _kwargs}
                data[key] = build_transform(conf, vars=_vars, iter=False)(**_kwargs)
    # If the dataset returns a None, don't overwrite the default kwargs.
    # This allow defaults to pass through if a dataset is specified as None.
    for key, val in data.items():
        if (key not in kwargs) or (val is not None):
            kwargs[key] = val
    return kwargs


def register_commands(register: Dict[str, str]) -> None:
    '''
    Register a new command to the command list.

    :arg dict register: keys are the command name. Values are the Python expression to run to apply
        the command. The expression can use 3 variables: ``shape`` (the Shape object to modify),
        ``spec`` (the configuration passed to your command) and ``data``.
    '''
    if not isinstance(register, dict):
        raise TypeError(f'register: must be a dict, not {type(register)}')
    for key, conf in register.items():
        commands.cmdlist[key] = build_transform(
            {'function': conf}, vars={'shape': None, 'spec': None, 'data': None}, iter=False
        )


def pick_only_slides(prs: Presentation, only: Union[int, List[int]] = None) -> list:
    '''
    Delete slides except those specified in ``only``.

    :arg Presentation prs: presentation to delete slides from
    :arg int/list only: slide number(s) to process. 1 is the first slide. [1, 3] is slides 1 & 3
    :return: list of slides to process. These slides are also DELETED in the original prs
    '''
    if only is None:
        return list(prs.slides)
    if isinstance(only, int):
        only = [only]
    if not isinstance(only, list):
        raise TypeError(f'pptgen(only=) takes slide number or list, not {type(only)}')
    all_slides = set(range(1, 1 + len(prs.slides)))
    for slide_num in sorted(all_slides - set(only), reverse=True):
        rid = prs.slides._sldIdLst[slide_num - 1].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[slide_num - 1]
    return list(prs.slides)


def slide_filter(slides, rule: dict, data: dict):
    '''
    Filter slides. Return iterable of (index, slide) for only those slides matching numbers/titles.

    :arg Slides slides: a Slides object (e.g. ``prs.slides``) to select slides from
    :arg dict rule: a dict that may have a ``slide-number`` or ``slide-title`` key to filter by.
        These can be expressions. ``slide-number`` is the slide number(s) to filter. 1 is the first
        slide. ``[1, 3]`` is slides 1 & 3. ``slide-title`` has the slide title pattern(s) to filter
        e.g. ``*Match*`` matches all slides with "match" anywhere in the title (case-insensitive).
    :arg dict data: data context for the rule
    :return: an iterable that yields (index, slide)
    '''
    numbers = commands.expr(rule.get('slide-number', []), data)
    titles = commands.expr(rule.get('slide-title', []), data)
    # numbers is the SET of 1-indexed slide numbers this rule applies to
    numbers = set(numbers if isinstance(numbers, list) else [numbers])
    # titles is a LIST of title patterns - ANY of which the slide title must match
    titles = titles if isinstance(titles, list) else [titles]
    # yield only the slides that match these conditions
    for index, slide in enumerate(slides):
        if numbers and (index + 1 not in numbers):
            continue
        title = slide.shapes.title.text if slide.shapes.title else ''
        if titles and (not any(fnmatchcase(title.lower(), pattern.lower()) for pattern in titles)):
            continue
        yield index, slide


def copy_slide(prs, source, target_index):
    '''
    Copy ``source`` slide from presentatation ``prs`` to appear at ``target_index``. python-pptx
    does not have this feature, so we tweak XML directly. Does not copy slides with diagrams or
    charts yet.

    :arg Presentation prs: presentation to copy the slide in
    :arg Slide source: slide to copy
    :arg target_index: location to copy into. 0 makes it the first slide
    '''
    # Append slide with source's layout. Then delete shapes to get a blank slide
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in dest.shapes:
        shp.element.getparent().remove(shp.element)
    # Copy background
    if source.background.element.bg is not None:
        dest.background.element.insert(0, copy.deepcopy(source.background.element.bg))
    # Copy shapes from source, in order
    for shape in source.shapes:
        new_shape = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    # Copy rels from source
    for key, val in source.part.rels.items():
        target = val._target
        if val.reltype == pptx.opc.constants.RELATIONSHIP_TYPE.NOTES_SLIDE:
            # This removes notes slide formatting and just copies the text
            dest.notes_slide.notes_text_frame.text = source.notes_slide.notes_text_frame.text
            continue
        elif 'chart' in val.reltype:
            # https://github.com/scanny/python-pptx/issues/132#issuecomment-414001942
            # Does not work for Treemap, modern charts.
            # TODO: Replace with underlying lxml code
            partname = target.package.next_partname(pptx.parts.chart.ChartPart.partname_template)
            xlsx_blob = target.chart_workbook.xlsx_part.blob
            target = pptx.parts.chart.ChartPart(
                partname,
                target.content_type,
                copy.deepcopy(target._element),
                package=target.package,
            )
            target.chart_workbook.xlsx_part = pptx.parts.chart.EmbeddedXlsxPart.new(
                xlsx_blob, target.package
            )
        # TODO: handle diagrams
        dest.part.rels.add_relationship(val.reltype, target, val.rId, val.is_external)
    # Move appended slide into target_index
    prs.slides.element.insert(target_index, prs.slides.element[-1])
    return dest


def transition(slide, spec: Union[str, dict], data: dict):
    '''
    Apply transition on slide based on spec. python-pptx does not have this feature, so we tweak
    XML directly.

    :arg Slide slide: slide to apply the transition to
    :arg dict/str spec: type of transition to apply. ``config.yaml`` lists transitions and options.
        It can also be a dict specifying ``type`` (type of transition), ``duration`` (in seconds)
        and ``advance`` (auto-advance after seconds)
    :arg dict data: data context for the ``spec`` expression
    '''
    if spec is None:
        return
    # Convert spec into this format: {type: ..., advance: ..., duration: ...}
    if isinstance(spec, str):
        spec = {'type': spec}
    if not isinstance(spec, dict):
        raise ValueError('transition: %r is not a str or dict' % spec)
    # conf is from config.yaml, and has all transition types and options
    conf = commands.conf
    type = commands.expr(spec.get('type', None), data)
    if type is not None:
        # Parse the type into OXML: "glitter diamond left" -> <glitter pattern="diamond" dir="r"/>
        tag, *options = type.split()
        attrs = {}
        if tag in conf['transition-alias']:
            attrs.update(conf['transition-alias'][tag])
            tag = attrs.pop('tag')
        if tag not in conf['transition']:
            raise ValueError(f'transition.type: {type} is an unknown transition')
        trans = conf['transition'][tag]
        options = trans['default'] if (not options and 'default' in trans) else options
        for option in options:
            if option not in trans:
                raise ValueError(f'transition.type: "{type}" has invalid option {option}')
            attrs.update(trans[option])
        # Remove existing transition
        el = slide.element.find(qn('mc:AlternateContent'))
        if el is not None:
            slide.element.remove(el)
        # Add transition OXML
        # TODO: fails on slides with equations, zoom, or any other mc:alternateContent
        if tag != 'none':
            ns = trans.get('ns', 'p')
            attrs = ' '.join(f'{k}="{v}"' for k, v in attrs.items())
            xml = conf['transition-tmpl'][ns].format(tag=tag, attrs=attrs)
            el = parse_xml(xml)[0]
            slide.element.append(el)
    # Add attributes for duration: and advance:
    for key, attr in (('duration', qn('p14:dur')), ('advance', 'advTm')):
        val = commands.expr(spec.get(key, None), data)
        if val is not None:
            trans = slide.element.find('mc:AlternateContent/mc:Choice/p:transition', _nsmap)
            if trans is not None:
                trans.set(attr, f'{float(val) * 1000:.0f}')


def iterate_on(spec, data: dict):
    '''
    ``clone:`` and ``copy:`` iterate on data to return a (key, val) pair. This method performs the
    iteration for different data types and returns (key, val) consistently.

    :arg expr spec: an expression to iterate on. It can be a dict, tuple, list, pd.Index,
        pd.DataFrame, or pd.DataFrameGroupBy
    :arg dict data: data context for the ``spec`` expression
    :return: an iterable that yields (key, val) tuples
    '''
    val = commands.expr(spec, data)
    # {x: 1} -> (x, 1)
    if isinstance(val, dict):
        return val.items()
    # [x, y] -> (0, x), (1, y)
    elif isinstance(val, (tuple, list, pd.Index)):
        return enumerate(val)
    # pd.Series({x: 1}) -> (x, 1)
    elif isinstance(val, pd.Series):
        return val.iteritems()
    # pd.DataFrame([{x:1, y:2], {x:3, y:4}]) -> (0, {x: 1, y:2}), ...
    elif isinstance(val, pd.DataFrame):
        return val.iterrows()
    # df.groupby(key) -> (key1, group_df1), (key2, group_df2), ...
    elif isinstance(val, pd.core.groupby.generic.DataFrameGroupBy):
        return val
    else:
        raise ValueError(f'Cannot iterate over {type(val)}: {val!r}')


def commandline(args=None):
    '''Generates target PPTX from a source PPTX, applying rules in config file and opens it.

    Usage

        slidesense [config.yaml] [url-name] [--source=...] [--target=...] [--data=...]

    If no config file is specified, uses `gramex.yaml` in the current directory.

    The config file can have a pptgen configuration like {source: ..., target: ..., rules: ...}
    or be a `gramex.yaml` with `url: {url-name: {handler: PPTXHandler, kwargs: {source: ...}}}`
    Rules are picked up from the first PPTXHandler URL that matches `url-name`,
    or the first PPTXHandler in `gramex.yaml`.

    --source=... overrides source PPTX path in config file
    --target=... overrides target PPTX path in config file (defaults to output.pptx)
    --data=...   overrides data file in config path
    --no-open    don't open target PPTX after generating it
    '''
    args = gramex.parse_command_line(sys.argv[1:] if args is None else args)

    if 'help' in args or (not args['_'] and not os.path.exists('gramex.yaml')):
        return gramex.console(dedent(commandline.__doc__).strip())

    config_file, *urls = args.pop('_') or ['gramex.yaml']
    conf = gramex.cache.open(config_file, 'config')

    if 'url' in conf:
        for key, spec in conf.url.items():
            if spec.handler == 'PPTXHandler' and (not urls or any(url in key for url in urls)):
                rules = spec.kwargs
                break
        else:
            return app_log.error(f'No PPTXHandler matched in file: {config_file}')
    elif any(key in conf for key in ('source', 'target', 'data', 'rules')):
        rules = conf
    else:
        return app_log.error(f'No rules found in file: {config_file}')

    gramex.config.merge(rules, args)
    rules.setdefault('target', 'output.pptx')
    rules.setdefault('mode', 'expr')
    # Allow importing python files in current directory
    sys.path.append('.')
    # Generate output
    gramex.pptgen2.pptgen(**rules)
    # If --no-open is specified, or the OS doesn't have startfile (e.g. Linux), stop here.
    # Otherwise, open the output PPTX created
    if not rules.get('no-open', False) and hasattr(os, 'startfile'):
        # B606:start_process_with_no_shell is safe -- it's a file we've explicitly created
        os.startfile(rules['target'])  # nosec B606
