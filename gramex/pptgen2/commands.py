'''
pptgen2.pptgen() modifies slides using commands (in ``__init__.py``). This file maps the commands
and their functions in ``cmdlist``.

Each command accepts ``(shape, spec, data)`` and modifies shape based on spec(ification) and data.
'''

import copy
import gramex.cache
import io
import matplotlib.cm
import matplotlib.colors
import os
import numpy as np
import pandas as pd
import pptx
import pptx.util
import re
import requests
from PIL import Image
from functools import partial
from gramex import console
from gramex.config import app_log, objectpath
from gramex.transforms import build_transform

# B410:import_lxml lxml.etree is safe on https://github.com/tiran/defusedxml/tree/main/xmltestdata
from lxml.html import fragments_fromstring, builder, HtmlElement  # nosec B410
from orderedattrdict import AttrDict
from pptx.chart import data as pptxchartdata
from pptx.dml.color import RGBColor
from pptx.dml.fill import FillFormat
from pptx.enum.base import EnumValue
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import _nsmap, qn
from pptx.oxml.simpletypes import ST_Percentage
from pptx.oxml.xmlchemy import OxmlElement
from pptx.text.text import _Run
from urllib.parse import urlparse
from typing import Union, List

conf = gramex.cache.open('config.yaml', rel=True)
_nsmap.update(conf['nsmap'])


# Expression utilities
# ---------------------------------------------------------------------
def expr(val, data: dict = {}):
    if data.get('_expr_mode'):
        if isinstance(val, dict) and 'value' in val:
            val = val['value']
            val = val.format(**data) if isinstance(val, str) else val
        elif isinstance(val, str):
            vars = {key: None for key in data}
            val = build_transform({'function': str(val)}, vars=vars, iter=False)(**data)
    else:
        if isinstance(val, dict) and 'expr' in val:
            vars = {key: None for key in data}
            val = build_transform({'function': str(val['expr'])}, vars=vars, iter=False)(**data)
        elif isinstance(val, str):
            val = val.format(**data)
    return val


def assign(convert, path: str):
    '''assign(int, 'font.3.size') returns a method(para, '10') -> para.font[3].size = int(10)'''
    path = path.split('.')

    def method(node, value, data):
        for p in path[:-1]:
            node = node[int(p)] if p.isdigit() else getattr(node, p)
        # To clear the bold, italic, etc, we need to set it to None. So don't convert None
        setattr(node, path[-1], None if value is None else convert(value))

    return method


# Length utilities
# ---------------------------------------------------------------------
length_unit = pptx.util.Inches  # The default length unit is inches. This is exposed
_length_expr = re.compile(
    r'''          # A length expression can be:ex
    ((?:[+-]?)[\d\.]+)                  #   Any integer or floating point (with + or -)
    \s*                                 #   optionally followed by spaces
    ("|in|inch|inches|cm|mm|pt|cp|centipoint|centipoints|emu|)  # and a unit that may be blank
    $                                   # with nothing after that
''',
    re.VERBOSE,
)
# Standardize the aliases into pptx.util attributes
length_alias = {
    '"': 'inches',
    'in': 'inches',
    'inch': 'inches',
    'centipoint': 'centipoints',
    'cp': 'centipoints',
}


def length_class(unit_str: str):
    '''Converts unit string like in, inch, cm, etc to pptx.util.Inches, pptx.util.Cm, etc.'''
    unit = length_alias.get(unit_str.lower(), unit_str.lower()).title()
    if hasattr(pptx.util, unit):
        return getattr(pptx.util, unit)
    elif not unit:
        return length_unit
    else:
        raise ValueError(f'Invalid unit: {unit_str}')


def length(val: Union[str, int, float]) -> pptx.util.Length:
    '''Converts 3.2, 3.2" or 3.2 inch to pptx.util.Inches(3.2) -- with specified or default unit'''
    if isinstance(val, str):
        match = _length_expr.match(val)
        if match:
            return length_class(match.group(2))(float(match.group(1)))
    elif isinstance(val, (int, float, np.number)):
        return length_unit(val)
    raise ValueError('Invalid length: %r' % val)


# Color utilities
# ---------------------------------------------------------------------
color_map = matplotlib.colors.get_named_colors_mapping()
# Theme colors can start with any of these
theme_color = re.compile(
    r'''
    (ACCENT|BACKGROUND|DARK|LIGHT|TEXT|HYPERLINK|FOLLOWED_HYPERLINK)
    _?(\d*)             # followed by _1, _2, etc (or 1, 2, etc.)
    (\+\d+|\-\d+)?      # There MAY be a +30 or -40 after that to adjust brightness%
''',
    re.VERBOSE,
)


def fill_color(fill: FillFormat, val: Union[str, tuple, list, None]) -> None:
    '''
    Set the FillFormat color to value specified as a:

    - a named color, like ``black``
    - a hex value, like ``#f80`` or ``#ff8800``
    - an RGB value, like ``rgb(255, 255, 0)`` or ``rgb(1, 0.5, 0.1)``
    - a tuple or list of RGB values, like ``(255, 255, 0)`` or ``[255, 255, 0]``
    - a theme color, like ``ACCENT_1``, ``ACCENT_2``, ``BACKGROUND_1``, ``DARK_1``, ``LIGHT_2``
    - a theme color with a brightness modifier, like ``ACCENT_1+40``, which is 40% brighter than
      Accent 1, or ``ACCENT_2-20`` which is 20% darker than Accent 2
    - ``'none'`` clears the color, i.e. makes it transparent
    '''
    fill.solid()
    if val == 'none':
        fill.background()
    elif isinstance(val, (list, tuple)):
        val = val[:3]
        if any(isinstance(v, float) for v in val) and all(0 <= v <= 1 for v in val):
            fill.fore_color.rgb = RGBColor(*(int(v * 256 if v < 1 else 255) for v in val))
        else:
            fill.fore_color.rgb = RGBColor(*val)
    elif isinstance(val, str):
        val = color_map.get(val, val).upper()
        if val.startswith('#'):
            if len(val) == 7:
                fill.fore_color.rgb = RGBColor.from_string(val[1:])
            elif len(val) == 4:
                fill.fore_color.rgb = RGBColor.from_string(val[1] * 2 + val[2] * 2 + val[3] * 2)
        elif val.startswith('RGB('):
            parts = re.findall(r'\d+', val)
            fill.fore_color.rgb = RGBColor(int(parts[0]), int(parts[1]), int(parts[2]))
        else:
            match = theme_color.match(val)
            if match:
                theme = match.group(1) + ('_' + match.group(2) if match.group(2) else '')
                fill.fore_color.theme_color = getattr(MSO_THEME_COLOR, theme)
                if match.group(3):
                    fill.fore_color.brightness = float(match.group(3)) / 100
                # else: No brightness adjustment required
            else:
                raise ValueError('Invalid color: %r' % val)


def fill_opacity(fill: FillFormat, val: Union[int, float, None]) -> None:
    '''Set the FillFormat opacity to a number'''
    if fill.type != MSO_FILL.SOLID:
        raise ValueError('Cannot set opacity: %r on non-solid fill type %r' % (val, fill.type))
    for tag in ('hslClr', 'sysClr', 'srgbClr', 'prstClr', 'scrgbClr', 'schemeClr'):
        color = fill._xPr.find('.//' + qn(f'a:{tag}'))
        if color is not None:
            alpha = color.find(qn('a:alpha'))
            if alpha is None:
                alpha = OxmlElement('a:alpha')
                color.append(alpha)
            alpha.set('val', ST_Percentage.convert_to_xml(val))
            break


# Other unit conversions
# ---------------------------------------------------------------------
def binary(val: Union[str, int, float]) -> bool:
    '''Convert string to boolean'''
    if val in {'y', 'yes', 'true', 'Y', 'Yes', 'YES', 'True', 'TRUE', 1, 1.0, True}:
        return True
    if val in {'n', 'no', 'false', 'N', 'No', 'NO', 'False', 'FALSE', 1, 1.0, False}:
        return False
    raise ValueError('Invalid boolean: %r' % val)


def alignment(enum, val: str) -> EnumValue:
    '''Convert string to alignment value. enum can be PP_ALIGN or MSO_VERTICAL_ANCHOR'''
    alignment = getattr(enum, val.upper(), None)
    if alignment is not None:
        return alignment
    raise ValueError(f'Invalid alignment: {val}')


# Basic command methods
# ---------------------------------------------------------------------
def name(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        shape.name = val


def print_command(shape, spec, data: dict):
    spec = spec if isinstance(spec, (list, tuple)) else [spec]
    for item in spec:
        console(f'    {item}: {expr(item, data)}')


# Position & style commands
# ---------------------------------------------------------------------
def set_size(prop, method, incr, shape, spec, data: dict):
    '''Generator for top, left, width, height, move-top, move-left, add-width, add-height'''
    val = expr(spec, data)
    if val is not None:
        setattr(shape, prop, method(val) + (getattr(shape, prop) if incr else 0))
        flatten_group_transforms(shape)


def flatten_group_transforms(shape, x=0, y=0, sx=1, sy=1):
    '''
    Groups re-scale child shapes. So changes to size and position of child have unexpected results.
    To avoid this, ensure child rect (a:chOff, a:chExt) is same as group rect (a:off, a:ext).
    See https://stackoverflow.com/a/56485145/100904
    '''
    if isinstance(shape, pptx.shapes.group.GroupShape):
        # Get group and child rect shapes
        xfrm = shape.element.find(qn('p:grpSpPr')).find(qn('a:xfrm'))
        off, ext = xfrm.find(qn('a:off')), xfrm.find(qn('a:ext'))
        choff, chext = xfrm.find(qn('a:chOff')), xfrm.find(qn('a:chExt'))
        # Calculate how much the child rect is transformed from the group rect
        tsx = sx * ext.cx / chext.cx
        tsy = sy * ext.cy / chext.cy
        tx = x + sx * (off.x - choff.x * ext.cx / chext.cx)
        ty = y + sy * (off.y - choff.y * ext.cy / chext.cy)
        # Transform child shapes using (x, y, sx, sy)
        for subshape in shape.shapes:
            flatten_group_transforms(subshape, tx, ty, tsx, tsy)
    # Transform shape positions using (x, y, sx, sy). But get all values before setting.
    # Placeholders need this: https://stackoverflow.com/a/49306814/100904
    left, top = int(shape.left * sx + x + 0.5), int(shape.top * sy + y + 0.5)
    width, height = int(shape.width * sx + 0.5), int(shape.height * sy + 0.5)
    shape.left, shape.top, shape.width, shape.height = left, top, width, height
    # Set child rect coords same as group rect (AFTER scaling the group)
    if isinstance(shape, pptx.shapes.group.GroupShape):
        choff.x, choff.y, chext.cx, chext.cy = off.x, off.y, ext.cx, ext.cy


def zoom(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        val = float(val)
        shape.left -= int((val - 1) * shape.width / 2)
        shape.top -= int((val - 1) * shape.height / 2)
        shape.width = int(shape.width * val)
        shape.height = int(shape.height * val)


def set_color(attr, shape, spec, data: dict):
    '''Generator for fill, stroke commands'''
    val = expr(spec, data)
    if val is not None:
        fill_color(objectpath(shape, attr), val)


def set_opacity(attr, shape, spec, data: dict):
    '''Generator for fill-opacity, stroke-opacity commands'''
    val = expr(spec, data)
    if val is not None:
        fill_opacity(objectpath(shape, attr), val)


def stroke_width(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        shape.line.width = length(val)


def set_adjustment(index, shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        shape.adjustments[index] = val


# Link commands
# python-pptx can't set slide links on runs, nor tooltips. Make our own
# ---------------------------------------------------------------------
def get_or_add_link(shape, type, event):
    # el is the lxml element of the shape. Can be ._element or ._run
    el = getattr(shape, conf['link-attrs'][type]['el'])
    # parent is the container inside which we insert the link. Can be p:cNvPr or a:rPr
    parent = el.find('.//' + conf['link-attrs'][type]['tag'], _nsmap)
    # Create link if required. (Else preserve all attributes one existing link, like tooltip)
    link = parent.find(qn(f'a:{event}'))
    if link is None:
        link = OxmlElement(f'a:{event}')
        parent.insert(0, link)
    return link


def set_link(type, event, prs, slide, shape, val):
    if val is None:
        return
    link = get_or_add_link(shape, type, event)
    # Set link's r:id= and action= based on the type of the link
    if val in conf['link-action']:  # it's a ppaction://
        rid = link.get(qn('r:id'), None)
        if rid in slide.part.rels:
            slide.part.drop_rel(rid)
        link.set(qn('r:id'), '')
        link.set('action', conf['link-action'][val])
    elif isinstance(val, int) or val.isdigit():  # it's a slide
        slide_part = prs.slides[int(val) - 1].part
        link.set(qn('r:id'), slide.part.relate_to(slide_part, RT.SLIDE, False))
        link.set('action', 'ppaction://hlinksldjump')
    elif urlparse(val).netloc:  # it's a URL
        link.set(qn('r:id'), slide.part.relate_to(val, RT.HYPERLINK, True))
        link.attrib.pop('action', '')
    elif os.path.splitext(val)[-1].lower() in conf['link-ppt-files']:  # it's a PPT
        link.set(qn('r:id'), slide.part.relate_to(val, RT.HYPERLINK, True))
        link.set('action', 'ppaction://hlinkpres?slideindex=1&slidetitle=')
    else:  # it's a file
        link.set(qn('r:id'), slide.part.relate_to(val, RT.HYPERLINK, True))
        link.set('action', 'ppaction://hlinkfile')
    return link


def set_tooltip(type, prs, slide, shape, val):
    if val is None:
        return
    link = get_or_add_link(shape, type, 'hlinkClick')
    rid, action = link.get(qn('r:id'), None), link.get('action', None)
    # Note: link was just created, or has no action, we can't set a tooltip. So create a next: link
    if rid is None or action is conf['link-action']['noaction']:
        link = set_link(type, 'hlinkClick', prs, slide, shape, 'next')
    link.set('tooltip', val)
    # PS: link: back doesn't work with tooltip


# Text commands
# ---------------------------------------------------------------------
def get_elements(children: List[Union[str, HtmlElement]], element_builder):
    '''
    Convert head & tail text into specified tag, so that we can just process all text and elements
    as a single list.
    '''
    elements, open = [], True
    for e in children:
        if isinstance(e, str):
            e = re.sub(r'\s+', ' ', e, re.DOTALL)
            if e:
                if open and len(elements):
                    elements[-1].text += e
                else:
                    elements.append(element_builder(e))
            open = True
        elif isinstance(e, HtmlElement):
            if e.tag != element_builder.args[0]:
                if open and len(elements):
                    elements[-1].append(e)
                else:
                    elements.append(element_builder(e))
                open = True
            else:
                elements.append(e)
                open = False
                tail = re.sub(r'\s+', ' ', e.tail, re.DOTALL) if e.tail else ''
                if tail.strip():
                    elements.append(element_builder(tail))
                    open = True

    return elements or [element_builder('')]


def baseline(run, val, data):
    if isinstance(val, str):
        val = val.strip().lower()
        # See https://github.com/scanny/python-pptx/pull/601 for unit values
        val = (
            0.3
            if val.startswith('sup')
            else -0.25
            if val.startswith('sub')
            else ST_Percentage._convert_from_percent_literal(val)
        )
    run.font._rPr.set('baseline', ST_Percentage.convert_to_xml(val))


def strike(run, val, data):
    val = (
        'sngStrike'
        if val.startswith('s')
        else 'dblStrike'
        if val.startswith('d')
        else 'noStrike'
        if val.startswith('n')
        else val
    )
    run.font._rPr.set('strike', val)


def get_text_frame(shape):
    # Get the text frame. Note: don't use .has_text_frame. This function should work for cells too
    try:
        return shape.text_frame
    except AttributeError:
        raise ValueError(f'Cannot set text on shape {shape.name} that has no text frame')


para_methods = {
    'align': assign(partial(alignment, PP_ALIGN), 'alignment'),
    'bold': assign(binary, 'font.bold'),
    'color': lambda p, v, d: fill_color(p.font.fill, v),
    'font-name': assign(str, 'font.name'),
    'font-size': assign(length, 'font.size'),
    'italic': assign(binary, 'font.italic'),
    'level': assign(int, 'level'),
    'line-spacing': assign(length, 'line_spacing'),
    'space-after': assign(length, 'space_after'),
    'space-before': assign(length, 'space_before'),
    'underline': assign(binary, 'font.underline'),
}
run_methods = {
    'baseline': baseline,
    'bold': assign(binary, 'font.bold'),
    'color': lambda r, v, d: fill_color(r.font.fill, v),
    'font-name': assign(str, 'font.name'),
    'font-size': assign(length, 'font.size'),
    'hover': lambda r, v, d: set_link('text', 'hlinkMouseOver', d['prs'], d['slide'], r, v),
    'italic': assign(binary, 'font.italic'),
    'link': lambda r, v, d: set_link('text', 'hlinkClick', d['prs'], d['slide'], r, v),
    'strike': strike,
    'tooltip': lambda r, v, d: set_tooltip('text', d['prs'], d['slide'], r, v),
    'underline': assign(binary, 'font.underline'),
}


def text(shape, spec, data):
    '''Set text on shape. Allows paragraph and run formatting with a HTML-like syntax'''
    val = expr(spec, data)
    if val is None:
        return
    frame = get_text_frame(shape)
    # Delete all but the first para (required), and its first run (to preserve formatting).
    # All new paras and runs get the same formatting as the first para & run, by default
    for p in frame.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for r in frame.paragraphs[0].runs[1:]:
        r._r.getparent().remove(r._r)
    para = frame.paragraphs[0]
    para_defaults = para._pPr.attrib
    run_defaults = para.runs[0].font._rPr.attrib if len(para.runs) > 0 else {}
    # Apply the text formatting
    tree = fragments_fromstring(str(val))
    for i, para in enumerate(get_elements(tree, builder.P)):
        # The first para is guaranteed to exist. After that, add a para
        p = frame.add_paragraph() if i >= len(frame.paragraphs) else frame.paragraphs[i]
        # Ensure that all paras have the same alignment, spacing, etc as the first para
        p._pPr.attrib.update(para_defaults)
        # Set specified para attributes
        for attr, val in para.attrib.items():
            if attr in para_methods:
                para_methods[attr](p, val, data)
        for j, run in enumerate(get_elements([para.text] + list(para), builder.A)):
            # The first run MAY exist, if there was text. Or not. Create if required
            r = p.add_run() if j >= len(p.runs) else p.runs[j]
            # Ensure that all runs have the same color, font, etc as the first run (if any)
            r.font._rPr.attrib.update(run_defaults)
            # Set specified run attributes
            r.text = run.text or ''
            for attr, val in run.attrib.items():
                if attr in run_methods:
                    run_methods[attr](r, val, data)


def replace(shape, spec, data):
    if not isinstance(spec, dict):
        raise ValueError('replace: needs a dict of {old: new} text, not %r' % spec)
    frame = get_text_frame(shape)

    def insert_run_after(r, p, original_r):
        new_r = copy.deepcopy(original_r)
        r._r.addnext(new_r)
        return _Run(new_r, p)

    spec = {re.compile(old): fragments_fromstring(expr(new, data)) for old, new in spec.items()}
    for p in frame.paragraphs:
        for r in p.runs:
            for old, tree in spec.items():
                match = old.search(r.text)
                if match:
                    original_r, prefix, suffix = (
                        r._r,
                        r.text[: match.start()],
                        r.text[match.end() :],
                    )
                    if prefix:
                        r.text = prefix
                        r = insert_run_after(r, p, original_r)
                    for j, run in enumerate(get_elements(tree, builder.A)):
                        r = insert_run_after(r, p, original_r) if j > 0 else r
                        r.text = run.text
                        for attr, val in run.attrib.items():
                            if attr in run_methods:
                                run_methods[attr](r, val, data)
                    if suffix:
                        # Ensure suffix has same attrs as original text
                        r = insert_run_after(r, p, original_r)
                        r.text = suffix


def set_text(attr, on_run, shape, spec, data):
    '''Generator for bold, italic, color, and other text commands'''
    val = expr(spec, data)
    if val is not None:
        frame = get_text_frame(shape)
        para_method, run_method = para_methods[attr], run_methods[attr]
        for para in frame.paragraphs:
            para_method(para, val, data)
            for run in para.runs:
                # If on_run=False, clear attr on every run (e.g. bold, italic)
                # Else, set attr on every run too (e.g. color)
                # (because PPT doesn't display use color on para.)
                run_method(run, val if on_run else None, data)


# Image commands
# ---------------------------------------------------------------------


def image(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        # Load image contents as a bytestring
        if urlparse(val).netloc:
            content = requests.get(val).content
        else:
            content = gramex.cache.open(val, 'bin')
        # Add the image part
        image_part, rid = shape.part.get_or_add_image_part(io.BytesIO(content))
        shape.element.blipFill.blip.rEmbed = rid
        # Preserve aspect ratio and width. Adjust height via a "cover" algorithm
        img = Image.open(io.BytesIO(content))
        img_width, img_height = img.size
        shape.height = int(shape.width * img_height / img_width)


def image_width(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        val = length(val)
        shape.width, shape.height = val, int(val * shape.height / shape.width)


def image_height(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        val = length(val)
        shape.width, shape.height = int(val * shape.width / shape.height), val


# Table commands
# ---------------------------------------------------------------------
def table_align(shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        val = alignment(PP_ALIGN, val)
        for para in shape.text_frame.paragraphs:
            para.alignment = val


def table_assign(convert, attr, shape, spec, data: dict):
    val = expr(spec, data)
    if val is not None:
        setattr(shape, attr, convert(val))


def table_width(table, col_index, spec, data):
    table._tbl.tblGrid.gridCol_lst[col_index].set('w', str(length(expr(spec, data))))


def _resize(elements, n):
    '''Ensure that a tr, tc or gridCol list has n children by cloning or deleting last element'''
    for index in range(len(elements), n, -1):
        elements[index - 1].delete()
    for index in range(len(elements), n):
        last_element = copy.deepcopy(elements[-1])
        # When copying the last element, remove the a:extLst. This may have a rowId and colId
        # Note: Anand isn't sure if this removes any useful properties. See https://bit.ly/2ZYDw7B
        for ext_lst in last_element.findall('.//' + qn('a:extLst')):
            ext_lst.getparent().remove(ext_lst)
        elements[-1].addnext(last_element)


table_cell_commands = {
    'text': text,
    'fill': partial(set_color, 'fill'),
    'fill-opacity': partial(set_opacity, 'fill'),
    'align': table_align,
    # TODO: table borders
    'vertical-align': partial(
        table_assign, partial(alignment, MSO_VERTICAL_ANCHOR), 'vertical_anchor'
    ),
    'margin-left': partial(table_assign, length, 'margin_left'),
    'margin-right': partial(table_assign, length, 'margin_right'),
    'margin-top': partial(table_assign, length, 'margin_top'),
    'margin-bottom': partial(table_assign, length, 'margin_bottom'),
    # TODO: width: column-wise width
    'bold': partial(set_text, 'bold', False),
    'color': partial(set_text, 'color', True),  # PPT needs colors on runs too, not only paras
    'font-name': partial(set_text, 'font-name', False),
    'font-size': partial(set_text, 'font-size', False),
    'italic': partial(set_text, 'italic', False),
    'underline': partial(set_text, 'underline', False),
}
table_col_commands = {
    'width': table_width,
}


def table(shape, spec, data: dict):
    if not shape.has_table:
        raise ValueError(f'Cannot run table commands on shape {shape.name} that is not a table')
    table = shape.table

    # Set or get table first/last row/col attributes
    header_row = expr(spec.get('header-row', table.first_row), data)
    table.first_row = bool(header_row)
    table.last_row = bool(expr(spec.get('total-row', table.last_row), data))
    table.first_col = bool(expr(spec.get('first-column', table.first_col), data))
    table.last_col = bool(expr(spec.get('last-column', table.last_col), data))

    # table data must be a DataFrame if specified. Else, create a DataFrame from existing text
    table_data = expr(spec.get('data', None), data)
    if table_data is not None and not isinstance(table_data, pd.DataFrame):
        raise ValueError(f'data on table {shape.name} must be a DataFrame, not {table_data!r}')
    # Extract data from table text if no data is specified.
    if table_data is None:
        table_data = pd.DataFrame([[cell.text for cell in row.cells] for row in table.rows])
        # If the PPTX table has a header row, set the first row as the DataFrame header too
        if table.first_row:
            table_data = table_data.T.set_index([0]).T

    # Adjust PPTX table size to data table size
    header_offset = 1 if table.first_row else 0
    _resize(table._tbl.tr_lst, len(table_data) + header_offset)
    data_cols = len(table_data.columns)
    _resize(table._tbl.tblGrid.gridCol_lst, data_cols)
    for row in table.rows:
        _resize(row._tr.tc_lst, data_cols)

    # Set header row text from header-row or from data column headers
    if table.first_row:
        header_columns = table_data.columns
        if isinstance(header_row, (list, tuple, pd.Index, pd.Series)):
            header_columns = header_row[: len(table_data.columns)]
        for j, column in enumerate(header_columns):
            text(table.cell(0, j), column, {'_expr_mode': False})

    # If `text` is not specified, just use the table value
    expr_mode = data.get('_expr_mode')
    spec.setdefault('text', 'cell.val' if expr_mode else {'expr': 'cell.val'})

    # TODO: Handle nans
    # Apply table commands. (Copy data to avoid modifying original. We'll add data['cell'] later)
    data = dict(data)
    columns = table_data.columns.tolist()
    for key, cmdspec in spec.items():
        # The command spec can be an expression, or a dict of expressions for each column.
        # Always convert into a {column: expression}.
        # But carefully, handling {value: ...} in expr mode and {expr: ...} in literal mode
        if (
            not isinstance(cmdspec, dict)
            or (expr_mode and 'value' in cmdspec)
            or (not expr_mode and 'expr' in cmdspec)
        ):
            cmdspec = {column: cmdspec for column in table_data.columns}
        # Apply commands that run on each cell
        if key in table_cell_commands:
            for i, (index, row) in enumerate(table_data.iterrows()):
                for j, (column, val) in enumerate(row.iteritems()):
                    data['cell'] = AttrDict(
                        val=val,
                        column=column,
                        index=index,
                        row=row,
                        data=table_data,
                        pos=AttrDict(row=i, column=j),
                    )
                    cell = table.cell(i + header_offset, j)
                    if column in cmdspec:
                        table_cell_commands[key](cell, cmdspec[column], data)
            for column in cmdspec:
                if column not in columns:
                    app_log.warn(f'pptgen2: No column: {column} in table: {shape.name}')
        # Apply commands that run on each column
        elif key in table_col_commands:
            for column, colspec in cmdspec.items():
                if column in columns:
                    col_index = columns.index(column)
                    data['cell'] = AttrDict(
                        val=column,
                        column=column,
                        index=None,
                        row=table.columns,
                        data=table_data,
                        pos=AttrDict(row=-1, column=col_index),
                    )
                    table_col_commands[key](table, col_index, colspec, data)
                else:
                    app_log.warn(f'pptgen2: No column: {column} in table: {shape.name}')


def chart(shape, spec, data: dict):
    # Ensure that the chart is of a type we support
    if not hasattr(shape, 'chart'):
        raise ValueError(f'Cannot run chart: on non-chart shape: {shape.name}')
    chart_tag = shape.chart.plots._plotArea.xCharts[0].tag
    if chart_tag not in conf['chart-type']:
        raise ValueError(f'Unsupported chart type {chart_tag} on shape: {shape.name}')

    # Set the chart data
    chart_data = expr(spec.get('data', None), data)
    # If it's specified, it must be a DataFrame
    if chart_data is not None and not isinstance(chart_data, pd.DataFrame):
        raise ValueError(f'Chart data {chart_data:r} is not a DataFrame on shape: {shape.name}')
    # If it's not specified, use the existing data
    if chart_data is None:
        chart_data = pd.read_excel(
            io.BytesIO(shape.chart.part.chart_workbook.xlsx_part.blob),
            index_col=0,
            engine='openpyxl',
        )
    else:
        chart_type = conf['chart-type'][chart_tag]
        # Create a new instance of CategoryChartData(), XYChartData() or BubbleChartData()
        new_chart_data = getattr(pptxchartdata, chart_type + 'ChartData')()
        new_chart_data.categories = chart_data.index
        if chart_type == 'Category':
            for name, col in chart_data.iteritems():
                new_chart_data.add_series(name, col.values)
        # TODO: This messes up the resulting Excel sheet, and is not extensible. Rewrite via lxml
        elif chart_type == 'Xy':
            for name, col in chart_data.iteritems():
                series = new_chart_data.add_series(name)
                for index, v in col.iteritems():
                    series.add_data_point(index, v)
        shape.chart.replace_data(new_chart_data)

    data = dict(**data, chartdata=chart_data)
    for attrname, method in (
        ('fill', partial(set_color, 'format.fill')),
        ('stroke', partial(set_color, 'format.line.fill')),
        ('text', lambda shape, spec, data: text(shape.data_label, spec, data)),
    ):
        val = expr(spec.get(attrname, None), data)
        if val is not None:
            for series in shape.chart.series:
                if series.name in val.columns:
                    vals = val[series.name].tolist()
                    for point_index, point in enumerate(series.points):
                        if point_index < len(vals) and vals[point_index] is not None:
                            method(point, vals[point_index], data)


cmdlist = {
    # Basic commands
    'name': name,
    'print': print_command,
    # Position
    'top': partial(set_size, 'top', length, False),
    'left': partial(set_size, 'left', length, False),
    'width': partial(set_size, 'width', length, False),
    'height': partial(set_size, 'height', length, False),
    'rotation': partial(set_size, 'rotation', float, False),
    'add-top': partial(set_size, 'top', length, True),
    'add-left': partial(set_size, 'left', length, True),
    'add-width': partial(set_size, 'width', length, True),
    'add-height': partial(set_size, 'height', length, True),
    'add-rotation': partial(set_size, 'rotation', float, True),
    # Shape
    'zoom': zoom,
    'adjustment1': partial(set_adjustment, 0),
    'adjustment2': partial(set_adjustment, 1),
    'adjustment3': partial(set_adjustment, 2),
    'adjustment4': partial(set_adjustment, 3),
    # Style
    'fill': partial(set_color, 'fill'),
    'stroke': partial(set_color, 'line.fill'),
    'fill-opacity': partial(set_opacity, 'fill'),
    'stroke-opacity': partial(set_opacity, 'line.fill'),
    'stroke-width': stroke_width,
    # Image
    'image': image,
    'image-width': image_width,
    'image-height': image_height,
    # Link
    'link': lambda shape, spec, data: (
        set_link('shape', 'hlinkClick', data['prs'], data['slide'], shape, expr(spec, data))
    ),
    'hover': lambda shape, spec, data: (
        set_link('shape', 'hlinkHover', data['prs'], data['slide'], shape, expr(spec, data))
    ),
    'tooltip': lambda shape, spec, data: (
        set_tooltip('shape', data['prs'], data['slide'], shape, expr(spec, data))
    ),
    # Text
    'replace': replace,
    'text': text,
    'bold': partial(set_text, 'bold', False),
    'color': partial(set_text, 'color', True),  # PPT needs colors on runs too, not only paras
    'font-name': partial(set_text, 'font-name', False),
    'font-size': partial(set_text, 'font-size', False),
    'italic': partial(set_text, 'italic', False),
    'underline': partial(set_text, 'underline', False),
    # Others
    'table': table,
    'chart': chart,
    'chart-data': lambda shape, spec, data: chart(shape, {'data': spec}, data),
    # Custom charts
    #   sankey
    #   bullet
    #   treemap
    #   heatgrid
    #   calendarmap
}
