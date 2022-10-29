'''Utility file.'''
import re
import ast
import copy
import platform
import numpy as np
import pandas as pd

# B410:import_lxml lxml.etree is safe on https://github.com/tiran/defusedxml/tree/main/xmltestdata
from lxml import objectify  # nosec B410
from lxml.builder import ElementMaker  # nosec B410
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.base import EnumValue
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from gramex.transforms import build_transform


def is_slide_allowed(change, slide, number):
    '''
    Given a change like one of the below::

        slide-number: 1
        slide-number: [1, 2, 3]
        slide-title: 'company'
        slide-title: ['company', 'industry']

    ... return True if:

    1. ``number`` matches a slide-number
    2. ``slide`` title matches a slide-title regex (case-insensitive)

    If none of these are specified, return True.
    '''
    match = True
    # Restrict to specific slide number(s), if specified
    if 'slide-number' in change:
        slide_number = change['slide-number']
        if isinstance(slide_number, (list, dict)):
            match = match and number in slide_number
        elif isinstance(slide_number, int):
            match = match and number == slide_number

    # Restrict to specific slide title(s), if specified
    if 'slide-title' in change:
        slide_title = change['slide-title']
        title = slide.shapes.title
        title = title.text if title is not None else ''
        if isinstance(slide_title, (list, dict)):
            match = match and any(re.search(expr, title, re.IGNORECASE) for expr in slide_title)
        elif isinstance(slide_title, str):
            match = match and re.search(slide_title, title, re.IGNORECASE)
    return match


def stack_elements(replica, shape, stack=False, margin=None):
    '''Function to extend elements horizontally or vertically.'''
    if not stack:
        return
    config = {
        'vertical': {'axis': 'y', 'attr': 'height'},
        'horizontal': {'axis': 'x', 'attr': 'width'},
    }
    grp_sp = shape.element
    # Adding a 15% default margin between original and new object.
    default_margin = 0.15
    margin = margin if margin else default_margin
    for index in range(replica):
        # Adding a cloned object to shape
        extend_shape = copy.deepcopy(grp_sp)
        # Getting attributes and axis values from config based on stack.
        attr = config.get(stack, {}).get('attr', 0)
        axis = config.get(stack, {}).get('axis', 0)
        # Taking width or height based on stack value and setting a margin.
        metric_val = getattr(shape, attr)
        axis_val = getattr(extend_shape, axis)
        # Setting margin accordingly either vertically or horizontally.
        axis_pos = metric_val * index
        set_attr = axis_val + axis_pos + (axis_pos * margin)
        # Setting graphic position of newly created object to slide.
        setattr(extend_shape, axis, int(set_attr))
        # Adding newly created object to slide.
        grp_sp.addprevious(extend_shape)
    shape.element.delete()


def stack_shapes(collection, change, data, handler):
    '''
    Function to stack Shapes if required.
    '''
    data_len = len(data)
    for shape in collection:
        if shape.name not in change:
            continue
        info = change[shape.name]
        if 'data' in info and info.get('stack') is not None:
            _vars = {'_color': None, 'data': None, 'handler': None}
            if not isinstance(info['data'], (dict,)):
                info['data'] = {'function': '{}'.format(info['data'])}
            elif isinstance(info['data'], (dict,)) and 'function' not in info['data']:
                info['data'] = {'function': '{}'.format(info['data'])}
            args = {'data': data, 'handler': handler}
            data_len = len(build_transform(info['data'], vars=_vars)(**args)[0])
        stack_elements(data_len, shape, stack=info.get('stack'), margin=info.get('margin'))


def delete_paragraph(paragraph):
    '''Delete a paragraph.'''
    p = paragraph._p
    parent_element = p.getparent()
    parent_element.remove(p)


def delete_run(run):
    '''Delete a run from paragraph.'''
    r = run._r
    r.getparent().remove(r)


def generate_slide(prs, source):
    '''Create a slide layout.'''
    layout_items_count = [len(layout.placeholders) for layout in prs.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return prs.slide_layouts[blank_layout_id]


def copy_slide_elem(shape, dest):
    '''
    Function to copy slide elements into a newly created slide.
    '''
    if dest is None:
        return
    new_elem = copy.deepcopy(shape.element)
    dest.shapes._spTree.insert_element_before(new_elem, 'p:extLst')


def add_new_slide(dest, source_slide):
    '''Function to add a new slide to presentation.'''
    if dest is None:
        return
    for value in source_slide.part.rels.values():
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" in value.reltype:
            continue
        dest.part.rels.add_relationship(value.reltype, value._target, value.rId)


def move_slide(presentation, old_index, new_index):
    '''Move a slide's index number.'''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def delete_slide(presentation, index):
    '''Delete a slide from Presentation.'''
    rid = presentation.slides._sldIdLst[index].rId
    presentation.part.drop_rel(rid)
    del presentation.slides._sldIdLst[index]


def manage_slides(prs, config):
    '''
    Delete not required slides from presentation.

    if `config.only` is present then remove the other slides apart from `config.only`
    slides from the presentation.
    `config.only` accepts a slide number or list of slide numbers starting from 1.
    '''
    slide_numbers = config.pop('only', None)
    if slide_numbers:
        if isinstance(slide_numbers, int):
            slide_numbers = {int(slide_numbers) - 1}
        elif isinstance(slide_numbers, list):
            slide_numbers = {int(i) - 1 for i in slide_numbers}
        else:
            raise ValueError('Slide numbers must be a list of integers or a single slide number.')
        slides = set(range(len(prs.slides)))
        for remove_status, slide_num in enumerate(sorted(slides - slide_numbers)):
            delete_slide(prs, slide_num - remove_status)
    return prs


def is_group(shape):
    # TODO: implement this
    return shape.element.tag.endswith('}grpSp')


def pixel_to_inch(pixel):
    '''Function to convert Pixel to Inches based on OS.'''
    linux_width = 72.0
    windows_width = 96.0
    os_name = platform.system().lower().strip()
    if os_name == 'windows':
        return Inches(pixel / windows_width)
    return Inches(pixel / linux_width)


def scale(series, lo=None, hi=None):
    '''
    Returns the values linearly scaled from 0 - 1.

    The lowest value becomes 0, the highest value becomes 1, and all other
    values are proportionally multiplied and have a range between 0 and 1.

    :arg Series series: Data to scale. Pandas Series, numpy array, list or iterable
    :arg float lo: Value that becomes 0. Values lower than ``lo`` in ``series``
        will be mapped to negative numbers.
    :arg float hi: Value that becomes 1. Values higher than ``hi`` in ``series``
        will be mapped to numbers greater than 1.

    Examples::

        >>> stats.scale([1, 2, 3, 4, 5])
        ... array([ 0.  ,  0.25,  0.5 ,  0.75,  1.  ])

        >>> stats.scale([1, 2, 3, 4, 5], lo=2, hi=4)
        ... array([-0.5,  0. ,  0.5,  1. ,  1.5])
    '''
    series = np.array(series, dtype=float)
    lo = np.nanmin(series) if lo is None or np.isnan(lo) else lo
    hi = np.nanmax(series) if hi is None or np.isnan(hi) else hi
    return (series - lo) / ((hi - lo) or np.nan)


def decimals(series):
    '''
    Given a ``series`` of numbers, returns the number of decimals
    *just enough* to differentiate between most numbers.

    :arg Series series: Pandas Series, numpy array, list or iterable.
        Data to find the required decimal precision for
    :return: The minimum number of decimals required to differentiate between
        most numbers

    Examples::

        stats.decimals([1, 2, 3])       # 0: All integers. No decimals needed
        stats.decimals([.1, .2, .3])    # 1: 1 decimal is required
        stats.decimals([.01, .02, .3])  # 2: 2 decimals are required
        stats.decimals(.01)             # 2: Only 1 no. of 2 decimal precision

    Note: This function first calculates the smallest difference between any pair
    of numbers (ignoring floating-point errors). It then finds the log10 of that
    difference, which represents the minimum decimals required to differentiate
    between these numbers.
    '''
    series = np.ma.masked_array(series, mask=np.isnan(series)).astype(float)
    series = series.reshape((series.size,))
    diffs = np.diff(series[series.argsort()])
    inf_diff = 1e-10
    min_float = 0.999999
    diffs = diffs[diffs > inf_diff]
    if len(diffs) > 0:
        smallest = np.nanmin(diffs.filled(np.Inf))
    else:
        nonnan = series.compressed()
        smallest = (abs(nonnan[0]) or 1) if len(nonnan) > 0 else 1
    return int(max(0, np.floor(min_float - np.log10(smallest))))


def convert_color_code(colorcode):
    '''Convert color code to valid PPTX color code.'''
    colorcode = colorcode.rsplit('#')[-1].lower()
    return colorcode + ('0' * (6 - len(colorcode)))


# Custom Charts Functions below(Sankey, Treemap, Calendarmap).


def apply_text_css(shape, run, paragraph, **kwargs):
    '''Apply css.'''
    pixcel_to_inch = 10000
    if kwargs.get('color'):
        rows_text = run.font.fill
        rows_text.solid()
        run.font.color.rgb = RGBColor.from_string(convert_color_code(kwargs['color']))
    if kwargs.get('font-family'):
        run.font.name = kwargs['font-family']
    if kwargs.get('font-size'):
        run.font.size = pixcel_to_inch * float(kwargs['font-size'])
    if kwargs.get('text-align'):
        if isinstance(kwargs['text-align'], EnumValue) or None:
            paragraph.alignment = kwargs['text-align']
        else:
            paragraph.alignment = getattr(PP_ALIGN, kwargs['text-align'].upper())
    for prop in {'bold', 'italic', 'underline'}:
        update_prop = kwargs.get(prop)
        if update_prop and not isinstance(update_prop, bool):
            update_prop = ast.literal_eval(update_prop)
        setattr(run.font, prop, update_prop)
    if kwargs.get('text-anchor'):
        shape.vertical_anchor = getattr(MSO_ANCHOR, shape['text-anchor'].upper())


def make_element():
    '''Function to create element structure.'''
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    a = ElementMaker(namespace=nsmap['a'], nsmap=nsmap)
    p = ElementMaker(namespace=nsmap['p'], nsmap=nsmap)
    r = ElementMaker(namespace=nsmap['r'], nsmap=nsmap)
    return {'nsmap': nsmap, 'a': a, 'p': p, 'r': r}


def fill_color(**kwargs):
    '''
    Return a new color object.

    You may use any one of the following ways of specifying colour:

        color(schemeClr='accent2')             # = second theme color
        color(prstClr='black')                 # = #000000
        color(hslClr=[14400000, 100.0, 50.0])  # = #000080
        color(sysClr='windowText')             # = window text color
        color(scrgbClr=(50000, 50000, 50000))  # = #808080
        color(srgbClr='aaccff')                # = #aaccff

    One or more of these modifiers may be specified:

    - alpha    : '10%' indicates 10% opacity
    - alphaMod : '10%' increased alpha by 10% (50% becomes 55%)
    - alphaOff : '10%' increases alpha by 10 points (50% becomes 60%)
    - blue     : '10%' sets the blue component to 10%
    - blueMod  : '10%' increases blue by 10% (50% becomes 55%)
    - blueOff  : '10%' increases blue by 10 points (50% becomes 60%)
    - comp     : True for opposite hue on the color wheel (e.g. red -> cyan)
    - gamma    : True for the sRGB gamma shift of the input color
    - gray     : True for the grayscale version of the color
    - green    : '10%' sets the green component to 10%
    - greenMod : '10%' increases green by 10% (50% becomes 55%)
    - greenOff : '10%' increases green by 10 points (50% becomes 60%)
    - hue      : '14400000' sets the hue component to 14400000
    - hueMod   : '600000' increases hue by 600000 (14400000 becomes 20000000)
    - hueOff   : '10%' increases hue by 10 points (50% becomes 60%)
    - inv      : True for the inverse color. R, G, B are all inverted
    - invGamma : True for the inverse sRGB gamma shift of the input color
    - lum      : '10%' sets the luminance component to 10%
    - lumMod   : '10%' increases luminance by 10% (50% becomes 55%)
    - lumOff   : '10%' increases luminance by 10 points (50% becomes 60%)
    - red      : '10%' sets the red component to 10%
    - redMod   : '10%' increases red by 10% (50% becomes 55%)
    - redOff   : '10%' increases red by 10 points (50% becomes 60%)
    - sat      : '100000' sets the saturation component to 100%
    - satMod   : '10%' increases saturation by 10% (50% becomes 55%)
    - satOff   : '10%' increases saturation by 10 points (50% becomes 60%)
    - shade    : '10%' is 10% of input color, 90% black
    - tint     : '10%' is 10% of input color, 90% white

    Refer
    <http://msdn.microsoft.com/en-in/library/documentformat.openxml.drawing(v=office.14).aspx>
    '''
    hslclr = kwargs.get('hslclr')
    sysclr = kwargs.get('sysclr')
    srgbclr = kwargs.get('srgbclr')
    prstclr = kwargs.get('prstclr')
    scrgbclr = kwargs.get('scrgbclr')
    schemeclr = kwargs.get('schemeclr')

    ns = xmlns('a')
    srgbclr = srgbclr.rsplit('#')[-1].lower()
    srgbclr = srgbclr + ('0' * (6 - len(srgbclr)))
    if schemeclr:
        s = f'<a:schemeClr {ns} val="{schemeclr}"/>'
    elif srgbclr:
        s = f'<a:srgbClr {ns} val="{srgbclr}"/>'
    elif prstclr:
        s = f'<a:prstClr {ns} val="{prstclr}"/>'
    elif hslclr:
        s = f'<a:hslClr {ns} hue="%.0f" sat="%.2f%%" lum="%.2f%%"/>' % tuple(hslclr)
    elif sysclr:
        s = f'<a:sysClr {ns} val="{sysclr}"/>'
    elif scrgbclr:
        s = f'<a:scrgbClr {ns} r="%.0f" g="%.0f" b="%.0f"/>' % tuple(scrgbclr)
    return objectify.fromstring(s)


def xmlns(*prefixes):
    '''XML ns.'''
    elem_schema = make_element()
    return ' '.join(f'xmlns:{pre}="{elem_schema["nsmap"][pre]}"' for pre in prefixes)


def call(val, g, group, default):
    '''Callback.'''
    if callable(val):
        return val(g)
    return default


def cust_shape(x, y, w, h, _id):
    '''Custom shapes.'''
    return objectify.fromstring(
        f'''
        <p:sp {xmlns("p", "a")}>
          <p:nvSpPr>
            <p:cNvPr id='{_id}' name='Freeform {_id}'/>
            <p:cNvSpPr/>
            <p:nvPr/>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm>
              <a:off x='{x}' y='{y}'/>
              <a:ext cx='{w}' cy='{h}'/>
            </a:xfrm>
            <a:custGeom>
              <a:avLst/>
              <a:gdLst/>
              <a:ahLst/>
              <a:cxnLst/>
              <a:rect l='0' t='0' r='0' b='0'/>
            </a:custGeom>
          </p:spPr>
        </p:sp>'''
    )


def draw_sankey(data, spec):
    '''Create sankey data logic.'''
    x0 = spec['x0']
    size = spec['size']
    group = spec['group']
    width = spec['width']
    default_color = '#ccfccf'
    default_stroke = '#ffffff'
    attrs = spec.get('attrs', {})
    sort = spec.get('sort', False)

    text = spec.get('text')
    order = spec.get('order')
    fill_color = spec.get('color')

    g = data.groupby(group)
    frame = pd.DataFrame(
        {
            'size': g[group[0]].count() if size is None else g[size].sum(),
            'seq': 0 if order is None else order(g),
        }
    )
    frame['width'] = frame['size'] / float(frame['size'].sum()) * width
    frame['fill'] = call(fill_color, g, group, default_color)
    result = call(text, g, group, '')
    frame['text'] = result
    # Add all attrs to the frame as well
    for key, val in attrs.items():
        frame[key] = call(val, g, group, None)
    if 'stroke' not in attrs:
        frame['stroke'] = default_stroke
    # Compute frame['x'] only after sorting
    if order and sort:
        frame.sort_values('seq', inplace=True)
    frame['x'] = x0 + frame['width'].cumsum() - frame['width']
    return frame


def squarified(x, y, w, h, data):
    '''
    Draw a squarified treemap.

    See <http://citeseerx.ist.psu.edu/viewdoc/summary?doi=10.1.1.36.6685>
    Returns a numpy array with (x, y, w, h) for each item in data.

    Examples::

        # The result is a 2x2 numpy array::
        >>> squarified(x=0, y=0, w=6, h=4, data=[6, 6, 4, 3, 2, 2, 1])
        array([[ 0.        ,  0.        ,  3.        ,  2.        ],
               [ 0.        ,  2.        ,  3.        ,  2.        ],
               [ 3.        ,  0.        ,  1.71428571,  2.33333333],
               [ 4.71428571,  0.        ,  1.28571429,  2.33333333],
               [ 3.        ,  2.33333333,  1.2       ,  1.66666667],
               [ 4.2       ,  2.33333333,  1.2       ,  1.66666667],
               [ 5.4       ,  2.33333333,  0.6       ,  1.66666667]])

        >>> squarified(x=0, y=0, w=1, h=1, data=[np.nan, 0, 1, 2])
        array([[ 0.        ,  0.        ,  0.        ,  0.        ],
               [ 0.        ,  0.        ,  0.        ,  0.        ],
               [ 0.        ,  0.        ,  0.33333333,  1.        ],
               [ 0.33333333,  0.        ,  0.66666667,  1.        ]])
    '''
    w, h = float(w), float(h)
    size = np.nan_to_num(np.array(data).astype(float))
    start, end = 0, len(size)
    result = np.zeros([end, 4])
    if w <= 0 or h <= 0:
        return result

    cumsize = np.insert(size.cumsum(), 0, 0)
    while start < end:
        # We lay out out blocks of rects on either the left or the top edge of
        # the remaining rectangle. But how many rects in the block? We take as
        # many as we can as long as the worst aspect ratio of the block's
        # rectangles keeps improving.

        # This section is (and should be) be heavily optimised. Each operation
        # is run on every element in data.
        last_aspect, newstart = np.Inf, start + 1
        startsize = cumsize[start]
        blockmin = blockmax = size[newstart - 1]
        blocksum = cumsize[newstart] - startsize
        datasum = cumsize[end] - startsize
        ratio = datasum * (h / w if w > h else w / h)
        while True:
            f = blocksum * blocksum / ratio
            aspect = blockmax / f if blockmax > f else f / blockmax
            aspect2 = blockmin / f if blockmin > f else f / blockmin
            if aspect2 > aspect:
                aspect = aspect2
            if aspect <= last_aspect:
                if newstart < end:
                    last_aspect = aspect
                    newstart += 1
                    val = size[newstart - 1]
                    if val < blockmin:
                        blockmin = val
                    if val > blockmax:
                        blockmax = val
                    blocksum += val
                else:
                    break
            else:
                if newstart > start + 1:
                    newstart = newstart - 1
                break

        # Now, lay out the block = start:newstart on the left or top edge.
        block = slice(start, newstart)
        blocksum = cumsize[newstart] - startsize
        scale = blocksum / datasum
        blockcumsize = cumsize[block] - startsize

        if w > h:
            # Layout left-edge, downwards
            r = h / blocksum
            result[block, 0] = x
            result[block, 1] = y + r * blockcumsize
            result[block, 2] = dx = w * scale
            result[block, 3] = r * size[block]
            x, w = x + dx, w - dx
        else:
            # Layout top-edge, rightwards
            r = w / blocksum
            result[block, 0] = x + r * blockcumsize
            result[block, 1] = y
            result[block, 2] = r * size[block]
            result[block, 3] = dy = h * scale
            y, h = y + dy, h - dy

        start = newstart

    return np.nan_to_num(result)


class SubTreemap:
    '''
    Yield a hierarchical treemap at multiple levels.

    Usage:
        SubTreemap(
            data=data,
            keys=['Parent', 'Child'],
            values={'Value':sum},
            size=lambda x: x['Value'],
            sort=None,
            padding=0,
            aspect=1)

    yields:
        x, y, w, h, (level, data)
    '''

    def __init__(self, **args):
        '''Default Constructor.'''
        self.args = args

    def draw(self, width, height, x=0, y=0, filter={}, level=0):
        '''Function to draw rectanfles.'''
        # We recursively into each column in `keys` and stop there
        if level >= len(self.args['keys']):
            return

        # Start with the base dataset. Filter by each key applied so far
        summary = self.args['data']
        for key in filter:
            summary = summary[summary[key] == filter[key]]

        # Aggregate by the key up to the current level
        summary = summary.groupby(self.args['keys'][: level + 1]).agg(self.args.get('values', {}))
        for key in self.args['keys'][: level + 1]:
            if hasattr(summary, 'reset_index'):
                # Just pop the key out. .reset_index(key) should do this.
                # But on Pandas 0.20.1, this fails
                summary = summary.reset_index([summary.index.names.index(key)])
            else:
                summary[key] = summary.index

        # If specified, sort the aggregated data
        if 'sort' in self.args and callable(self.args['sort']):
            summary = self.args['sort'](summary)

        pad = self.args.get('padding', 0)
        aspect = self.args.get('aspect', 1)

        # Find the positions of each box at this level
        key = self.args['keys'][level]
        rows = summary.to_records() if hasattr(summary, 'to_records') else summary

        rects = squarified(x, y * aspect, width, height * aspect, self.args['size'](rows))
        for i2, (x2, y2, w2, h2) in enumerate(rects):
            v2 = rows[i2]
            y2, h2 = y2 / aspect, h2 / aspect
            # Ignore invalid boxes generated by Squarified
            if (
                np.isnan([x2, y2, w2, h2]).any()
                or np.isinf([x2, y2, w2, h2]).any()
                or w2 < 0
                or h2 < 0
            ):
                continue

            # For each box, dive into the next level
            filter2 = dict(filter)
            filter2.update({key: v2[key]})
            yield from self.draw(
                w2 - 2 * pad, h2 - 2 * pad, x=x2 + pad, y=y2 + pad, filter=filter2, level=level + 1
            )
            # Once we've finished yielding smaller boxes, yield the parent box
            yield x2, y2, w2, h2, (level, v2)


class TableProperties:
    '''Get/Set Table's properties.'''

    def extend_table(self, shape, data, total_rows, total_columns):
        '''Function to extend table rows and columns if required.'''
        avail_rows = len(shape.table.rows)
        avail_cols = len(shape.table.columns)

        col_width = shape.table.columns[0].width
        row_height = shape.table.rows[0].height
        # Extending Table Rows if required based on the data
        while avail_rows < total_rows:
            shape.table.rows._tbl.add_tr(row_height)
            avail_rows += 1
        # Extending Table Columns if required based on the data
        while avail_cols < total_columns:
            shape.table._tbl.tblGrid.add_gridCol(col_width)
            avail_cols += 1

    def get_default_css(self, shape):
        '''Function to get Table style for rows and columns.'''
        pixel_inch = 10000
        tbl_style = {}
        mapping = {0: 'header', 1: 'row'}
        for row_num in range(len(list(shape.table.rows)[:2])):
            style = {}
            txt = shape.table.rows[row_num].cells[0].text_frame.paragraphs[0]

            if txt.alignment:
                style['text-align'] = '{}'.format(txt.alignment).split()[0]

            if not hasattr(txt, 'runs'):
                txt.add_run()
            if txt.runs:
                txt = txt.runs[0].font
                style['bold'] = txt.bold
                style['italic'] = txt.italic
                style['font-size'] = (txt.size / pixel_inch) if txt.size else txt.size
                style['font-family'] = txt.name
                style['underline'] = txt.underline
            tbl_style[mapping[row_num]] = style

        if 'row' not in tbl_style or not len(tbl_style['row']):
            tbl_style['row'] = copy.deepcopy(tbl_style['header'])

        if 'font-size' not in tbl_style['row']:
            tbl_style['row']['font-size'] = tbl_style['header'].get('font-size', None)

        return tbl_style

    def get_css(self, info, column_list, data):
        '''Get Table CSS from config.'''
        columns = info.get('columns', {})
        table_css = {}
        for col in column_list:
            common_css = copy.deepcopy(info.get('style', {}))
            common_css.update(columns.get(col, {}))
            if 'gradient' in common_css:
                common_css['min'] = common_css.get('min', data[col].min())
                common_css['max'] = common_css.get('max', data[col].max())
            table_css[col] = common_css
        return table_css

    def apply_table_css(self, cell, paragraph, run, info):
        '''Apply Table style.'''
        if info.get('fill'):
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = RGBColor.from_string(convert_color_code(info['fill']))
        apply_text_css(cell, run, paragraph, **info)
