'''PPTGen module.'''
import os
import sys
import copy
import collections
from pptx import Presentation
from pptx.shapes.shapetree import SlideShapes
from orderedattrdict import AttrDict
import pandas as pd
import gramex.data
import gramex.cache
from gramex.config import merge
from gramex import parse_command_line
from gramex.transforms import build_transform
from . import commands
from .utils import stack_shapes, delete_slide, generate_slide, manage_slides
from .utils import is_slide_allowed, is_group, add_new_slide, copy_slide_elem

COMMANDS_LIST = commands.cmdlist


def commandline():
    '''
    Runs PPTGen from the command line.
    This is called via setup.py console_scripts.
    Though a trivial function, is is kept different from run_commands to allow
    unit testing of run_commands.
    '''
    run_commands(sys.argv[1:], pptgen)


def run_commands(commands, callback):
    '''
    For example::

        run_commands(['a.yaml', 'b.yaml', '--x=1'], method)

    will do the following:

    - Load a.yaml into config
        - Set config['a'] = 1
        - Change to directory where a.yaml is
        - Call method(config)
    - Load b.yaml into config
        - Set config['a'] = 1
        - Change to directory where b.yaml is
        - Call method(config)

    Command line arguments are passed as ``commands``.
    Callback is a function that is called for each config file.
    '''
    args = parse_command_line(commands)
    original_path = os.getcwd()
    for config_file in args.pop('_'):
        config = gramex.cache.open(config_file, 'config')
        config = merge(old=config, new=args, mode='overwrite')
        os.chdir(os.path.dirname(os.path.abspath(config_file)))
        try:
            callback(**config)
        finally:
            os.chdir(original_path)


def load_data(data_config, handler=None):
    '''
    Loads data using gramex cache.
    '''
    if not isinstance(
        data_config,
        (
            dict,
            AttrDict,
        ),
    ):
        raise ValueError('Data argument must be a dict like object.')

    data = {}
    for key, conf in data_config.items():
        if isinstance(
            conf,
            (
                dict,
                AttrDict,
            ),
        ):
            if 'function' in conf:
                data[key] = build_transform(conf, vars={'handler': None})(handler=handler)[0]
            elif conf.get('ext') in {'yaml', 'yml', 'json'}:
                data[key] = gramex.cache.open(conf.pop('url'), conf.pop('ext'), **dict(conf))
            elif 'url' in conf:
                data[key] = gramex.data.filter(conf.pop('url'), **dict(conf))
        else:
            data[key] = conf
    return data


def replicate_slides(data, prs, change, slide, slides_to_remove, index, handler):
    '''
    Function to replicate slides.
    '''
    if isinstance(data, pd.DataFrame):
        data = data.to_dict(orient='records')
    copy_slide = copy.deepcopy(slide)
    slides_to_remove.append(index)
    # Stacking shapes if required.
    stack_shapes(copy_slide.shapes, change, data, handler)
    new_slide = generate_slide(prs, copy_slide)
    args = {'prs': prs, 'copy_slide': True, 'source_slide': slide, 'new_slide': new_slide}
    change_shapes(copy_slide.shapes, change, data, handler, **args)


def register(config):
    '''Function to register a new `command` to command list.'''
    global COMMANDS_LIST
    resister_command = config.pop('register', {})
    if not isinstance(resister_command, (dict,)):
        raise ValueError('Register should be a dict like object')
    for command_name, command_function in resister_command.items():
        if command_name not in COMMANDS_LIST:
            if not isinstance(command_function, (dict,)):
                command_function = {'function': command_function}
            _vars = {'shape': None, 'spec': None, 'data': None}
            COMMANDS_LIST[command_name] = build_transform(command_function, vars=_vars)


def pptgen(source, target=None, **config):
    '''
    Process a configuration. This loads a Presentation from source, applies the
    (optional) configuration changes and saves it into target.
    '''
    # Config was being over written using PPTXHandler and data key was being
    # removed from yaml config.
    handler = config.pop('handler', None)
    _config = copy.deepcopy(config)
    if _config.get('is_formhandler', False):
        data = _config.pop('data')
        _config.pop('is_formhandler')
    else:
        data = AttrDict(load_data(_config.pop('data', {}), handler=handler))
    # Register a `command` if present in configuration
    register(_config)

    # Loading input template
    prs = Presentation(source)
    # Removing not required slides from presentation.
    prs = manage_slides(prs, _config)
    slides = prs.slides
    # Loop through each change configuration
    slides_to_remove = []
    manage_slide_order = collections.defaultdict(list)

    for key, change in _config.items():
        # Apply it to every slide
        slide_data = copy.deepcopy(data)
        if 'data' in change and change['data'] is not None:
            if not isinstance(change['data'], (dict,)):
                change['data'] = {'function': change.pop('data')}
            slide_data = build_transform(change['data'], vars={'data': None})(slide_data)[0]

        for index, slide in enumerate(slides):
            # Restrict to specific slides, if specified
            if not is_slide_allowed(change, slide, index + 1):
                continue
            if change.get('replicate'):
                is_grp = isinstance(slide_data, pd.core.groupby.DataFrameGroupBy)
                if isinstance(slide_data, collections.Iterable):
                    for _slide_data in slide_data:
                        _slide_data = _slide_data[1] if is_grp is True else _slide_data
                        replicate_slides(
                            _slide_data, prs, change, slide, slides_to_remove, index, handler
                        )
                        # Creating dict mapping to order slides.
                        manage_slide_order[index + 1].append(len(prs.slides))
                else:
                    raise NotImplementedError()
            else:
                # Stacking shapes if required.
                stack_shapes(slide.shapes, change, slide_data, handler)
                change_shapes(slide.shapes, change, slide_data, handler)

    indexes = []
    slides_to_remove = list(set(slides_to_remove))
    for key in sorted(manage_slide_order.keys()):
        indexes.append(manage_slide_order[key])

    matrix = list(map(list, zip(*indexes)))

    for indx_lst in matrix:
        for idx in indx_lst:
            src = prs.slides[idx - 1]
            slides_to_remove.append(idx - 1)
            copy_slide = copy.deepcopy(src)
            new_slide = generate_slide(prs, copy_slide)
            dest = prs.slides.add_slide(new_slide)
            for shape in copy_slide.shapes:
                copy_slide_elem(shape, dest)
            add_new_slide(dest, src)
    for removed_status, sld_idx in enumerate(set(slides_to_remove)):
        delete_slide(prs, (sld_idx - removed_status))
        for slide_num in manage_slide_order:
            manage_slide_order[slide_num] = [(i - 1) for i in manage_slide_order[slide_num]]
    if target is None:
        return prs
    else:
        prs.save(target)


def change_shapes(collection, change, data, handler, **kwargs):
    '''
    Apply changes to a collection of shapes in the context of data.
    ``collection`` is a slide.shapes or group shapes.
    ``change`` is typically a dict of <shape-name>: commands.
    ``data`` is a dictionary passed to the template engine.
    '''
    prs = kwargs.get('prs')
    new_slide = kwargs.get('new_slide')
    copy_slide = kwargs.get('copy_slide', False)
    source_slide = kwargs.get('source_slide')

    dest = prs.slides.add_slide(new_slide) if copy_slide else None
    mapping = {}
    for shape in collection:
        if shape.name not in change:
            copy_slide_elem(shape, dest)
            continue

        spec = change[shape.name]
        if shape.name not in mapping:
            mapping[shape.name] = 0

        if spec.get('data'):
            if not isinstance(spec['data'], (dict,)):
                spec['data'] = {
                    'function': '{}'.format(spec['data'])
                    if not isinstance(spec['data'], str)
                    else spec['data']
                }
            shape_data = build_transform(spec['data'], vars={'data': None, 'handler': None})(
                data=data, handler=handler
            )[0]
        else:
            if (
                isinstance(
                    data,
                    (
                        dict,
                        AttrDict,
                    ),
                )
                and 'handler' in data
            ):
                data.pop('handler')
            shape_data = copy.deepcopy(data)

        if isinstance(
            shape_data,
            (
                dict,
                AttrDict,
            ),
        ):
            shape_data['handler'] = handler

        if spec.get('stack'):
            shape_data = shape_data[mapping[shape.name]]
        mapping[shape.name] = mapping[shape.name] + 1
        # If the shape is a group, apply spec to each sub-shape
        if is_group(shape):
            sub_shapes = SlideShapes(shape.element, collection)
            change_shapes(sub_shapes, spec, shape_data, handler)
        # Add args to shape_data
        if hasattr(handler, 'args'):
            args = {k: v[0] for k, v in handler.args.items() if len(v) > 0}
            shape_data['args'] = args
        # Run commands in the spec
        for cmd, method in COMMANDS_LIST.items():
            if cmd in spec:
                method(shape, spec, shape_data)
        copy_slide_elem(shape, dest)
    add_new_slide(dest, source_slide)
