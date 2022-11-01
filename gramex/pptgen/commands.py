'''Python-PPTX customized module.'''
import os
import copy
import logging
import requests
import tempfile
import operator
import collections
import numpy as np
import pandas as pd
import matplotlib.cm
import matplotlib.colors

# B410:import_lxml lxml.etree is safe on https://github.com/tiran/defusedxml/tree/main/xmltestdata
from lxml.etree import fromstring  # nosec B410
from tornado.template import Template
from tornado.escape import to_unicode
from pptx.chart import data as pptxcd
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from urllib.parse import urlparse
from gramex.transforms import build_transform
from . import utils
from . import fontwidth
from . import color as _color


_template_cache = {}


def template(tmpl, data):
    '''Execute tornado template'''
    if tmpl not in _template_cache:
        _template_cache[tmpl] = Template(tmpl, autoescape=None)
    return to_unicode(_template_cache[tmpl].generate(**data))


def text(shape, spec, data):
    '''Replace entire text of shape with spec['text']'''
    if not shape.has_text_frame:
        logging.error(f'"{shape.name}" is not a TextShape to apply text:')
        return
    if not isinstance(data, (dict,)):
        data = {'data': data}
    handler = data.pop('handler', None)

    style = copy.deepcopy(spec.get('style', {}))
    style = generate_style(style, data, handler)
    pixel_inch = 10000
    # Get paragraph
    paragraph = shape.text_frame.paragraphs[0]
    # Removing the extra paragraphs if more than one present
    for para in shape.text_frame.paragraphs[1:]:
        utils.delete_paragraph(para)
    # Removing the extra run in paragraph if more than one present
    for _run in paragraph.runs[1:]:
        utils.delete_run(_run)

    theme, brightness, default_css = None, None, {}
    # Calculatiing default style
    clrtype = paragraph.runs[0].font.color.type
    if clrtype and 'rgb' in '{}'.format(clrtype).lower().split():
        default_css['color'] = '{}'.format(paragraph.runs[0].font.color.rgb)
    elif clrtype and 'scheme' in '{}'.format(clrtype).lower().split():
        theme = paragraph.runs[0].font.color.theme_color
        brightness = paragraph.runs[0].font.color.brightness
    for prop in {'bold', 'italic', 'underline'}:
        default_css[prop] = getattr(paragraph.runs[0].font, prop)
    if paragraph.runs[0].font.size:
        default_css['font-size'] = paragraph.runs[0].font.size / pixel_inch
    default_css['text-align'] = paragraph.alignment
    default_css['font-family'] = paragraph.runs[0].font.name
    # Updating default css with css from config.
    default_css.update(style)
    default_css['color'] = default_css.get('color', '#0000000')
    # B320: lxml.etree is safe on https://github.com/tiran/defusedxml/tree/main/xmltestdata
    update_text = fromstring('<root>{}</root>'.format(template(spec['text'], data)))  # nosec B320
    paragraph.runs[0].text = update_text.text if update_text.text else ''
    utils.apply_text_css(shape, paragraph.runs[0], paragraph, **default_css)
    index = 1
    for child in update_text.getchildren():
        if not child.tag == 'text':
            raise ValueError('XML elemet must contain only "text" tag.')
        # Adding runs for each text item
        for idx, new_txt in enumerate([child.text, child.tail]):
            if not new_txt:
                continue
            common_css = {key: val for key, val in default_css.items()}
            paragraph.add_run()
            # Adding text to run
            paragraph.runs[index].text = new_txt
            # Updating the text css
            common_css.update(dict(child.items())) if not idx else None
            if theme:
                paragraph.runs[index].font.color.theme_color = theme
            if brightness:
                paragraph.runs[index].font.color.brightness = brightness
            utils.apply_text_css(shape, paragraph.runs[index], paragraph, **common_css)
            index += 1


def replace(shape, spec, data):
    '''Replace keywords in shape using the dictionary at spec['replace']'''
    if not shape.has_text_frame:
        logging.error(f'"{shape.name}" is not a TextShape to apply text:')
        return
    if not isinstance(data, (dict,)):
        data = {'data': data}

    handler = spec.get('handler', None)
    common_css = spec.get('style', {})
    style = {}
    for old, new in spec['replace'].items():
        _style = common_css.pop(old, {})
        style[old] = generate_style(_style, {'val': template(new, data), 'data': data}, handler)

    common_css = generate_style(common_css, data, handler)
    for key, val in style.items():
        css = copy.deepcopy(common_css)
        css.update(val)
        style[key] = css

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old, new in spec['replace'].items():
                run.text = run.text.replace(old, template(new, data))
                utils.apply_text_css(shape, run, paragraph, **style)


def image(shape, spec, data):
    '''Replace image with a different file specified in spec['image']'''
    image = template(spec['image'], data)
    # If it's a URL, use the requests library's raw stream as a file-like object
    if urlparse(image).netloc:
        r = requests.get(image)
        with tempfile.NamedTemporaryFile(delete=False) as handle:
            handle.write(r.content)
        new_img_part, new_rid = shape.part.get_or_add_image_part(handle.name)
        os.unlink(handle.name)
    else:
        new_img_part, new_rid = shape.part.get_or_add_image_part(image)
    shape._pic.blipFill.blip.rEmbed = new_rid
    shape.part.related_parts[new_rid].blob = new_img_part.blob


def generate_style(style, data, handler):
    '''Function to conpile style section from kwargs.'''
    for key, value in style.items():
        if isinstance(value, (dict,)) and 'function' in value:
            result = compile_function(style, key, data, handler)
            style[key] = result(data) if callable(result) else result
    return style


def rect_css(shape, **kwargs):
    '''Function to add text to shape.'''
    for key in {'fill', 'stroke'}:
        if kwargs.get(key):
            fill = shape.fill if key == 'fill' else shape.line.fill
            rectcss = kwargs[key].rsplit('#')[-1].lower()
            rectcss = rectcss + ('0' * (6 - len(rectcss)))
            chart_css(fill, kwargs, rectcss)


def add_text_to_shape(shape, textval, **kwargs):
    '''Function to add text to shape.'''
    min_inc = 13000
    pixel_inch = 10000
    if (kwargs.get('font-size', 14) * pixel_inch) < min_inc:
        return
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.add_run()
    for run in paragraph.runs:
        run.text = textval
        shape_txt = run.font.fill
        shape_txt.solid()
        utils.apply_text_css(shape, run, paragraph, **kwargs)


def scale_data(data, lo, hi, factor=None):
    '''Function to scale data.'''
    data = np.array(data, dtype=float)
    return ((data - lo) / ((hi - lo) or np.nan)) * (factor or 1.0)


def rect(shape, x, y, width, height):
    '''Add rectangle to slide.'''
    return shape.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)


def _update_chart(info, data, chart_data, series_columns, chart='ChartData'):
    '''Updating Chart data.'''
    if chart == 'ChartData':
        chart_data.categories = data[info['x']].dropna().unique().tolist()
        for series in series_columns:
            if np.issubdtype(data[series].dtype, np.number):
                chart_data.add_series(series, tuple(data[series].fillna(0).values.tolist()))
        return chart_data
    series_dict = {}
    columns = data.columns.difference([info['x']])
    is_numeric_x = np.issubdtype(data[info['x']].dtype, np.number)
    xindex = data[info['x']].astype(float) if is_numeric_x else pd.Series(data.index + 1)
    for index, row in data.fillna(0).iterrows():
        x = xindex.loc[index]
        for col in series_columns:
            if col not in columns or not np.issubdtype(data[col].dtype, np.number):
                continue
            serieslist = [series.name for series in chart_data._series]
            if col not in serieslist:
                series_dict[col] = chart_data.add_series(col)
            if chart == 'XyChartData':
                series_dict[col].add_data_point(x, row[col])
            elif chart == 'BubbleChartData':
                bubble_size = row[info['size']] if info.get('size') else 1
                if col != info.get('size'):
                    series_dict[col].add_data_point(x, row[col], bubble_size)
    return chart_data


def chart_css(fill, style, color):
    '''Function to add opacity to charts.'''
    fill.solid()
    pix_to_inch = 100000
    fill.fore_color.rgb = RGBColor.from_string(utils.convert_color_code(color))
    solid_fill = fill.fore_color._xFill
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '%d' % (pix_to_inch * style.get('opacity', 1.0)))
    solid_fill.srgbClr.append(alpha)
    return fill


def compile_function(spec, key, data, handler):
    '''A function to compile configuration.'''
    if key not in spec:
        return None
    _vars = {'_color': None, 'data': None, 'handler': None}
    if not isinstance(spec[key], (dict,)):
        spec[key] = {'function': '{}'.format(spec[key])}
    elif isinstance(spec[key], (dict,)) and 'function' not in spec[key]:
        spec[key] = {'function': '{}'.format(spec[key])}
    args = {'data': data, 'handler': handler, '_color': _color}
    return build_transform(spec[key], vars=_vars)(**args)[0]


def table(shape, spec, data):
    '''Update an existing Table shape with data.'''
    if not shape.has_table:
        raise AttributeError('Shape must be a table object.')
    if not spec.get('table', {}).get('data'):
        return
    spec = copy.deepcopy(spec['table'])
    handler = data.pop('handler') if 'handler' in data else None
    data = compile_function(spec, 'data', data, handler)
    if not len(data):
        return
    data_cols = data.columns
    data_cols_len = len(data_cols)
    table_properties = utils.TableProperties()
    # Extending table if required.
    table_properties.extend_table(shape, data, len(data) + 1, data_cols_len)
    # Fetching Table Style for All Cells and texts.
    tbl_style = table_properties.get_default_css(shape)
    input_colspec = spec.get('columns', {})
    input_cols = input_colspec.keys()
    if all(isinstance(x, int) for x in input_cols):
        for x in list(input_cols):
            if x < data_cols_len:
                input_colspec[data_cols[x]] = input_colspec.pop(x)
    styled_cols = data.columns.intersection(input_cols or data_cols)
    cell_style = table_properties.get_css(spec, styled_cols, data)
    data = data.to_dict(orient='records')
    for row_num, row in enumerate(shape.table.rows):
        cols = len(row.cells._tr.tc_lst)
        # Extending cells in newly added rows.
        while cols < len(data_cols):
            row.cells._tr.add_tc()
            cols += 1
        for col_num, cell in enumerate(row.cells):
            colname = data_cols[col_num]
            for paragraph in cell.text_frame.paragraphs:
                if not paragraph.text.strip():
                    paragraph.add_run()
                for run in paragraph.runs:
                    txt = colname if row_num == 0 else data[row_num - 1][colname]
                    run.text = '{}'.format(txt)
                    cellcss = {} if row_num == 0 else copy.deepcopy(cell_style.get(colname, {}))
                    txt_css = copy.deepcopy(tbl_style.get('header' if row_num == 0 else 'row', {}))
                    if row_num > 0 and 'gradient' in cellcss:
                        grad_txt = scale_data(txt, cellcss['min'], cellcss['max'])
                        gradient = matplotlib.cm.get_cmap(cellcss['gradient'])
                        cellcss['fill'] = matplotlib.colors.to_hex(gradient(grad_txt))
                    if cellcss.get('fill'):
                        txt_css['color'] = cellcss.get('color', _color.contrast(cellcss['fill']))
                        cellcss['color'] = cellcss.get('color', _color.contrast(cellcss['fill']))
                    txt_css.update(cellcss)
                    table_properties.apply_table_css(cell, paragraph, run, txt_css)


def chart(shape, spec, data):
    '''Replacing chart Data.'''
    if not shape.has_chart:
        raise AttributeError('Shape must be a chart object.')

    chart_type = None
    if hasattr(shape.chart, 'chart_type'):
        chart_type = '{}'.format(shape.chart.chart_type).split()[0]

    chart_types = {
        'ChartData': {
            'AREA',
            'AREA_STACKED',
            'AREA_STACKED_100',
            'BAR_CLUSTERED',
            'BAR_OF_PIE',
            'BAR_STACKED',
            'BAR_STACKED_100',
            'COLUMN_CLUSTERED',
            'COLUMN_STACKED',
            'COLUMN_STACKED_100',
            'LINE',
            'LINE_MARKERS',
            'LINE_MARKERS_STACKED',
            'LINE_MARKERS_STACKED_100',
            'LINE_STACKED',
            'LINE_STACKED_100',
            'RADAR_MARKERS',
            'RADAR',
            'RADAR_FILLED',
            'PIE',
            'PIE_EXPLODED',
            'PIE_OF_PIE',
            'DOUGHNUT',
            'DOUGHNUT_EXPLODED',
        },
        'XyChartData': {
            'XY_SCATTER',
            'XY_SCATTER_LINES',
            'XY_SCATTER_LINES_NO_MARKERS',
            'XY_SCATTER_SMOOTH',
            'XY_SCATTER_SMOOTH_NO_MARKERS',
        },
        'BubbleChartData': {'BUBBLE', 'BUBBLE_THREE_D_EFFECT'},
    }

    if not chart_type:
        raise NotImplementedError()

    info = copy.deepcopy(spec['chart'])
    # Load data
    handler = data.pop('handler') if 'handler' in data else None
    for prop in {'x', 'size', 'usecols'}:
        if prop in info and isinstance(info[prop], (dict,)):
            if 'function' not in info[prop]:
                info[prop]['function'] = '{}'.format(info[prop])
            info[prop] = compile_function(info, prop, data, handler)

    style = {
        'color': info.pop('color', None),
        'opacity': info.pop('opacity', None),
        'stroke': info.pop('stroke', None),
    }

    for key in {'color', 'stroke', 'opacity'}:
        if key in style and isinstance(style[key], (dict,)):
            if 'function' not in style[key]:
                style[key]['function'] = '{}'.format(style[key])
            style[key] = compile_function(style, key, data, handler)

    data = compile_function(info, 'data', data, handler)
    # Getting subset of data if `usecols` is defined.
    change_data = data.reset_index(drop=True)[info.get('usecols', data.columns)]
    series_cols = change_data.columns.drop(info['x']) if info['x'] else change_data.columns
    chart_name = next((k for k in chart_types if chart_type in chart_types[k]), None)

    if not chart_name:
        raise NotImplementedError('Input Chart Type {} is not supported'.format(chart_type))

    chart_data = _update_chart(
        info, change_data, getattr(pptxcd, chart_name)(), series_cols, chart=chart_name
    )

    shape.chart.replace_data(chart_data)
    if chart_name == 'scatter' and not style.get('color'):
        series_names = [series.name for series in shape.chart.series]
        style['color'] = dict(zip(series_names, _color.distinct(len(series_names))))

    if style.get('color'):
        color_mapping = {
            'XyChartData': 'point.marker.format',
            'ChartData': 'point.format',
            'BubbleChartData': 'point.format',
            'area': 'series.format',
        }
        is_area = {'AREA', 'AREA_STACKED', 'AREA_STACKED_100'}
        chart_name = 'area' if chart_type in is_area else chart_name
        is_donut = {'PIE', 'PIE_EXPLODED', 'PIE_OF_PIE', 'DOUGHNUT', 'DOUGHNUT_EXPLODED'}
        for series in shape.chart.series:
            for index, point in enumerate(series.points):
                row = data.loc[index]
                args = {
                    'handler': handler,
                    'row': row.to_dict(),
                    'name': row[info['x']] if chart_type in is_donut else series.name,
                }
                point_css = {}
                for key in {'opacity', 'color', 'stroke'}:
                    if style.get(key) is None:
                        continue
                    point_css[key] = style[key]
                    if callable(style[key]):
                        prop = style[key](index)
                        point_css[key] = prop(**args) if callable(prop) else prop
                    elif isinstance(style[key], (dict)):
                        point_css[key] = style[key].get(args['name'], '#cccccc')
                # Evaluatinig line and shape color format
                for series_point in {'point', 'series'}:
                    # Replacing point with series to change color in legend
                    fillpoint = color_mapping[chart_name].replace('point', series_point)
                    # B307:eval this is safe since `expr` is written by app developer
                    chart_css(eval(fillpoint).fill, point_css, point_css['color'])  # nosec B307
                    # Will apply on outer line of chart shape line(like stroke in html)
                    _stroke = point_css.get('stroke', point_css['color'])
                    chart_css(eval(fillpoint).line.fill, point_css, _stroke)  # nosec B307


# Custom Charts Functions below(Sankey, Treemap, Calendarmap).


def sankey(shape, spec, data):
    '''Draw sankey in PPT.'''
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    # Getting parent shapes
    pxl_to_inch = 10000
    default_thickness = 40
    spec = copy.deepcopy(spec['sankey'])
    handler = data.pop('handler') if 'handler' in data else None
    y0 = shape.top
    x0 = shape.left
    width = shape.width
    height = shape.height
    shapes = shape._parent
    shape_ids = {'shape': 0}

    groups = compile_function(spec, 'groups', data, handler)
    thickness = spec.get('thickness', default_thickness) * pxl_to_inch
    h = (height - (thickness * len(groups))) / (len(groups) - 1) + thickness
    frames = {}
    # Sankey Rectangles and texts.
    sankey_conf = {}
    for k in ['size', 'order', 'text', 'color']:
        sankey_conf[k] = compile_function(spec, k, data, handler)
    sankey_conf['x0'] = x0
    sankey_conf['width'] = width
    sankey_conf['attrs'] = spec.get('attrs', {})
    sankey_conf['sort'] = spec.get('sort')
    stroke = spec.get('stroke', '#ffffff')
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    elem_schema = utils.make_element()
    data = compile_function(spec, 'data', data, handler)
    for ibar, group in enumerate(groups):
        y = y0 + h * ibar
        sankey_conf['group'] = [group]
        df = frames[group] = utils.draw_sankey(data, sankey_conf)
        # Adding rectangle
        for key, row in df.iterrows():
            shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, row['x'], y, row['width'], thickness)
            rectstyle = {'fill': row['fill'], 'stroke': stroke}
            rect_css(shp, **rectstyle)
            text_style = {'color': _color.contrast(row['fill'])}
            text_style.update(spec.get('style', {}))
            add_text_to_shape(shp, row['text'], **text_style)

    # Sankey Connection Arcs.
    for ibar, (group1, group2) in enumerate(zip(groups[:-1], groups[1:])):
        sankey_conf['group'] = [group1, group2]
        sankey_conf['sort'] = False
        df = utils.draw_sankey(data, sankey_conf)
        pos = collections.defaultdict(float)
        for key1, row1 in frames[group1].iterrows():
            for key2, row2 in frames[group2].iterrows():
                if (key1, key2) in df.index:
                    row = df.loc[(key1, key2)]
                    y1, y2 = y0 + h * ibar + thickness, y0 + h * (ibar + 1)
                    ym = (y1 + y2) / 2
                    x1 = row1['x'] + pos[0, key1]
                    x2 = row2['x'] + pos[1, key2]

                    _id = shape_ids['shape'] = shape_ids['shape'] + 1
                    shp = utils.cust_shape(
                        0, 0, '{:.0f}'.format(row['width']), '{:.0f}'.format(ym), _id
                    )
                    path = elem_schema['a'].path(
                        w='{:.0f}'.format(row['width']), h='{:.0f}'.format(ym)
                    )
                    shp.find('.//a:custGeom', namespaces=elem_schema['nsmap']).append(
                        elem_schema['a'].pathLst(path)
                    )
                    path.append(
                        elem_schema['a'].moveTo(
                            elem_schema['a'].pt(
                                x='{:.0f}'.format(x1 + row['width']), y='{:.0f}'.format(y1)
                            )
                        )
                    )

                    path.append(
                        elem_schema['a'].cubicBezTo(
                            elem_schema['a'].pt(
                                x='{:.0f}'.format(x1 + row['width']), y='{:.0f}'.format(ym)
                            ),
                            elem_schema['a'].pt(
                                x='{:.0f}'.format(x2 + row['width']), y='{:.0f}'.format(ym)
                            ),
                            elem_schema['a'].pt(
                                x='{:.0f}'.format(x2 + row['width']), y='{:.0f}'.format(y2)
                            ),
                        )
                    )

                    path.append(
                        elem_schema['a'].lnTo(
                            elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(y2))
                        )
                    )

                    path.append(
                        elem_schema['a'].cubicBezTo(
                            elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(ym)),
                            elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(ym)),
                            elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(y1)),
                        )
                    )

                    path.append(elem_schema['a'].close())
                    shp.spPr.append(
                        elem_schema['a'].solidFill(utils.fill_color(srgbclr=row['fill']))
                    )
                    shapes._spTree.append(shp)
                    pos[0, key1] += row['width']
                    pos[1, key2] += row['width']


def treemap(shape, spec, data):
    '''Function to download data as ppt.'''
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    shapes = shape._parent
    x0 = shape.left
    y0 = shape.top
    width = shape.width
    height = shape.height
    spec = copy.deepcopy(spec['treemap'])
    stroke = spec.get('stroke', '#ffffff')
    # Load data
    handler = data.pop('handler') if 'handler' in data else None
    for k in ['keys', 'values', 'size', 'sort', 'color', 'text', 'data']:
        spec[k] = compile_function(spec, k, data, handler)
    # Getting rectangle's width and height using `squarified` algorithm.
    treemap_data = utils.SubTreemap(**spec)
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    font_aspect = 14.5
    pixel_inch = 10000
    default_rect_color = '#cccccc'
    for x, y, w, h, (level, v) in treemap_data.draw(width, height):
        if level == 0:
            shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, x + x0, y + y0, w, h)
            rect_color = default_rect_color
            if spec.get('color'):
                rect_color = spec['color'](v) if callable(spec['color']) else spec['color']
            if spec.get('text'):
                text = spec['text'](v) if callable(spec['text']) else spec['text']
            else:
                text = '{}'.format(v[1])
            rectstyle = {'fill': rect_color, 'stroke': stroke}
            rect_css(shp, **rectstyle)
            font_size = min(h, w * font_aspect / fontwidth.fontwidth('{}'.format(text)), np.Inf)
            text_style = {'color': _color.contrast(rect_color)}
            text_style.update(spec.get('style', {}))
            text_style['font-size'] = font_size / pixel_inch
            # Adding text inside rectangles
            add_text_to_shape(shp, text, **text_style)


def calendarmap(shape, spec, data):
    '''Draw calendar map in PPT.'''
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()

    shapes = shape._parent
    spec = copy.deepcopy(spec['calendarmap'])
    handler = data.get('handler')
    # Load data
    data = compile_function(spec, 'data', data, handler).fillna(0)
    startdate = compile_function(spec, 'startdate', data, handler)

    pixel_inch = 10000
    size = spec.get('size', None)

    label_top = spec.get('label_top', 0) * pixel_inch
    label_left = spec.get('label_left', 0) * pixel_inch

    width = spec['width'] * pixel_inch
    shape_top = label_top + shape.top
    shape_left = label_left + shape.left
    y0 = width + shape_top
    x0 = width + shape_left

    # Deleting the shape
    shape.element.delete()
    # Style
    default_color = '#ffffff'
    default_line_color = '#787C74'
    default_txt_color = '#000000'
    style = copy.deepcopy(spec.get('style', {}))
    style = generate_style(style, data, handler)

    font_size = style.get('font-size', 12)
    stroke = style.get('stroke', '#ffffff')
    fill_rect = style.get('fill', '#cccccc')
    text_color = style.get('color', '#000000')
    # Treat infinities as nans when calculating scale
    scaledata = pd.Series(data).replace([np.inf, -np.inf], np.nan)
    for key in {'lo', 'hi', 'weekstart'}:
        if isinstance(spec.get(key), (dict,)) and 'function' in spec.get(key):
            spec[key] = compile_function(spec, key, data, handler)

    lo_data = spec.get('lo', scaledata.min())
    range_data = spec.get('hi', scaledata.max()) - lo_data
    gradient = matplotlib.cm.get_cmap(spec.get('gradient', 'RdYlGn'))
    color = style.get(
        'fill',
        lambda v: matplotlib.colors.to_hex(gradient((float(v) - lo_data) / range_data))
        if not pd.isnull(v)
        else default_color,
    )

    startweekday = (startdate.weekday() - spec.get('weekstart', 0)) % 7
    # Weekday Mean and format
    weekday_mean = pd.Series([scaledata[(x - startweekday) % 7 :: 7].mean() for x in range(7)])
    weekday_format = spec.get('format', '{:,.%df}' % utils.decimals(weekday_mean.values))
    # Weekly Mean and format
    weekly_mean = pd.Series(
        [scaledata[max(0, x) : x + 7].mean() for x in range(-startweekday, len(scaledata), 7)]
    )
    weekly_format = spec.get('format', '{:,.%df}' % utils.decimals(weekly_mean.values))
    # Scale sizes as square roots from 0 to max (not lowest to max -- these
    # should be an absolute scale)
    sizes = (
        width * utils.scale([v**0.5 for v in size], lo=0)
        if size is not None
        else [width] * len(scaledata)
    )
    for i, val in enumerate(data):
        nx = (i + startweekday) // 7
        ny = (i + startweekday) % 7
        d = startdate + pd.DateOffset(days=i)
        fill = '#cccccc'
        if not pd.isnull(val):
            fill = color(val) if callable(color) else color

        shp = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x0 + (width * nx) + (width - sizes[i]) / 2,
            y0 + (width * ny) + (width - sizes[i]) / 2,
            sizes[i],
            sizes[i],
        )
        rectstyle = {'fill': fill, 'stroke': stroke(val) if callable(stroke) else stroke}
        rect_css(shp, **rectstyle)
        text_style = {
            'color': style.get('color')(val)
            if callable(style.get('color'))
            else spec.get('color', _color.contrast(fill)),
            'font-size': font_size(val) if callable(font_size) else font_size,
        }
        for k in ['bold', 'italic', 'underline', 'font-family']:
            text_style[k] = style.get(k)
        add_text_to_shape(shp, '%02d' % d.day, **text_style)

        # Draw the boundary lines between months
        if i >= 7 and d.day == 1 and ny > 0:
            border = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0 + width * nx, y0 + (width * ny), width, 2 * pixel_inch
            )
            border.name = 'border'
            rectstyle = {'fill': default_line_color, 'stroke': default_line_color}
            rect_css(border, **rectstyle)
        if i >= 7 and d.day <= 7 and nx > 0:
            border = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0 + (width * nx), y0 + (width * ny), 2 * pixel_inch, width
            )
            border.name = 'border'
            rectstyle = {'fill': default_line_color, 'stroke': default_line_color}
            rect_css(border, **rectstyle)
        # Adding weekdays text to the chart (left side)
        if i < 7:
            txt = shapes.add_textbox(
                x0 - (width / 2), y0 + (width * ny) + (width / 2), width, width
            )
            text_style['color'] = default_txt_color
            add_text_to_shape(txt, d.strftime('%a')[0], **text_style)
        # Adding months text to the chart (top)
        if d.day <= 7 and ny == 0:
            txt = shapes.add_textbox(x0 + (width * nx), y0 - (width / 2), width, width)
            text_style['color'] = default_txt_color
            add_text_to_shape(txt, d.strftime('%b %Y'), **text_style)
    if label_top:
        lo_weekly = spec.get('lo', weekly_mean.min())
        range_weekly = spec.get('hi', weekly_mean.max()) - lo_weekly
        for nx, val in enumerate(weekly_mean.fillna(0)):
            w = label_top * ((val - lo_weekly) / range_weekly)
            px = x0 + (width * nx)
            bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, px, shape_top - w, width, w)
            bar.name = 'summary.top.bar'
            rectstyle = {
                'fill': fill_rect(val) if callable(fill_rect) else fill_rect,
                'stroke': stroke(val) if callable(stroke) else stroke,
            }
            rect_css(bar, **rectstyle)
            label = shapes.add_textbox(px, shape_top - width, width, width)
            label.name = 'summary.top.label'
            text_style['color'] = text_color(val) if callable(text_color) else text_color
            add_text_to_shape(label, weekly_format.format(weekly_mean[nx]), **text_style)
    if label_left:
        lo_weekday = spec.get('lo', weekday_mean.min())
        range_weekday = spec.get('hi', weekday_mean.max()) - lo_weekday
        for ny, val in enumerate(weekday_mean.fillna(0)):
            w = label_left * ((val - lo_weekday) / range_weekday)
            bar = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, shape_left - w, y0 + (width * ny), w, width
            )
            bar.name = 'summary.left.bar'
            rectstyle = {
                'fill': fill_rect(val) if callable(fill_rect) else fill_rect,
                'stroke': stroke(val) if callable(stroke) else stroke,
            }
            rect_css(bar, **rectstyle)
            label = shapes.add_textbox(shape_left - width, y0 + (width * ny), w, width)
            label.name = 'summary.left.label'
            text_style['color'] = text_color(val) if callable(text_color) else text_color
            add_text_to_shape(label, weekday_format.format(weekday_mean[ny]), **text_style)


def bullet(shape, spec, data):
    '''Function to plot bullet chart.'''
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    spec = copy.deepcopy(spec['bullet'])

    orient = spec.get('orient', 'horizontal')
    if orient not in {'horizontal', 'vertical'}:
        raise NotImplementedError()

    font_aspect = 5
    pixel_inch = 10000
    x = shape.left
    y = shape.top

    handler = data.get('handler')
    for met in ['poor', 'average', 'good', 'target']:
        spec[met] = compile_function(spec, met, data, handler) if spec.get(met) else np.nan

    height = shape.height if orient == 'horizontal' else shape.width
    width = shape.width if orient == 'horizontal' else shape.height
    if spec.get('max-width'):
        max_width = compile_function(spec, 'max-width', data, handler)
        width = max_width(width) if callable(max_width) else max_width * width

    spec['data'] = compile_function(spec, 'data', data, handler)
    gradient = spec.get('gradient', 'RdYlGn')
    shapes = shape._parent
    shape._sp.delete()
    lo = spec.get('lo', 0)
    hi = spec.get(
        'hi',
        np.nanmax([spec['data'], spec['target'], spec['poor'], spec['average'], spec['good']]),
    )
    style = {}
    common_style = copy.deepcopy(spec.get('style', {}))
    data_text = common_style.get('data', {}).pop('text', spec.get('text', True))
    target_text = common_style.get('target', {}).pop('text', spec.get('text', True))
    if data_text:
        data_text = compile_function({'text': data_text}, 'text', data, handler)
    if target_text:
        target_text = compile_function({'text': target_text}, 'text', data, handler)

    css = {
        'data': common_style.pop('data', {}),
        'target': common_style.pop('target', {}),
        'poor': common_style.pop('poor', {}),
        'good': common_style.pop('good', {}),
        'average': common_style.pop('average', {}),
    }

    for key, val in css.items():
        _style = copy.deepcopy(common_style)
        _style.update(val)
        for css_prop, css_val in _style.items():
            if isinstance(css_val, (dict,)) and 'function' in css_val:
                _style[css_prop] = compile_function(_style, css_prop, data, handler)
        style[key] = _style

    gradient = matplotlib.cm.get_cmap(gradient)
    percentage = {'good': 0.125, 'average': 0.25, 'poor': 0.50, 'data': 1.0, 'target': 1.0}
    for index, metric in enumerate(['good', 'average', 'poor']):
        scaled = scale_data(spec.get(metric, np.nan), lo, hi, factor=width)
        if not np.isnan(scaled):
            _width = scaled if orient == 'horizontal' else height
            _hight = height if orient == 'horizontal' else scaled
            yaxis = y if orient == 'horizontal' else y + (width - scaled)
            _rect = rect(shapes, x, yaxis, _width, _hight)
            fill = style.get(metric, {})
            stroke = fill.get('stroke')
            fill = fill.get('fill', matplotlib.colors.to_hex(gradient(percentage[metric])))
            rect_css(_rect, fill=fill, stroke=stroke or fill)

    getmax = {key: spec.get(key, np.nan) for key in ['data', 'target', 'good', 'average', 'poor']}
    max_data_val = percentage[max(getmax.items(), key=operator.itemgetter(1))[0]]

    scaled = scale_data(spec['data'], lo, hi, factor=width)
    if not np.isnan(scaled):
        _width = scaled if orient == 'horizontal' else height / 2.0
        yaxis = y + height / 4.0 if orient == 'horizontal' else y + (width - scaled)
        xaxis = x if orient == 'horizontal' else x + height / 4.0
        _hight = height / 2.0 if orient == 'horizontal' else scaled
        data_rect = rect(shapes, xaxis, yaxis, _width, _hight)
        fill = style.get('data', {})
        stroke = fill.get('stroke')
        fill = fill.get('fill', matplotlib.colors.to_hex(gradient(1.0)))
        rect_css(data_rect, fill=fill, stroke=stroke or fill)

    if data_text and not np.isnan(scaled):
        if callable(data_text):
            _data_text = '{}'.format(data_text(spec['data']))
        else:
            _data_text = '{}'.format(spec['data']) if data_text is True else data_text
        parent = data_rect._parent
        text_width = (_width if orient == 'vertical' else _hight * 2) * len(_data_text)
        _xaxis = xaxis if orient == 'vertical' else x + scaled - text_width
        parent = parent.add_textbox(_xaxis, yaxis, text_width, text_width / len(_data_text))
        data_txt_style = style.get('data', {})
        data_txt_style['color'] = data_txt_style.get('color', _color.contrast(fill))
        default_align = 'left' if orient == 'vertical' else 'right'
        data_txt_style['text-align'] = data_txt_style.get('text-align', default_align)
        # Setting default font-size
        font_size = (
            (text_width / pixel_inch) * font_aspect / fontwidth.fontwidth('{}'.format(_data_text))
        )
        font_size = min(text_width / pixel_inch, font_size, np.Inf)
        data_txt_style['font-size'] = data_txt_style.get('font-size', font_size)
        add_text_to_shape(parent, _data_text, **data_txt_style)

    scaled = scale_data(spec['target'], lo, hi, factor=width)
    if not np.isnan(scaled):
        line_hight = 10000
        _width = line_hight if orient == 'horizontal' else height
        _hight = height if orient == 'horizontal' else line_hight
        yaxis = y if orient == 'horizontal' else (width - scaled) + y
        xaxis = x + scaled if orient == 'horizontal' else x
        target_line = rect(shapes, xaxis, yaxis, _width, _hight)
        fill = style.get('target', {})
        stroke = fill.get('stroke')
        fill_target_rect = fill.get('fill', matplotlib.colors.to_hex(gradient(1.0)))
        rect_css(target_line, fill=fill_target_rect, stroke=stroke or fill_target_rect)
        if target_text:
            if callable(target_text):
                _target_text = '{}'.format(target_text(spec['target']))
            else:
                _target_text = '{}'.format(spec['target']) if target_text is True else target_text
            handler = data.get('handler')
            parent = target_line._parent
            yaxis = yaxis - (_width / 2) if orient == 'vertical' else yaxis
            text_width = (_width if orient == 'vertical' else _hight) * len(_target_text)
            parent = parent.add_textbox(xaxis, yaxis, text_width, text_width / len(_target_text))
            target_txt_style = style.get('target', {})
            fill_max = fill.get('fill', matplotlib.colors.to_hex(gradient(max_data_val)))
            target_txt_style['color'] = target_txt_style.get('color', _color.contrast(fill_max))
            # Setting default font-size
            font_size = font_aspect / fontwidth.fontwidth('{}'.format(_target_text))
            font_size = min(text_width / pixel_inch, (text_width / pixel_inch) * font_size, np.Inf)
            target_txt_style['font-size'] = target_txt_style.get('font-size', font_size)
            add_text_to_shape(parent, _target_text, **target_txt_style)


def heatgrid(shape, spec, data):
    '''Create a heat grid.'''
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()

    spec = copy.deepcopy(spec['heatgrid'])

    top = shape.top
    left = shape.left
    width = shape.width
    pixel_inch = 10000
    default_height = 20
    height = spec.get('cell-height', default_height) * pixel_inch
    parent = shape._parent
    shape.element.delete()

    # Loading config
    handler = data.pop('handler') if 'handler' in data else None
    for key in ['row', 'column', 'value', 'column-order', 'row-order']:
        if key not in spec:
            continue
        if isinstance(spec[key], (dict,)) and 'function' in spec[key]:
            spec[key] = compile_function(spec, key, data, handler)
    # Loading data
    data = compile_function(spec, 'data', data, handler)
    data = data.sort_values(by=[spec['column']])
    rows = spec.get('row-order') or sorted(data[spec['row']].unique().tolist())
    columns = spec.get('column-order') or sorted(data[spec['column']].unique().tolist())

    left_margin = width * spec.get('left-margin', 0.15)
    padding = spec.get('style', {}).get('padding', 5)
    if not isinstance(padding, (dict,)):
        padding = {'left': padding, 'right': padding, 'top': padding, 'bottom': padding}

    styles = copy.deepcopy(spec.get('style', {}))

    if styles.get('gradient'):
        _min, _max = data[spec['value']].min(), data[spec['value']].max()
    # Compiling style elements if required
    for key in ['gradient', 'color', 'fill', 'font-size', 'font-family', 'stroke']:
        if isinstance(styles.get(key), (dict,)) and 'function' in styles[key]:
            prop = compile_function(styles, key, data, handler)
            styles[key] = prop(data=data, handler=handler) if callable(prop) else prop
    # Calculating cell's width based on config
    _width = (width - left_margin) / float(len(columns)) / pixel_inch
    _width = spec.get('cell-width', _width) * pixel_inch
    # Adding Columns to the HeatGrid.
    for idx, column in enumerate(columns):
        txt = parent.add_textbox(left + _width * idx + left_margin, top - height, _width, height)
        add_text_to_shape(txt, '{}'.format(column), **styles)
    # Cell width
    for index, row in enumerate(rows):
        _data = data[data[spec['row']] == row].dropna()
        _data = pd.merge(
            pd.DataFrame({spec['column']: list(columns)}),
            _data,
            left_on=spec['column'],
            right_on=spec['column'],
            how='left',
        ).reset_index(drop=True)

        for _idx, _row in _data.iterrows():
            style = copy.deepcopy(styles)
            # Setting callable padding args
            _vars = {'handler': None, 'row': None, 'column': None, 'value': None}
            args = {
                'handler': handler,
                'row': row,
                'column': _row[spec['column']],
                'value': _row[spec['value']],
            }
            # Setting padding if callable.
            _pad = copy.deepcopy(padding)
            for key, val in _pad.items():
                if isinstance(val, (dict,)) and 'function' in val:
                    _pad[key] = build_transform(val, vars=_vars)(**args)[0]
            top_pad = _pad.get('top', 5) * pixel_inch
            left_pad = _pad.get('left', 5) * pixel_inch
            right_pad = _pad.get('right', 5) * pixel_inch
            bottom_pad = _pad.get('bottom', 5) * pixel_inch

            # Adding cells
            xaxis = left + (_width * _idx) + left_margin + left_pad
            yaxis = top + (height * index) + (top_pad) * index
            _rect = rect(parent, xaxis, yaxis, _width - left_pad - right_pad, height - top_pad)
            # Adding color gradient to cell if gradient is True
            if style.get('gradient'):
                grad_txt = scale_data(_row[spec['value']], _min, _max)
                gradient = matplotlib.cm.get_cmap(style['gradient'])
                style['fill'] = matplotlib.colors.to_hex(gradient(grad_txt))
                style['color'] = _color.contrast(style['fill'])
            if np.isnan(_row[spec['value']]) and spec.get('na-color'):
                style['fill'] = spec.get('na-color')
                style['color'] = _color.contrast(style['fill'])

            style['stroke'] = style.get('stroke', style['fill'])
            rect_css(_rect, **style)
            # Adding text to cells if required.
            if spec.get('text'):
                _txt = parent.add_textbox(
                    xaxis, yaxis, _width - left_pad - right_pad, height - top_pad - bottom_pad
                )
                if isinstance(spec['text'], dict) and 'function' in spec['text']:
                    cell_txt = compile_function(spec, 'text', _row, handler)
                else:
                    cell_txt = '{}'.format(_row[spec['value']])
                if pd.isnull(cell_txt) and spec.get('na-text'):
                    cell_txt = spec.get('na-text')
                add_text_to_shape(_txt, cell_txt, **style)
        # Adding row's text in left side
        txt = parent.add_textbox(
            left, top + (height * index) + top_pad * index, _width + left_margin, height
        )
        add_text_to_shape(txt, row, **styles)


def css(shape, spec, data):
    '''Function to modify a rectangle's property in PPT.'''
    pxl_to_inch = 10000
    handler = data.pop('handler') if 'handler' in data else None
    spec = copy.deepcopy(spec['css'])
    data = compile_function(spec, 'data', data, handler)
    style = copy.deepcopy(spec.get('style', {}))
    shape_prop = {'width', 'height', 'top', 'left'}
    for prop in shape_prop:
        setprop = style.get(prop)
        if setprop:
            if not isinstance(style[prop], (dict,)):
                style[prop] = {
                    'function': '{}'.format(style[prop])
                    if not isinstance(style[prop], str)
                    else style[prop]
                }
            setprop = compile_function(style, prop, data, handler)
            setprop = setprop * pxl_to_inch
        else:
            setprop = getattr(shape, prop)
        setattr(shape, prop, setprop)

    _style = {}
    for key, val in style.items():
        if key not in shape_prop:
            _style[key] = val
            if isinstance(val, (dict,)):
                _style[key] = compile_function(style, key, data, handler)
            _style[key] = _style[key](data) if callable(_style[key]) else _style[key]
    rect_css(shape, **_style)


cmdlist = {
    'css': css,
    'text': text,
    'image': image,
    'chart': chart,
    'table': table,
    'sankey': sankey,
    'bullet': bullet,
    'replace': replace,
    'treemap': treemap,
    'heatgrid': heatgrid,
    'calendarmap': calendarmap,
}
